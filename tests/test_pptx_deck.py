"""Tests de `paginer_items` (pptx_deck.py) — bin-packing pur, sans
dépendance à python-pptx/DB/HTTP, utilisé par la pagination auto de
l'export PPT (voir pptx_export.py).

Complété (arbitrage 2026-09-03) par les tests des « helpers durcis deck
binaire » remontés depuis VSCode4 (clear_slides, supprimer_slide,
purger_rels_slides_orphelines, trouver_slide_par_titre, sans_puce) et des
formes/texte « riches » (add_forme, definir_geometrie, configurer_text_frame,
definir_paragraphes, add_text_runs, trouver_cadre_layout) — même esprit que le
reste du fichier : aucune dépendance au domaine métier/DB de ce projet, la
bibliothèque doit rester réutilisable telle quelle par la flotte."""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches

from app.services.pptx_deck import (
    _normaliser,
    add_forme,
    add_rect,
    add_text_runs,
    clear_slides,
    configurer_text_frame,
    definir_geometrie,
    definir_paragraphes,
    paginer_items,
    purger_rels_slides_orphelines,
    sans_puce,
    supprimer_slide,
    trouver_cadre_layout,
    trouver_slide_par_titre,
)


def test_paginer_items_splits_when_capacity_exceeded() -> None:
    assert paginer_items([1, 1, 1, 1], lambda x: x, 2) == [[1, 1], [1, 1]]


def test_paginer_items_single_page_when_everything_fits() -> None:
    assert paginer_items([1, 1, 1], lambda x: x, 10) == [[1, 1, 1]]


def test_paginer_items_oversized_item_alone_on_its_page() -> None:
    assert paginer_items([3], lambda x: x, 2) == [[3]]
    assert paginer_items([1, 3, 1], lambda x: x, 2) == [[1], [3], [1]]


def test_paginer_items_empty_input_returns_one_empty_page() -> None:
    assert paginer_items([], lambda x: x, 2) == [[]]


def test_paginer_items_preserves_order_and_drops_nothing() -> None:
    items = list(range(10))
    pages = paginer_items(items, lambda x: 1, 3)
    assert [x for page in pages for x in page] == items


# --------------------------------------------------------------------------- #
# Helpers durcis deck binaire (clear_slides, supprimer_slide,
# purger_rels_slides_orphelines, trouver_slide_par_titre) — priorité de test
# demandée : ce sont les 4 qui manipulent des relations OOXML/slides.
# --------------------------------------------------------------------------- #


def _slide_rid(prs, slide) -> str:
    """rId de `slide` dans sldIdLst — pour vérifier qu'une relation est bien
    lâchée (ou au contraire orpheline) après manipulation."""
    for sld_id in prs.slides._sldIdLst:
        if int(sld_id.get("id")) == slide.slide_id:
            return sld_id.get(qn("r:id"))
    raise AssertionError("slide absente de sldIdLst")


def _slide_reltypes(prs) -> list[str]:
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    return [rid for rid, rel in prs.part.rels.items() if rel.reltype == RT.SLIDE]


def _new_prs_with_slides(n: int, titres=None):
    prs = Presentation()
    layout = prs.slide_layouts[6]  # layout vierge
    slides = []
    for i in range(n):
        slide = prs.slides.add_slide(layout)
        texte = (titres[i] if titres else f"Slide {i}")
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(1))
        box.text_frame.text = texte
        slides.append(slide)
    return prs, slides


def test_clear_slides_retire_toutes_les_slides_et_lache_les_rels() -> None:
    prs, _ = _new_prs_with_slides(3)
    assert len(prs.slides) == 3
    assert len(_slide_reltypes(prs)) == 3
    clear_slides(prs)
    assert len(prs.slides) == 0
    assert list(prs.slides._sldIdLst) == []
    # Le point qui compte : plus AUCUNE relation de type slide ne traîne
    # (sinon PowerPoint refuse d'ouvrir le fichier réserialisé).
    assert _slide_reltypes(prs) == []


def test_supprimer_slide_ne_retire_que_la_slide_ciblee_et_lache_son_rel() -> None:
    prs, (s0, s1) = _new_prs_with_slides(2, ["Garder", "Retirer"])
    rid_retire = _slide_rid(prs, s1)
    supprimer_slide(prs, s1)
    assert len(prs.slides) == 1
    assert list(prs.slides)[0] is s0
    assert rid_retire not in prs.part.rels, "relation de la slide supprimée non lâchée (drop_rel)"
    # La slide restante garde sa propre relation intacte.
    assert _slide_rid(prs, s0) in prs.part.rels


def test_supprimer_slide_leve_value_error_si_slide_deja_absente() -> None:
    """Un 2e appel sur la même slide échoue — mais le message vient de
    python-pptx, pas de notre garde explicite : `slide.slide_id` (1re ligne de
    supprimer_slide) recherche la part dans sldIdLst par identité et lève déjà
    ValueError une fois la relation lâchée par le 1er appel ; la garde
    `"slide id=... absente de sldIdLst"` de supprimer_slide n'est donc jamais
    atteinte par cette voie (elle documente une invariant "ne devrait jamais
    arriver" plutôt qu'un cas réellement accessible via l'API publique)."""
    prs, (s0,) = _new_prs_with_slides(1)
    supprimer_slide(prs, s0)
    with pytest.raises(ValueError):
        supprimer_slide(prs, s0)


def test_purger_rels_slides_orphelines_nettoie_sans_toucher_les_relations_actives() -> None:
    prs, (s0, s1) = _new_prs_with_slides(2)
    rid_orphelin = _slide_rid(prs, s1)
    # Reproduit la corruption historique : retirer l'entrée sldIdLst SANS
    # lâcher la relation (l'erreur que supprimer_slide/clear_slides évitent).
    for sld_id in list(prs.slides._sldIdLst):
        if sld_id.get(qn("r:id")) == rid_orphelin:
            prs.slides._sldIdLst.remove(sld_id)
    assert rid_orphelin in prs.part.rels, "précondition : la relation orpheline existe encore"
    purges = purger_rels_slides_orphelines(prs)
    assert purges == 1
    assert rid_orphelin not in prs.part.rels
    # La relation de la slide restante (toujours référencée) n'est pas touchée.
    assert _slide_rid(prs, s0) in prs.part.rels
    assert purger_rels_slides_orphelines(prs) == 0, "un deck sain ne doit plus rien purger"


def test_trouver_slide_par_titre_trouve_par_egalite_normalisee() -> None:
    prs, (s0, s1) = _new_prs_with_slides(2, ["  Slide   Alpha ", "Slide Beta"])
    idx, slide = trouver_slide_par_titre(prs, "slide alpha")  # casse + espaces différents
    assert idx == 0
    assert slide is s0


def test_trouver_slide_par_titre_leve_si_aucun_ou_plusieurs_matches() -> None:
    prs, _ = _new_prs_with_slides(2, ["Un titre", "Un autre titre"])
    with pytest.raises(ValueError, match="0 slide"):
        trouver_slide_par_titre(prs, "Titre absent")
    prs2, _ = _new_prs_with_slides(2, ["Titre dupliqué", "Titre dupliqué"])
    with pytest.raises(ValueError, match="2 slide"):
        trouver_slide_par_titre(prs2, "Titre dupliqué")


def test_normaliser_replie_espaces_et_ignore_la_casse() -> None:
    assert _normaliser("  Un   Titre \n de Slide ") == _normaliser("un titre de slide")


# --------------------------------------------------------------------------- #
# sans_puce
# --------------------------------------------------------------------------- #


def test_sans_puce_retire_indentation_et_force_buNone() -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = "01"
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(int(Inches(0.5))))
    pPr.set("indent", str(int(Inches(-0.25))))
    pPr.append(pPr.makeelement(qn("a:buChar"), {"char": "•"}))

    sans_puce(p)

    pPr = p._p.get_or_add_pPr()
    assert pPr.get("marL") == "0"
    assert pPr.get("indent") == "0"
    assert pPr.findall(qn("a:buChar")) == []
    assert pPr.findall(qn("a:buAutoNum")) == []
    assert pPr.findall(qn("a:buNone")) != [], "buNone absent — la puce peut survivre"


# --------------------------------------------------------------------------- #
# Formes/texte « riches » : add_forme, definir_geometrie, configurer_text_frame,
# definir_paragraphes, add_text_runs, trouver_cadre_layout.
# --------------------------------------------------------------------------- #


def test_add_forme_applique_preset_adjustments_et_alpha() -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shp = add_forme(slide, "roundRect", 1, 1, 2, 1, fill="#2c5cc5", adj=[0.2],
                     fill_alpha=50)
    spPr = shp._element.spPr
    geom = spPr.find(qn("a:prstGeom"))
    assert geom.get("prst") == "roundRect"
    gd = geom.find(qn("a:avLst")).find(qn("a:gd"))
    assert gd.get("fmla") == "val 20000"  # adj[0]=0.2 -> échelle OOXML 100000
    srgb = spPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
    alpha = srgb.find(qn("a:alpha"))
    assert alpha is not None and alpha.get("val") == "50000"  # 50 * 1000


def test_add_forme_prst_inconnu_leve_key_error() -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    with pytest.raises(KeyError):
        add_forme(slide, "prst_inexistant", 1, 1, 1, 1)


def test_definir_geometrie_pose_position_et_taille() -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shp = add_rect(slide, 0, 0, 1, 1)
    definir_geometrie(shp, 2.5, 1.25, 4.0, 0.5)
    assert shp.left == Inches(2.5)
    assert shp.top == Inches(1.25)
    assert shp.width == Inches(4.0)
    assert shp.height == Inches(0.5)


def test_configurer_text_frame_ne_touche_que_les_champs_fournis() -> None:
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.3)

    configurer_text_frame(tf, anchor=MSO_ANCHOR.MIDDLE)  # wrap/autosize/margins omis

    assert tf.vertical_anchor == MSO_ANCHOR.MIDDLE
    assert tf.word_wrap is True, "word_wrap non fourni : ne doit pas être écrasé"
    assert tf.auto_size == MSO_AUTO_SIZE.NONE
    assert tf.margin_left == Inches(0.3), "marge non fournie : ne doit pas être écrasée"


def test_definir_paragraphes_pose_puce_reelle_avec_retrait_suspendu() -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    definir_paragraphes(box.text_frame, [
        ([("Un point important", {"size": 12, "bold": True})],
         {"bullet": {"char": "•", "color": "#2c5cc5"}}),
    ])
    p = box.text_frame.paragraphs[0]
    assert p.runs[0].text == "Un point important"
    assert p.runs[0].font.bold is True
    pPr = p._p.get_or_add_pPr()
    assert pPr.get("marL") == str(int(Inches(0.1875)))
    assert pPr.get("indent") == str(int(Inches(-0.1875)))
    buChar = pPr.find(qn("a:buChar"))
    assert buChar is not None and buChar.get("char") == "•"


def test_add_text_runs_melange_plusieurs_styles_dans_un_paragraphe() -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = add_text_runs(slide, 1, 1, 3, 1, [
        ([("Normal ", {}), ("gras", {"bold": True, "size": 14})], {}),
    ])
    p = box.text_frame.paragraphs[0]
    assert [r.text for r in p.runs] == ["Normal ", "gras"]
    assert p.runs[0].font.bold is not True
    assert p.runs[1].font.bold is True


def test_trouver_cadre_layout_desambigue_par_largeur_minimale() -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_forme(slide, "round2DiagRect", 0.5, 0.5, 1.0, 1.0)  # petit cadre, à ignorer
    grand = add_forme(slide, "round2DiagRect", 2.0, 2.0, 3.0, 2.0)  # grand cadre, attendu

    trouve = trouver_cadre_layout(slide.shapes, "round2DiagRect", largeur_min_in=2.0)
    assert trouve is not None
    left, top, width, height, geom, flip = trouve
    assert Emu(left).inches == pytest.approx(2.0)
    assert Emu(width).inches == pytest.approx(3.0)
    assert geom.get("prst") == "round2DiagRect"
    assert flip == (False, False)
    assert grand.left == left


def test_trouver_cadre_layout_renvoie_none_si_aucun_preset_ne_matche() -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_forme(slide, "ellipse", 0, 0, 1, 1)
    assert trouver_cadre_layout(slide.shapes, "round2DiagRect") is None
