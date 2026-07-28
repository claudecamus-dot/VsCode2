"""Test « rendu de qualité » du deck de restitution (P6 refonte, 2026-07-22).

Consolide les contrôles qualité des slides au-delà du seul `verifier_geometrie`
(formes hors-cadre) : STRUCTURE narrative (couverture de marque, 4 intercalaires de
chapitre, slides de contenu titrées), CHARTE (police de marque Outfit réellement
appliquée au texte dessiné), et le garde-fou géométrie (build lève sinon). C'est le
filet automatisé complémentaire à `pptx-verify` (rendu visuel humain) : ce dernier
voit ce qu'un test ne voit pas (collisions, vides), ce test verrouille ce que l'œil
ne re-vérifie pas à chaque fois (structure/charte/géométrie).
"""
from __future__ import annotations

from app.db import DB_PATH, SessionLocal, engine, init_db
from app.models import (
    GlobalSynthesis, Interview, Mission, MissionDifficulty, MissionExecutiveSummary,
    MissionSwot, Question, Recommendation, RecommendationAxis, Theme, Trame, Verbatim,
)
from app.services import pptx_deck as D
from app.services.pptx_export import build_presentation


def setup_module() -> None:
    try:
        engine.dispose()
    except Exception:
        pass
    if DB_PATH.exists():
        DB_PATH.unlink()
    init_db()


def teardown_module() -> None:
    try:
        engine.dispose()
    except Exception:
        pass
    if DB_PATH.exists():
        DB_PATH.unlink()


def _mission_complete() -> int:
    """Mission couvrant les 4 chapitres (exec summary, synthèse, difficultés, SWOT,
    verbatims, axes+recos) — de quoi exercer toute la structure du deck."""
    db = SessionLocal()
    try:
        m = Mission(name="Audit qualité — Deck complet")
        db.add(m); db.flush()
        tr = Trame(mission_id=m.id); db.add(tr); db.flush()
        th = Theme(trame_id=tr.id, title="T", position=0); db.add(th); db.flush()
        q = Question(theme_id=th.id, label="Q", qtype="open", position=0); db.add(q); db.flush()
        iv = Interview(mission_id=m.id, interviewee_name="Témoin", status="done"); db.add(iv); db.flush()
        v = Verbatim(interview_id=iv.id, question_id=q.id, quote="Une citation représentative.")
        db.add(v); db.flush()
        m.restitution_verbatim_ids = [v.id]
        db.add(GlobalSynthesis(
            mission_id=m.id, status="generated", contexte="- C", culture_adn="- C",
            forces_succes="- F", points_amelioration="- Silos\n- Dette", aspirations="- A"))
        db.add(MissionExecutiveSummary(
            mission_id=m.id, status="generated", headline="Un titre-claim de synthèse",
            points="- Point un\n- Point deux", key_message="Le message à retenir."))
        m.difficulties = [MissionDifficulty(position=0, label="Silos entre équipes")]
        db.add(MissionSwot(
            mission_id=m.id, status="generated", forces="- F", faiblesses="- Fa",
            opportunites="- O", menaces="- Me"))
        ax = RecommendationAxis(mission_id=m.id, title="Axe 1", position=0); db.add(ax); db.flush()
        # Volumes réalistes (pas « - a ») : le filet verifier_debordements_texte
        # exécuté sur cette fixture doit stresser de VRAIES hauteurs de texte
        # (revue adversariale — une fixture quasi vide ne vérifiait rien).
        db.add(Recommendation(
            axis_id=ax.id, position=0,
            title="Mettre en place une instance de gouvernance data décisionnelle",
            objectif="Créer un mécanisme régulier et légitime pour arbitrer les "
                     "priorités et valider les règles communes de gestion.",
            acteurs="CDO, DSI, représentants métiers",
            valeur=5, complexite=2,
            proposition_valeur="Réduire les délais de décision et concentrer les "
                               "capacités sur les cas d'usage les plus porteurs.",
            plan_actions="- Définir le mandat et les droits de décision\n"
                         "- Nommer les membres permanents\n"
                         "- Organiser un comité mensuel et publier les décisions",
            resultats_attendus="- Priorités partagées, diminution des escalades et "
                               "meilleure allocation des ressources data."))
        db.commit()
        return m.id
    finally:
        db.close()


def _layouts(prs) -> list[str]:
    return [(s.slide_layout.name or "").lower() for s in prs.slides]


def test_deck_structure_narrative() -> None:
    """Structure : couverture de marque, sommaire, 4 intercalaires de chapitre,
    chaque slide de contenu (« titre seul ») porte un titre non vide."""
    db = SessionLocal()
    try:
        prs = build_presentation(db.get(Mission, _mission_complete()))
    finally:
        db.close()
    layouts = _layouts(prs)
    assert any("couverture" in l for l in layouts), "couverture de marque OCTO manquante"
    assert sum(1 for l in layouts if "chapitre" in l) == 4, "4 intercalaires de chapitre attendus"
    # Chaque slide de contenu a du texte (titre) — pas de slide vide.
    for s in prs.slides:
        if "titre seul" in (s.slide_layout.name or "").lower():
            has_text = any(
                sh.has_text_frame and sh.text_frame.text.strip() for sh in s.shapes
            )
            assert has_text, "slide de contenu sans aucun texte"


def test_deck_charte_police_du_theme_appliquee_pas_outfit() -> None:
    """Charte / rendu : le deck applique au texte la police du THÈME (Arial sur OCTO,
    une police système donc rendue telle quelle), et surtout PAS l'Outfit des
    placeholders — qui, n'étant pas installée, serait rendue en substitution (c'était
    la cause du « la police ne matche pas la référence » ; cf. police_theme vs
    police_marque dans pptx_deck). Ce test est le garde-fou anti-régression du fix."""
    db = SessionLocal()
    try:
        prs = build_presentation(db.get(Mission, _mission_complete()))
    finally:
        db.close()
    theme = D.police_theme(prs)
    assert theme, "police du thème non détectée sur le template"
    fonts = {
        r.font.name
        for s in prs.slides for sh in s.shapes if sh.has_text_frame
        for p in sh.text_frame.paragraphs for r in p.runs
        if r.font.name  # ignorer les runs sans police explicite (héritage)
    }
    assert theme in fonts, f"police du thème {theme!r} non appliquée au texte dessiné"
    # Garde-fou dur : l'Outfit des placeholders (non installée) ne doit JAMAIS être
    # forcée sur le texte — sinon rendu en substitution, le bug d'origine.
    placeholder = D.police_marque(prs)  # 'Outfit' sur le template OCTO
    if placeholder and placeholder != theme:
        assert placeholder not in fonts, (
            f"la police des placeholders {placeholder!r} (non installée) est forcée sur "
            "le texte — elle serait rendue en substitution ; utiliser la police du thème"
        )


def test_deck_charte_pas_d_ombre_sur_les_cartes() -> None:
    """Charte (règle dure OCTO : différenciation par bordure, jamais d'ombre) — les
    autoshapes dessinées (cartes/rectangles) ont l'héritage d'ombre coupé (_no_shadow)."""
    from pptx.oxml.ns import qn
    db = SessionLocal()
    try:
        prs = build_presentation(db.get(Mission, _mission_complete()))
    finally:
        db.close()
    # Aucune forme dessinée ne doit porter un effet d'ombre explicite (<a:effectLst>
    # avec une ombre) — _no_shadow coupe l'héritage du thème sur add_card/add_rect.
    for s in prs.slides:
        for sh in s.shapes:
            spPr = getattr(sh._element, "spPr", None)
            if spPr is None:
                continue
            effect = spPr.find(qn("a:effectLst"))
            if effect is not None:
                assert effect.find(qn("a:outerShdw")) is None, "ombre portée détectée (charte OCTO)"


# ---- Revue de design automatisée (demande 2026-07-22, skill deck-design-review) ----
# Chaque invariant ci-dessous verrouille un défaut TROUVÉ par une revue visuelle
# réelle — le test empêche sa réapparition, le skill (rendu + œil) trouve les suivants.


def _prs_complete():
    db = SessionLocal()
    try:
        return build_presentation(db.get(Mission, _mission_complete()))
    finally:
        db.close()


def _slide_titre(s) -> str:
    t = s.shapes.title
    return (t.text_frame.text if t is not None and t.has_text_frame else "").strip()


def test_design_matrice_dessinee_pas_de_scatter_excel() -> None:
    """La matrice de priorisation est DESSINÉE (skill priority-matrix) : plus aucun
    graphique scatter natif (marqueurs Excel gris + légende ◆■▲ illisibles), et la
    slide porte ses invariants — 4 libellés de quadrant + une bulle par reco."""
    from pptx.oxml.ns import qn
    prs = _prs_complete()
    matrice = None
    for s in prs.slides:
        if "Matrice de priorisation" in _slide_titre(s):
            matrice = s
        for gf in s.shapes:
            if gf.has_chart:
                xml = gf.chart._chartSpace.xml
                assert "scatterChart" not in xml, "scatter Excel réintroduit (priority-matrix)"
    assert matrice is not None, "slide matrice de priorisation absente"
    textes = " ".join(sh.text_frame.text for sh in matrice.shapes if sh.has_text_frame)
    for lbl in ("QUICK WINS", "CHANTIERS DE FOND", "OPPORTUNISTES", "À DIFFÉRER"):
        assert lbl in textes, f"libellé de quadrant manquant : {lbl}"
    assert "1.1" in textes, "bulle de reco absente de la matrice"


def test_design_legende_matrice_porte_toutes_les_recos() -> None:
    """La légende encadrée de la matrice liste TOUTES les recos (shrink-to-fit :
    réduire la taille plutôt que droper). À taille fixe, l'estimation pessimiste
    des hauteurs s'accumulait sur 8 items à intitulés longs et la dernière reco
    (« 4.2 ») sautait silencieusement de la légende (vu au rendu réel 2026-07-22
    — sa bulle restait pourtant sur la matrice)."""
    titres_longs = [
        "Mettre en place une instance de gouvernance data décisionnelle",
        "Formaliser l'ownership des domaines et actifs de données",
        "Structurer un portefeuille data unique et priorisé",
        "Généraliser les squads mixtes sur les cas d'usage prioritaires",
        "Relancer le catalogue comme produit opérationnel",
        "Déployer un socle de self-service gouverné",
        "Lancer un programme ciblé de data literacy et de culture produit",
        "Intégrer la donnée aux rituels de pilotage métier",
    ]
    db = SessionLocal()
    try:
        m = Mission(name="Audit qualité — Légende matrice")
        db.add(m); db.flush()
        for i in range(4):
            ax = RecommendationAxis(mission_id=m.id, title=f"Axe {i + 1}", position=i)
            db.add(ax); db.flush()
            for j in range(2):
                db.add(Recommendation(
                    axis_id=ax.id, position=j, title=titres_longs[i * 2 + j],
                    valeur=5 - i, complexite=1 + i))
        db.commit()
        prs = build_presentation(db.get(Mission, m.id))
    finally:
        db.close()
    matrice = next((s for s in prs.slides if "Matrice de priorisation" in _slide_titre(s)), None)
    assert matrice is not None
    textes = " ".join(sh.text_frame.text for sh in matrice.shapes if sh.has_text_frame)
    for i in range(1, 5):
        for j in range(1, 3):
            idx = f"{i}.{j}"
            # L'index seul est la bulle ; « i.j  Intitulé » est l'item de légende.
            assert f"{idx}  {titres_longs[(i - 1) * 2 + (j - 1)][:20]}" in textes, (
                f"reco {idx} absente de la légende de la matrice (dropée)"
            )


def test_design_legende_matrice_jamais_droper_meme_sur_titres_extremes() -> None:
    """Régression (revue adversariale 2026-07-22) : l'ancienne boucle shrink
    `while t_leg > 7.5` sortait SANS avoir évalué le plancher 7.5 pt — sur des
    titres extrêmes le garde-fou du rendu dropait alors les dernières recos en
    silence, exactement le contrat que la légende prétend tenir. La cascade
    corrigée finit à 1 ligne/item tronquée : titre coupé, reco JAMAIS absente.
    12 recos à titres très longs : à 2 lignes/item RIEN ne tient même à 7.5 pt
    (l'ancien code dropait ~4 items) — seul le cran 1 ligne/item les loge toutes."""
    queue = " des données puis de l'organisation cible sur l'ensemble du périmètre" * 3
    db = SessionLocal()
    try:
        m = Mission(name="Audit qualité — Légende extrême")
        db.add(m); db.flush()
        for i in range(4):
            ax = RecommendationAxis(mission_id=m.id, title=f"Axe {i + 1}", position=i)
            db.add(ax); db.flush()
            for j in range(3):
                db.add(Recommendation(
                    axis_id=ax.id, position=j,
                    title=f"Chantier{i + 1}{j + 1}{queue}",
                    valeur=5 - i, complexite=1 + i))
        db.commit()
        prs = build_presentation(db.get(Mission, m.id))
    finally:
        db.close()
    matrice = next((s for s in prs.slides if "Matrice de priorisation" in _slide_titre(s)), None)
    assert matrice is not None
    textes = " ".join(sh.text_frame.text for sh in matrice.shapes if sh.has_text_frame)
    for i in range(1, 5):
        for j in range(1, 4):
            assert f"{i}.{j}  Chantier{i}{j}" in textes, (
                f"reco {i}.{j} dropée de la légende sur titres extrêmes"
            )


def test_verifier_debordements_estime_au_plus_grand_run() -> None:
    """Régression (revue adversariale 2026-07-22) : le vérificateur prenait la
    taille du PREMIER run stylé — un préfixe 8 pt devant un corps 16 pt faisait
    estimer tout le paragraphe à 8 pt (faux négatif, contraire au contrat
    pessimiste). Le paragraphe mixte ci-dessous tient à 8 pt mais pas à 16 :
    il DOIT être signalé."""
    from pptx import Presentation
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.3))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "1.1  "; r1.font.size = Pt(8)
    r2 = p.add_run()
    r2.text = "Un corps de texte volontairement long pour la calibration mixte."
    r2.font.size = Pt(16)
    problemes = D.verifier_debordements_texte(prs)
    assert problemes, (
        "paragraphe mixte estimé au 1er run (8 pt) — le débordement réel à 16 pt "
        "n'est pas détecté"
    )


def test_fiche_reco_pas_de_bandeau_pour_resultats_reduits_a_des_puces_vides() -> None:
    """Régression (revue adversariale 2026-07-22) : `resultats_attendus` réduit à
    des marqueurs de puce vides (« - ») réservait ET dessinait un bandeau
    « RÉSULTATS ATTENDUS — » vide de 0.72in, amputant d'autant les cartes. Le
    prédicat suit désormais le contenu réellement rendu (_bullet_lines)."""
    db = SessionLocal()
    try:
        m = Mission(name="Audit qualité — Bandeau vide")
        db.add(m); db.flush()
        ax = RecommendationAxis(mission_id=m.id, title="Axe 1", position=0)
        db.add(ax); db.flush()
        db.add(Recommendation(
            axis_id=ax.id, position=0, title="Reco sans résultats réels",
            objectif="O", acteurs="A", valeur=3, complexite=3,
            proposition_valeur="P", plan_actions="- a",
            resultats_attendus="- \n- "))
        db.commit()
        prs = build_presentation(db.get(Mission, m.id))
    finally:
        db.close()
    fiche = next((s for s in prs.slides if _slide_titre(s).startswith("1.1")), None)
    assert fiche is not None
    textes = " ".join(sh.text_frame.text for sh in fiche.shapes if sh.has_text_frame)
    assert "RÉSULTATS ATTENDUS" not in textes, (
        "bandeau résultats dessiné alors que le contenu se réduit à des puces vides"
    )


def test_design_encart_a_retenir_reduit_avant_de_tronquer() -> None:
    """Régression (batterie design 2026-07-22) : l'encart « à retenir » des slides
    de synthèse tronquait son claim en plein mot à h3 fixe/2 lignes (« …la DSI
    Groupe, le Marketing, la… ») — un so-what coupé ne dit plus rien. La police
    descend désormais h3→body→small avant l'ellipse : ce claim réaliste de
    ~110 caractères doit apparaître EN ENTIER."""
    claim = ("L'analyse couvre dix entretiens représentant la DSI Groupe, le "
             "Marketing, la Finance et les Opérations sur deux mois.")
    db = SessionLocal()
    try:
        m = Mission(name="Audit qualité — Encart à retenir")
        db.add(m); db.flush()
        db.add(Interview(mission_id=m.id, interviewee_name="Témoin", status="done"))
        db.add(GlobalSynthesis(
            mission_id=m.id, status="generated",
            contexte=f"- {claim}\n- Une seconde puce pour activer l'encart.",
            culture_adn="- C", forces_succes="- F",
            points_amelioration="- P", aspirations="- A"))
        db.commit()
        prs = build_presentation(db.get(Mission, m.id))
    finally:
        db.close()
    contexte = next((s for s in prs.slides if "Contexte" in _slide_titre(s)), None)
    assert contexte is not None
    textes = " ".join(sh.text_frame.text for sh in contexte.shapes if sh.has_text_frame)
    assert claim in textes, "claim de l'encart « à retenir » tronqué au lieu d'être réduit"


def test_design_bulles_matrice_jamais_superposees_meme_a_9_par_ligne() -> None:
    """Régression (defer revue adversariale corrigé) : sur une ligne de valeur
    saturée (9 recos de même valeur), le recalage `max(pl+0.02, x-depassement)`
    RE-SUPERPOSAIT toutes les bulles écrêtées au bord gauche — numéros masqués.
    La répartition uniforme garde tous les centres distincts."""
    from pptx.util import Emu as _Emu
    db = SessionLocal()
    try:
        m = Mission(name="Audit qualité — Ligne de bulles saturée")
        db.add(m); db.flush()
        ax = RecommendationAxis(mission_id=m.id, title="Axe 1", position=0)
        db.add(ax); db.flush()
        for j in range(9):
            db.add(Recommendation(
                axis_id=ax.id, position=j, title=f"Reco {j + 1}",
                valeur=5, complexite=5))
        db.commit()
        prs = build_presentation(db.get(Mission, m.id))
    finally:
        db.close()
    matrice = next((s for s in prs.slides if "Matrice de priorisation" in _slide_titre(s)), None)
    assert matrice is not None
    # Les bulles sont les formes au texte « 1.N » (badges ronds de la matrice).
    import re
    lefts = sorted(
        _Emu(sh.left).inches
        for sh in matrice.shapes
        if sh.has_text_frame and re.fullmatch(r"1\.\d", sh.text_frame.text.strip())
    )
    assert len(lefts) == 9, f"9 bulles attendues, {len(lefts)} trouvées"
    ecarts = [b - a for a, b in zip(lefts, lefts[1:])]
    assert min(ecarts) > 0.25, (
        f"bulles superposées sur ligne saturée (écart min {min(ecarts):.2f}in)"
    )


def test_emit_bullet_overflow_termine_sur_puce_insecable_geante() -> None:
    """Régression (revue adversariale 2026-07-22, boucle infinie PROUVÉE : 1726
    slides en 15 s) : une puce SANS retour à la ligne plus haute qu'une slide de
    suite entière revenait intégralement en overflow à chaque itération — export
    pendu, mémoire non bornée. La garde de progression coupe : une seule slide de
    suite, contenu tronqué à l'ellipse. Sur l'ancien code ce test ne terminait
    pas (vérifié par le script de preuve de la revue, pas rejoué ici)."""
    db = SessionLocal()
    try:
        m = Mission(name="Audit qualité — Puce géante")
        db.add(m); db.flush()
        ax = RecommendationAxis(mission_id=m.id, title="Axe 1", position=0)
        db.add(ax); db.flush()
        db.add(Recommendation(
            axis_id=ax.id, position=0, title="Reco au plan insécable",
            objectif="O", acteurs="A", valeur=3, complexite=3,
            proposition_valeur="P",
            plan_actions="Une action unique interminable " * 120,  # ~3 700 car., aucun \n
            resultats_attendus="- r"))
        db.commit()
        prs = build_presentation(db.get(Mission, m.id))
    finally:
        db.close()
    suites = [s for s in prs.slides if "(suite" in _slide_titre(s)]
    assert len(suites) <= 2, (
        f"{len(suites)} slides de suite pour UNE puce insécable — la garde de "
        "progression de _emit_bullet_overflow ne coupe plus"
    )


def test_fiche_reco_realiste_sans_slide_de_suite() -> None:
    """R4 sur le correctif principal de la refonte fiche (revue adversariale) :
    la motivation « plus de slide de suite systématique » (25 slides au lieu de
    33) doit être un invariant testé — une fiche à volumes RÉALISTES (objectif
    3 lignes, 4 puces de plan, résultats ~150 car.) tient sur UNE slide."""
    db = SessionLocal()
    try:
        m = Mission(name="Audit qualité — Fiche réaliste")
        db.add(m); db.flush()
        ax = RecommendationAxis(mission_id=m.id, title="Axe 1", position=0)
        db.add(ax); db.flush()
        db.add(Recommendation(
            axis_id=ax.id, position=0,
            title="Mettre en place une instance de gouvernance data décisionnelle",
            objectif="Créer un mécanisme régulier et légitime pour arbitrer les "
                     "priorités, valider les règles communes et rendre les décisions "
                     "data visibles de tous.",
            acteurs="CDO, DSI, représentants métiers et data office",
            valeur=5, complexite=3,
            proposition_valeur="Réduire les délais de décision, rendre les arbitrages "
                               "transparents et concentrer les capacités sur les cas "
                               "d'usage les plus porteurs.",
            plan_actions="- Définir le mandat et les droits de décision\n"
                         "- Nommer les membres permanents\n"
                         "- Établir des critères communs de valeur, risque et effort\n"
                         "- Organiser un comité mensuel et publier les décisions",
            resultats_attendus="- Priorités partagées, diminution des escalades au "
                               "COMEX, meilleure allocation des ressources et "
                               "visibilité accrue sur le portefeuille data."))
        db.commit()
        prs = build_presentation(db.get(Mission, m.id))
    finally:
        db.close()
    suites = [s for s in prs.slides if "(suite" in _slide_titre(s)]
    assert suites == [], "une fiche à volumes réalistes repart en slide de suite"


def test_design_swot_matrice_axes_explicites() -> None:
    """La SWOT reste une MATRICE (skill swot-matrix) : axes Interne/Externe ×
    Favorable/Défavorable présents — pas quatre cartes flottantes."""
    prs = _prs_complete()
    swot = next((s for s in prs.slides if "SWOT" in _slide_titre(s)), None)
    assert swot is not None
    textes = " ".join(sh.text_frame.text for sh in swot.shapes if sh.has_text_frame)
    for lbl in ("FAVORABLE", "DÉFAVORABLE", "INTERNE", "EXTERNE"):
        assert lbl in textes, f"axe de la matrice SWOT manquant : {lbl}"


def test_design_fiche_reco_encarts_arrondis() -> None:
    """Les fiches reco posent leur contenu dans des ENCARTS ARRONDIS format OCTO
    (demande 2026-07-22) : cartes roundRect + l'encart gris « proposition »."""
    prs = _prs_complete()
    fiche = next((s for s in prs.slides if _slide_titre(s).startswith("1.1")), None)
    assert fiche is not None, "fiche 1.1 absente"
    xml = fiche._element.xml
    assert xml.count('prst="roundRect"') >= 3, "cartes arrondies manquantes sur la fiche"
    assert D.ENCART_BG.lstrip("#").upper() in xml.upper(), "encart gris proposition absent"


def test_design_titres_a_l_echelle() -> None:
    """Tous les titres de slides de contenu sont à D.TYPE['title'] (20pt, aligné
    référence) — pas de taille de titre ad hoc réintroduite. Un titre long peut
    être RÉDUIT (jamais tronqué, 2026-07-23) mais reste borné entre le plancher
    h3 et l'échelle : jamais au-dessus, jamais illisible."""
    from pptx.util import Pt
    prs = _prs_complete()
    for s in prs.slides:
        lay = (s.slide_layout.name or "").lower()
        if "titre seul" not in lay:
            continue  # couverture/chapitres : tailles du template de marque
        t = s.shapes.title
        if t is None or not t.has_text_frame:
            continue
        for p in t.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size is not None:
                    assert Pt(D.TYPE["h3"]) <= r.font.size <= Pt(D.TYPE["title"]), (
                        f"titre hors échelle sur « {_slide_titre(s)} » : {r.font.size}"
                    )


def test_design_titres_jamais_tronques() -> None:
    """Aucun titre de slide n'est tronqué à l'ellipse (demande 2026-07-23) : un
    titre long est rendu en police réduite (et replie au besoin), jamais coupé.
    Aurait échoué avant : la fiche reco tronquait « N — titre » à 1 ligne 20pt,
    et _new_slide coupait tout titre au-delà de 2 lignes."""
    titre_long = (
        "Mettre en place une gouvernance des données décisionnelle, outillée et "
        "partagée entre les directions métiers, la DSI et les équipes terrain"
    )
    db = SessionLocal()
    try:
        m = db.get(Mission, _mission_complete())
        m.recommendation_axes[0].recommendations[0].title = titre_long
        db.commit()
        prs = build_presentation(m)
    finally:
        db.close()
    # Le titre complet apparaît tel quel sur la fiche (aucune ellipse).
    fiche = next((s for s in prs.slides if _slide_titre(s).startswith("1.1")), None)
    assert fiche is not None, "fiche 1.1 absente"
    assert titre_long in _slide_titre(fiche), "titre de fiche coupé au lieu d'être réduit"
    for s in prs.slides:
        titre = _slide_titre(s)
        assert not titre.endswith("…"), f"titre de slide tronqué : « {titre} »"


def test_design_chips_fiche_ne_chevauchent_jamais_le_label_criteres() -> None:
    """Sur une fiche raccourcie (titre de slide sur 2-3 lignes depuis le
    non-tronquage 2026-07-23), les chips Valeur/Complexité ne peignent JAMAIS
    par-dessus le label « CRITÈRES DE PRIORISATION » (constat rendu réel) : soit
    le bloc label+chips tient, soit le label saute — pas de chevauchement."""
    titre_long = (
        "Mettre en place une gouvernance des données décisionnelle, outillée et "
        "partagée entre les directions métiers, la DSI et les équipes terrain"
    )
    db = SessionLocal()
    try:
        m = db.get(Mission, _mission_complete())
        m.recommendation_axes[0].recommendations[0].title = titre_long
        db.commit()
        prs = build_presentation(m)
    finally:
        db.close()
    fiche = next((s for s in prs.slides if _slide_titre(s).startswith("1.1")), None)
    assert fiche is not None, "fiche 1.1 absente"
    label = next((sh for sh in fiche.shapes if sh.has_text_frame
                  and sh.text_frame.text.strip() == "CRITÈRES DE PRIORISATION"), None)
    chips = [sh for sh in fiche.shapes if sh.has_text_frame
             and sh.text_frame.text.strip().startswith(("Valeur ", "Complexité "))]
    assert chips, "chips Valeur/Complexité absents de la fiche"
    if label is not None:
        label_bas = label.top + label.height
        for chip in chips:
            assert chip.top >= label_bas, (
                "chip peint par-dessus le label CRITÈRES DE PRIORISATION"
            )


def test_design_le_texte_tient_dans_sa_boite() -> None:
    """Filet « texte hors cadre » (demande utilisateur 2026-07-22, revue de la
    couverture : le défaut visuel récurrent — texte qui peint par-dessus le bord
    d'une carte — n'était couvert par AUCUN test ; verifier_geometrie ne voit que
    les bords des formes). verifier_debordements_texte estime le contenu de
    chaque zone dessinée avec une calibration PESSIMISTE et doit rester à zéro."""
    prs = _prs_complete()
    problemes = D.verifier_debordements_texte(prs)
    assert problemes == [], "texte dépassant sa boîte détecté :\n" + "\n".join(problemes)


def test_design_cache_image_proc_ne_bloque_pas_la_photo(tmp_path, monkeypatch) -> None:
    """Le repli procédural (tests/offline) n'occupe JAMAIS le slot photo : un run
    en ligne ultérieur doit retenter le vrai fetch (cause racine des images
    « générées » servies à vie sur les slides synthèse, 2026-07-22)."""
    # Les internes image vivent dans le sous-module `images` depuis le découpage
    # du gros module pptx_export (audit 2026-07-24) — le monkeypatch de
    # _IMG_CACHE doit viser le module où _resoudre_image_cachee lit sa globale.
    import app.services.pptx_export.images as px
    if not px._FRAMED_OK:  # infra image absente : rien à vérifier ici
        return
    monkeypatch.setattr(px, "_IMG_CACHE", tmp_path)
    monkeypatch.setenv("PPTX_NO_PHOTO_FETCH", "1")
    p = px._resoudre_image_cachee("t_1_10x10", "mountains", 1, 1.0, 10, 10, requete="x")
    assert p.name.endswith("_proc.png") and p.exists()
    assert not list(tmp_path.glob("*_photo.png")), (
        "le repli procédural a écrit dans le slot photo — il le monopoliserait à vie"
    )


def test_design_visuel_de_synthese_suit_la_key_pas_le_libelle() -> None:
    """Invariant R4 du défaut trouvé par le diagnostic superviseur du 2026-07-28.

    Depuis que les axes d'étude sont configurables (2026-07-27), le libellé est
    RENOMMABLE et seule la `key` est stable — mais le visuel de la slide de synthèse
    restait indexé par libellé, avec un repli constant. Deux conséquences visibles au
    rendu : un axe renommé perdait son image métier, et plusieurs axes sur mesure
    recevaient LA MÊME photo (même scène, même seed), ce qui se lit comme un
    décor et non comme du sens."""
    from app.services.pptx_export.slides_diagnostic import _SYNTHESE_VISUEL, _visuel_axe

    # 1. Un axe RENOMMÉ garde le visuel de sa key (le libellé ne compte plus).
    assert _visuel_axe("contexte") == _SYNTHESE_VISUEL["contexte"]

    # 2. Deux axes SUR MESURE ne partagent ni (scène, seed) ni requête photo. La scène
    #    seule peut se répéter (il n'y en a que six), mais la requête porte alors un mot
    #    tiré de la key — sinon deux axes de même scène tapent la même page de résultats
    #    et peuvent recevoir le MÊME cliché (`search_photo` indexe `seed % len(results)`).
    sur_mesure = ["outillage_donnees", "gouvernance", "relation_client", "delais",
                  "securite", "formation"]
    visuels = [_visuel_axe(k) for k in sur_mesure]
    assert len({(v[0], v[2]) for v in visuels}) == len(sur_mesure), (
        f"axes sur mesure au visuel identique : {visuels}"
    )
    assert len({v[1] for v in visuels}) == len(sur_mesure), (
        f"axes sur mesure à la même requête photo : {[v[1] for v in visuels]}"
    )

    # 3. Sur un large échantillon, la scène n'est PAS déterminée par le seed (tranches
    #    de hash distinctes). Quand elle l'était (une seule empreinte, `6 | 900`),
    #    l'espace des signatures s'effondrait à 900 et 3000 clés ne donnaient que
    #    853 visuels distincts ; l'espace réel est désormais 6×997, dont ~2350
    #    signatures distinctes attendues sur 3000 tirages (paradoxe des anniversaires).
    larges = [_visuel_axe(f"axe_{i}") for i in range(3000)]
    assert len({(v[0], v[2]) for v in larges}) >= 2000

    # 4. Reproductible d'un export à l'autre (pas de hash randomisé par processus).
    assert visuels == [_visuel_axe(k) for k in sur_mesure]


def test_design_deck_a_axes_sur_mesure_ne_repete_pas_la_meme_image() -> None:
    """Le couplage COMPLET, de `build_presentation` à l'image posée — l'invariant
    précédent n'exerce que `_visuel_axe` en isolation, si bien qu'un retour à l'appel
    à trois arguments (le bug d'origine) le laisserait vert. C'est aussi la seule
    couverture du paramètre `axes_etude`, jusqu'ici exercé par aucun test.

    On rend une mission à axes SUR MESURE et on compare les images réellement
    embarquées : avant le correctif, toutes les slides de synthèse partageaient le
    même repli, donc le même blob."""
    from app.services.pptx_export import build_presentation

    class _Axe:
        def __init__(self, key, label):
            self.key, self.label, self.hint = key, label, ""

    axes = [
        _Axe("contexte", "Contexte & historique"),      # renommé, key stable
        _Axe("outillage_donnees", "Outillage & données"),  # sur mesure
        _Axe("relation_client", "Relation client"),        # sur mesure
    ]
    db = SessionLocal()
    mission = db.query(Mission).first()
    gs = db.query(GlobalSynthesis).filter_by(mission_id=mission.id).first()
    for axe in axes:
        gs.set_contenu(axe.key, f"{axe.label}\n- Un constat.\n- Un autre constat.")
    db.flush()  # visible pour le build, jamais commité : la base du module reste intacte

    prs = build_presentation(
        mission, include_sommaire=False, include_executive_summary=False,
        include_synthese=True, include_difficultes=False, include_swot=False,
        include_verbatims=False, include_axes_overview=False, include_matrix=False,
        include_axis_ids=set(), axes_etude=axes,
    )
    db.rollback()
    db.close()

    # Les slides de synthèse SEULES (le deck émet aussi couverture et intercalaire,
    # qui portent leurs propres images).
    synthese = [
        slide for slide in prs.slides
        if any(sh.has_text_frame and sh.text_frame.text.startswith("Synthèse globale —")
               for sh in slide.shapes)
    ]
    assert len(synthese) == len(axes), f"{len(synthese)} slide(s) pour {len(axes)} axes"
    images = [
        pic.image.blob
        for slide in synthese
        for pic in slide.shapes
        if pic.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
    ]
    import app.services.pptx_export.images as px

    if not px._FRAMED_OK:  # infra image absente : rien à vérifier ici
        return
    assert len(images) == len(axes), f"{len(images)} image(s) pour {len(axes)} axes"
    assert len(set(images)) == len(images), (
        "deux slides de synthèse portent la MÊME image — le visuel est retombé "
        "sur un repli commun (régression du couplage key → visuel)"
    )


def test_design_sommaire_ne_promet_pas_une_synthese_sans_slide() -> None:
    """Parité sommaire/PPT (trouvée en revue adversariale) : `gs.has_content` répond sur
    `valeurs` ET les 5 colonnes historiques, donc reste vrai pour un axe SUPPRIMÉ depuis.
    Le sommaire annonçait alors « Synthèse globale » et l'intercalaire « Le diagnostic »
    s'ouvrait sur zéro slide — même classe de défaut que celui corrigé le 2026-07-22 sur
    les difficultés."""
    from app.services.pptx_export import build_presentation

    class _Axe:
        def __init__(self, key, label):
            self.key, self.label, self.hint = key, label, ""

    db = SessionLocal()
    mission = db.query(Mission).first()
    gs = db.query(GlobalSynthesis).filter_by(mission_id=mission.id).first()
    gs.set_contenu("contexte", "Du contenu sur un axe que la mission n'étudie plus.")
    db.flush()
    # La mission n'étudie plus QUE cet axe-là : sa matière est hors périmètre.
    prs = build_presentation(
        mission, include_sommaire=True, include_executive_summary=False,
        include_synthese=True, include_difficultes=False, include_swot=False,
        include_verbatims=False, include_axes_overview=False, include_matrix=False,
        include_axis_ids=set(), axes_etude=[_Axe("axe_neuf", "Axe neuf")],
    )
    db.rollback()
    db.close()

    textes = [
        sh.text_frame.text for slide in prs.slides for sh in slide.shapes
        if sh.has_text_frame
    ]
    assert not any(t.startswith("Synthèse globale —") for t in textes)  # aucune slide…
    assert not any("Synthèse globale" in t for t in textes), (  # …donc pas au sommaire
        "le sommaire annonce une section sans la moindre slide"
    )


def test_design_matrice_aucune_bulle_sur_la_frontiere_des_quadrants() -> None:
    """Invariant R4 (rendu 2026-07-28) : la règle linéaire `(s-0.5)/5` posait le score 3
    EXACTEMENT sur la frontière des quadrants — la bulle chevauchait deux quadrants de
    sens opposé (« OPPORTUNISTES » / « QUICK WINS ») et recouvrait leur libellé. Le
    découpage par moitié (1-3 en bas, 4-5 en haut) garantit qu'aucune bulle ne touche la
    ligne médiane ni la bande des libellés."""
    from app.services.pptx_export.slides_trajectoire import _fraction_score

    demi_bulle = 0.46 / 2 / 3.95  # rayon de bulle rapporté à la largeur de plot
    for score in (1, 2, 3, 4, 5):
        f = _fraction_score(score)
        assert abs(f - 0.5) > demi_bulle, (
            f"score {score} : bulle à cheval sur la frontière des quadrants (f={f:.3f})"
        )
        assert demi_bulle < f < 1 - demi_bulle, f"score {score} hors zone de tracé"
        assert (f < 0.5) == (score <= 3), f"score {score} du mauvais côté de la médiane"

    # Bande réservée aux libellés en tête de CHAQUE moitié : ni le score 5 (haut du
    # graphe) ni le score 3 (haut de la moitié basse) ne doit y entrer.
    bande = 0.38 / 4.0
    for score, plafond in ((3, 0.5), (5, 1.0)):
        f = _fraction_score(score, bande)
        assert f + demi_bulle < plafond - bande, (
            f"score {score} : la bulle mord sur le libellé de quadrant"
        )

    # Ordre strictement croissant conservé : la matrice reste lisible comme un
    # classement, on n'a pas seulement écarté les bulles de la ligne.
    fractions = [_fraction_score(s) for s in range(1, 6)]
    assert fractions == sorted(fractions)


def test_design_matrice_legende_dimensionnee_au_contenu() -> None:
    """Invariant R4 : la carte de légende occupait TOUTE la hauteur du graphe quel que
    soit le nombre de recos — trois lignes en haut d'un grand rectangle vide, défaut
    « panneau étiré » de la checklist pptx-verify."""
    from pptx.util import Emu

    from app.services.pptx_export import build_presentation

    db = SessionLocal()
    mission = db.query(Mission).first()
    prs = build_presentation(
        mission, include_sommaire=False, include_executive_summary=False,
        include_synthese=False, include_difficultes=False, include_swot=False,
        include_verbatims=False, include_axes_overview=False, include_matrix=True,
        include_axis_ids=set(),
    )
    db.close()

    matrice = [
        slide for slide in prs.slides
        if any(sh.has_text_frame and sh.text_frame.text.startswith("Matrice de priorisation")
               for sh in slide.shapes)
    ]
    assert len(matrice) == 1
    # La carte de légende est la forme la plus à droite ; elle ne doit pas descendre
    # jusqu'au bas de la zone de tracé (le label d'axe X, lui, y reste).
    formes = [sh for sh in matrice[0].shapes if sh.width and sh.left]
    droite = max(formes, key=lambda sh: sh.left)
    hauteur_in = Emu(droite.height).inches
    assert hauteur_in < 3.5, (
        f"la légende occupe {hauteur_in:.2f}in de haut pour quelques recos — "
        "carte étirée au lieu d'être dimensionnée au contenu"
    )
