import io
import json
import os
import re
import subprocess
from pathlib import Path

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Emu, Inches

from app.main import app
from app.db import DB_PATH, SessionLocal, engine, init_db
from app.importers.docx_trame import parse_docx_bytes
from app.models import Answer, Interview, Mission, Question, Recommendation, RecommendationAxis, Theme, Trame
from app.routers.trames import _parsed_to_json
from app.services import audio_transcribe
from app.services.openhub_agents import invoke_skill
from sqlalchemy import select


def setup_module() -> None:
    # Ensure test DB is isolated and fresh.
    if DB_PATH.exists():
        DB_PATH.unlink()
    init_db()


def teardown_module() -> None:
    try:
        engine.dispose()
    except Exception:
        pass
    if DB_PATH.exists():
        DB_PATH.unlink()


@pytest.fixture
def client() -> TestClient:
    return TestClient(app)


def test_create_mission_and_view_trame(client: TestClient) -> None:
    response = client.post(
        "/missions", data={
            "name": "Mission Test",
            "description": "Description de test",
        },
        follow_redirects=False,
    )
    assert response.status_code == 303
    assert "/missions/" in response.headers["location"]

    mission_url = response.headers["location"]
    mission_id = mission_url.rsplit("/", 1)[-1]
    assert mission_id.isdigit()

    response = client.get(f"/missions/{mission_id}/trame")
    assert response.status_code == 200
    assert "Trame d'entretien" in response.text
    assert "Mission Test" in response.text
    assert "/missions/{mission_id}/trame/themes" in response.text or "Ajouter un thème" in response.text


def test_add_theme_and_question(client: TestClient) -> None:
    response = client.post(
        "/missions", data={
            "name": "Mission Test 2",
            "description": "Description de test 2",
        },
        follow_redirects=False,
    )
    assert response.status_code == 303
    mission_id = response.headers["location"].rsplit("/", 1)[-1]

    response = client.post(
        f"/missions/{mission_id}/trame/themes",
        data={"title": "Thème Test"},
        follow_redirects=False,
    )
    assert response.status_code == 303

    session = SessionLocal()
    try:
        theme = session.scalars(
            select(Theme).where(Theme.title == "Thème Test")
        ).one()
    finally:
        session.close()

    response = client.post(
        f"/missions/{mission_id}/trame/themes/{theme.id}/questions",
        data={
            "label": "Question de test ?",
            "qtype": "open",
            "scale_min": "1",
            "scale_max": "5",
            "choices": "",
            "help_text": "Texte d'aide",
        },
        follow_redirects=False,
    )
    assert response.status_code == 303

    response = client.get(f"/missions/{mission_id}/trame")
    assert response.status_code == 200
    assert "Thème Test" in response.text
    assert "Question de test ?" in response.text
    assert "Texte d'aide" in response.text


def _build_trame_docx(with_intro: bool = True) -> bytes:
    """Un .docx simple avec un thème et deux questions (convention OCTO)."""
    from docx import Document
    import io

    document = Document()
    if with_intro:
        document.add_heading("Objectifs et principes", level=1)
        document.add_paragraph("Introduction DOCX")
    document.add_heading("Thème importé", level=1)
    document.add_paragraph("Question A ?")
    document.add_paragraph("Question B ?")

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def test_import_docx_creates_theme_and_questions(client: TestClient) -> None:
    response = client.post(
        "/missions",
        data={
            "name": "Mission DOCX",
            "description": "Description DOCX",
        },
        follow_redirects=False,
    )
    assert response.status_code == 303
    mission_id = response.headers["location"].rsplit("/", 1)[-1]

    docx_bytes = _build_trame_docx()

    # Étape 1 : l'import ne fait plus que proposer une revue — rien n'est
    # encore écrit en base.
    response = client.post(
        f"/missions/{mission_id}/trame/import",
        files={"file": ("import.docx", docx_bytes, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
        follow_redirects=False,
    )
    assert response.status_code == 200
    assert "Revue avant import" in response.text

    response = client.get(f"/missions/{mission_id}/trame")
    assert "Thème importé" not in response.text

    # Étape 2 : validation -> fusion effective.
    parsed = parse_docx_bytes(docx_bytes)
    keep = [
        f"{ti}-{qi}"
        for ti, theme in enumerate(parsed.themes)
        for qi in range(len(theme.questions))
    ]
    response = client.post(
        f"/missions/{mission_id}/trame/import/confirm",
        data={"parsed": _parsed_to_json(parsed), "keep": keep},
        follow_redirects=False,
    )
    assert response.status_code == 303
    assert response.headers["location"].endswith(f"/missions/{mission_id}/trame/preview")

    response = client.get(f"/missions/{mission_id}/trame")
    assert response.status_code == 200
    assert "Thème importé" in response.text
    assert "Question A ?" in response.text
    assert "Question B ?" in response.text


def test_trame_import_review_can_exclude_a_question(client: TestClient) -> None:
    response = client.post(
        "/missions",
        data={
            "name": "Mission DOCX Exclude",
            "description": "Description DOCX Exclude",
        },
        follow_redirects=False,
    )
    assert response.status_code == 303
    mission_id = response.headers["location"].rsplit("/", 1)[-1]

    # Sans intro : évite que le texte "Question B ?" ne réapparaisse dans le
    # bloc d'introduction repris tel quel (heuristique de `extract_intro`),
    # ce qui fausserait l'assertion d'exclusion ci-dessous.
    docx_bytes = _build_trame_docx(with_intro=False)
    parsed = parse_docx_bytes(docx_bytes)

    # Ne garde que "Question A ?" (theme 0, question 0) -> "Question B ?" exclue.
    response = client.post(
        f"/missions/{mission_id}/trame/import/confirm",
        data={"parsed": _parsed_to_json(parsed), "keep": "0-0"},
        follow_redirects=False,
    )
    assert response.status_code == 303

    response = client.get(f"/missions/{mission_id}/trame")
    assert "Question A ?" in response.text
    assert "Question B ?" not in response.text


def test_agents_page_lists_available_skills(client: TestClient) -> None:
    response = client.post(
        "/missions",
        data={
            "name": "Mission Skills",
            "description": "Description Skills",
        },
        follow_redirects=False,
    )
    assert response.status_code == 303
    mission_id = response.headers["location"].rsplit("/", 1)[-1]

    response = client.get(f"/missions/{mission_id}/agents")
    assert response.status_code == 200
    assert "Skills disponibles" in response.text
    assert "beads-dev" in response.text or "Workflow exécuteur Beads" in response.text


def test_skill_invocation_from_agents_page(client: TestClient) -> None:
    response = client.post(
        "/missions",
        data={
            "name": "Mission Skill Invoke",
            "description": "Description Skill Invoke",
        },
        follow_redirects=False,
    )
    assert response.status_code == 303
    mission_id = response.headers["location"].rsplit("/", 1)[-1]

    response = client.post(
        f"/missions/{mission_id}/skills/beads-dev/invoke",
        follow_redirects=False,
    )
    assert response.status_code == 200
    assert "Skill OpenHub" in response.text
    assert "beads-dev" in response.text


def test_dynamic_skill_execution_uses_opencode_runtime(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    python_path = tmp_path / "opencode.py"
    python_path.write_text(
        "import sys\n"
        "print('skill-ok')\n"
        "sys.exit(0)\n",
        encoding="utf-8",
    )

    wrapper_path = tmp_path / "opencode.cmd"
    wrapper_path.write_text(
        "@echo off\n"
        f"python \"{python_path}\" %*\n",
        encoding="utf-8",
    )

    monkeypatch.setenv("PATH", f"{tmp_path}{os.pathsep}{os.environ.get('PATH', '')}")

    mission = type("Mission", (), {"name": "Mission", "description": "Desc", "trame": None, "interviews": []})()
    result = invoke_skill("beads-dev", mission)

    assert result["runtime_available"] is True
    assert "Skill OpenHub" in result["output"]
    assert "skill-ok" in result["output"]


def test_interview_capture_and_save_answer(client: TestClient) -> None:
    response = client.post(
        "/missions",
        data={
            "name": "Mission Interview",
            "description": "Description Interview",
        },
        follow_redirects=False,
    )
    assert response.status_code == 303
    mission_id = response.headers["location"].rsplit("/", 1)[-1]

    response = client.post(
        f"/missions/{mission_id}/trame/themes",
        data={"title": "Interview Thème"},
        follow_redirects=False,
    )
    assert response.status_code == 303

    session = SessionLocal()
    try:
        theme = session.scalars(
            select(Theme).where(Theme.title == "Interview Thème")
        ).one()
    finally:
        session.close()

    response = client.post(
        f"/missions/{mission_id}/trame/themes/{theme.id}/questions",
        data={
            "label": "Question interview ?",
            "qtype": "open",
            "scale_min": "1",
            "scale_max": "5",
            "choices": "",
            "help_text": "Note interview",
        },
        follow_redirects=False,
    )
    assert response.status_code == 303

    # Créer un entretien et vérifier la saisie d'une réponse.
    response = client.post(
        f"/missions/{mission_id}/interviews",
        data={
            "interviewee_name": "Jean Dupont",
            "interviewee_role": "Responsable",
            "interviewee_entity": "Produit",
            "interview_date": "2026-06-30",
            "reference_text": "Référence d'entretien",
        },
        follow_redirects=False,
    )
    assert response.status_code == 303
    interview_id = response.headers["location"].rsplit("/", 1)[-1]

    session = SessionLocal()
    try:
        question = session.scalars(
            select(Question).where(
                Question.label == "Question interview ?",
                Question.theme_id == theme.id,
            )
        ).one()
    finally:
        session.close()

    response = client.post(
        f"/interviews/{interview_id}/answers/{question.id}",
        data={"text": "Réponse test"},
        follow_redirects=False,
    )
    assert response.status_code == 200
    assert "Répondue" in response.text or "enr." in response.text

    response = client.get(f"/interviews/{interview_id}")
    assert response.status_code == 200
    assert "Question interview ?" in response.text
    assert "Réponse test" in response.text


def test_interview_preview_shows_all_themes_answers_verbatims_and_notes(
    client: TestClient,
) -> None:
    """US2 (évol) — la page /interviews/{id}/preview doit rester lisible en
    un coup d'œil quel que soit le type de question, le statut de réponse,
    la présence de verbatims/notes libres, et fonctionner même sans thème."""
    response = client.post("/missions", data={"name": "Mission Preview"}, follow_redirects=False)
    mission_id = response.headers["location"].rsplit("/", 1)[-1]

    client.post(f"/missions/{mission_id}/trame/themes", data={"title": "Organisation"}, follow_redirects=False)
    client.post(f"/missions/{mission_id}/trame/themes", data={"title": "Culture"}, follow_redirects=False)

    session = SessionLocal()
    try:
        theme_org = session.scalars(select(Theme).where(Theme.title == "Organisation")).one()
        theme_culture = session.scalars(select(Theme).where(Theme.title == "Culture")).one()
        theme_org_id, theme_culture_id = theme_org.id, theme_culture.id
    finally:
        session.close()

    client.post(
        f"/missions/{mission_id}/trame/themes/{theme_org_id}/questions",
        data={"label": "Question ouverte ?", "qtype": "open"},
        follow_redirects=False,
    )
    client.post(
        f"/missions/{mission_id}/trame/themes/{theme_org_id}/questions",
        data={"label": "Question échelle ?", "qtype": "scale", "scale_min": "1", "scale_max": "5"},
        follow_redirects=False,
    )
    client.post(
        f"/missions/{mission_id}/trame/themes/{theme_culture_id}/questions",
        data={"label": "Question choix ?", "qtype": "choice", "choices": "Collaboratif\nHiérarchique"},
        follow_redirects=False,
    )
    client.post(
        f"/missions/{mission_id}/trame/themes/{theme_culture_id}/questions",
        data={"label": "Question non répondue ?", "qtype": "open"},
        follow_redirects=False,
    )

    session = SessionLocal()
    try:
        q_open = session.scalars(select(Question).where(Question.label == "Question ouverte ?")).one()
        q_scale = session.scalars(select(Question).where(Question.label == "Question échelle ?")).one()
        q_choice = session.scalars(select(Question).where(Question.label == "Question choix ?")).one()
        q_unanswered = session.scalars(select(Question).where(Question.label == "Question non répondue ?")).one()
        q_open_id, q_scale_id, q_choice_id, q_unanswered_id = q_open.id, q_scale.id, q_choice.id, q_unanswered.id
    finally:
        session.close()

    response = client.post(
        f"/missions/{mission_id}/interviews",
        data={"interviewee_name": "Marie Dupont", "interviewee_role": "DRH"},
        follow_redirects=False,
    )
    interview_id = response.headers["location"].rsplit("/", 1)[-1]

    client.post(f"/interviews/{interview_id}/answers/{q_open_id}", data={"text": "Réponse ouverte détaillée."})
    client.post(f"/interviews/{interview_id}/answers/{q_scale_id}", data={"value": "4"})
    client.post(f"/interviews/{interview_id}/answers/{q_choice_id}", data={"value": "Collaboratif"})
    client.post(f"/interviews/{interview_id}/answers/{q_unanswered_id}/status", data={"status": "revisit"})
    # Pas d'apostrophe dans la citation : l'autoescape Jinja la rendrait en
    # `&#39;` dans le HTML source (correct à l'affichage navigateur, mais
    # `in response.text` compare le HTML brut, pas le rendu décodé).
    client.post(f"/interviews/{interview_id}/verbatims/{q_open_id}", data={"quote": "On adapte tout en continu."})
    client.post(f"/interviews/{interview_id}/notes", data={"free_notes": "Note libre hors trame."})

    response = client.get(f"/interviews/{interview_id}/preview")
    assert response.status_code == 200
    text = response.text

    # Identité + couverture.
    assert "Marie Dupont" in text
    assert "3/4" in text

    # Les deux thèmes et les 4 questions apparaissent.
    assert "Organisation" in text
    assert "Culture" in text
    assert "Question ouverte ?" in text
    assert "Question échelle ?" in text
    assert "Question choix ?" in text
    assert "Question non répondue ?" in text

    # Chaque type de réponse est rendu (texte, échelle, choix).
    assert "Réponse ouverte détaillée." in text
    assert ">4<" in text
    assert "Collaboratif" in text

    # Question non répondue : badge de statut + pas de contenu fantôme.
    assert "À revoir" in text
    assert "sans réponse" in text

    # Verbatim et notes libres visibles.
    assert "On adapte tout en continu." in text
    assert "Note libre hors trame." in text

    # Pas de bouton/champ de saisie sur cette page en lecture seule.
    assert "<textarea" not in text
    assert "hx-post" not in text


def test_interview_preview_without_trame_or_answers_does_not_crash(client: TestClient) -> None:
    """Garde-fou : ni une trame vide, ni un entretien sans aucune réponse, ne
    doivent faire planter la page (juste un aperçu vide/à sans-réponse)."""
    response = client.post("/missions", data={"name": "Mission Preview Vide"}, follow_redirects=False)
    mission_id = response.headers["location"].rsplit("/", 1)[-1]

    # Entretien créé alors que la trame est encore vide (aucun thème).
    response = client.post(
        f"/missions/{mission_id}/interviews",
        data={"interviewee_name": "Sans Thème"},
        follow_redirects=False,
    )
    interview_id = response.headers["location"].rsplit("/", 1)[-1]

    response = client.get(f"/interviews/{interview_id}/preview")
    assert response.status_code == 200
    assert "Sans Thème" in response.text

    # Puis avec un thème/question mais aucune réponse saisie.
    client.post(f"/missions/{mission_id}/trame/themes", data={"title": "Thème Vide"}, follow_redirects=False)
    session = SessionLocal()
    try:
        theme = session.scalars(select(Theme).where(Theme.title == "Thème Vide")).one()
        theme_id = theme.id
    finally:
        session.close()
    client.post(
        f"/missions/{mission_id}/trame/themes/{theme_id}/questions",
        data={"label": "Jamais répondue ?", "qtype": "open"},
        follow_redirects=False,
    )

    response = client.get(f"/interviews/{interview_id}/preview")
    assert response.status_code == 200
    assert "Jamais répondue ?" in response.text
    assert "sans réponse" in response.text


def test_interview_import_from_document_prefills_to_review(
    client: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    response = client.post(
        "/missions",
        data={
            "name": "Mission Import Entretien",
            "description": "Description Import",
        },
        follow_redirects=False,
    )
    assert response.status_code == 303
    mission_id = response.headers["location"].rsplit("/", 1)[-1]

    response = client.post(
        f"/missions/{mission_id}/trame/themes",
        data={"title": "Thème Import"},
        follow_redirects=False,
    )
    assert response.status_code == 303

    session = SessionLocal()
    try:
        theme = session.scalars(
            select(Theme).where(Theme.title == "Thème Import")
        ).one()
    finally:
        session.close()

    response = client.post(
        f"/missions/{mission_id}/trame/themes/{theme.id}/questions",
        data={
            "label": "Question importée ?",
            "qtype": "open",
            "scale_min": "1",
            "scale_max": "5",
            "choices": "",
            "help_text": "",
        },
        follow_redirects=False,
    )
    assert response.status_code == 303

    session = SessionLocal()
    try:
        question = session.scalars(
            select(Question).where(
                Question.label == "Question importée ?",
                Question.theme_id == theme.id,
            )
        ).one()
        question_id = question.id
    finally:
        session.close()

    # Extraction IA mockée : pas d'appel réseau dans les tests.
    def fake_extract(questions, text):
        assert any(q.id == question_id for q in questions)
        return {question_id: {"text": "Réponse extraite du document", "verbatims": ["Citation extraite"]}}

    monkeypatch.setattr(
        "app.routers.interviews.extract_answers_from_text", fake_extract
    )

    docx_bytes = _build_trame_docx()
    response = client.post(
        f"/missions/{mission_id}/interviews/import",
        files={"file": ("entretien.docx", docx_bytes, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
        data={"interviewee_name": "Alice Martin"},
        follow_redirects=False,
    )
    assert response.status_code == 200
    assert "Revue avant import" in response.text
    assert "Réponse extraite du document" in response.text

    # Rien n'est encore créé avant confirmation.
    session = SessionLocal()
    try:
        assert session.scalars(select(Answer).where(Answer.question_id == question_id)).first() is None
    finally:
        session.close()

    proposed_json = json.dumps(
        {
            "identity": {"interviewee_name": "Alice Martin"},
            "answers": [
                {
                    "question_id": question_id,
                    "text": "Réponse extraite du document",
                    "verbatims": ["Citation extraite"],
                }
            ],
        }
    )
    response = client.post(
        f"/missions/{mission_id}/interviews/import/confirm",
        data={"proposed": proposed_json, "keep": str(question_id)},
        follow_redirects=False,
    )
    assert response.status_code == 303
    interview_id = response.headers["location"].rsplit("/", 1)[-1]

    session = SessionLocal()
    try:
        answer = session.scalars(
            select(Answer).where(Answer.question_id == question_id)
        ).one()
        assert answer.status == "to_review"
        assert answer.text == "Réponse extraite du document"
    finally:
        session.close()

    response = client.get(f"/interviews/{interview_id}")
    assert response.status_code == 200
    assert "Réponse extraite du document" in response.text


def test_interview_record_transcribes_and_prefills_to_review(
    client: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    response = client.post(
        "/missions",
        data={
            "name": "Mission Enregistrement",
            "description": "Description Enregistrement",
        },
        follow_redirects=False,
    )
    assert response.status_code == 303
    mission_id = response.headers["location"].rsplit("/", 1)[-1]

    response = client.post(
        f"/missions/{mission_id}/trame/themes",
        data={"title": "Thème Enregistrement"},
        follow_redirects=False,
    )
    assert response.status_code == 303

    session = SessionLocal()
    try:
        theme = session.scalars(
            select(Theme).where(Theme.title == "Thème Enregistrement")
        ).one()
    finally:
        session.close()

    response = client.post(
        f"/missions/{mission_id}/trame/themes/{theme.id}/questions",
        data={
            "label": "Question enregistrée ?",
            "qtype": "open",
            "scale_min": "1",
            "scale_max": "5",
            "choices": "",
            "help_text": "",
        },
        follow_redirects=False,
    )
    assert response.status_code == 303

    session = SessionLocal()
    try:
        question_id = session.scalars(
            select(Question).where(
                Question.label == "Question enregistrée ?",
                Question.theme_id == theme.id,
            )
        ).one().id
    finally:
        session.close()

    # La transcription se fait désormais au fil de l'eau côté client (par
    # segments envoyés à /audio/transcribe-segment) : cette route ne reçoit
    # plus que le texte déjà assemblé. Seule l'extraction IA doit être
    # mockée ici.
    def fake_extract(questions, text):
        assert text == "Transcription simulée."
        assert any(q.id == question_id for q in questions)
        return {question_id: {"text": "Réponse transcrite", "verbatims": []}}

    monkeypatch.setattr(
        "app.routers.interviews.extract_answers_from_text", fake_extract
    )

    response = client.post(
        f"/missions/{mission_id}/interviews/record",
        data={"interviewee_name": "Bruno Petit", "transcript": "Transcription simulée."},
        follow_redirects=False,
    )
    assert response.status_code == 200
    assert "Revue avant import" in response.text
    assert "Réponse transcrite" in response.text


def test_notes_record_appends_transcript_and_proposes_answers(
    client: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    response = client.post(
        "/missions",
        data={
            "name": "Mission Notes Enregistrement",
            "description": "Description Notes Enregistrement",
        },
        follow_redirects=False,
    )
    assert response.status_code == 303
    mission_id = response.headers["location"].rsplit("/", 1)[-1]

    response = client.post(
        f"/missions/{mission_id}/trame/themes",
        data={"title": "Thème Notes"},
        follow_redirects=False,
    )
    assert response.status_code == 303

    session = SessionLocal()
    try:
        theme_id = session.scalars(
            select(Theme).where(Theme.title == "Thème Notes")
        ).one().id
    finally:
        session.close()

    for label in ("Question déjà répondue ?", "Question encore ouverte ?"):
        response = client.post(
            f"/missions/{mission_id}/trame/themes/{theme_id}/questions",
            data={
                "label": label,
                "qtype": "open",
                "scale_min": "1",
                "scale_max": "5",
                "choices": "",
                "help_text": "",
            },
            follow_redirects=False,
        )
        assert response.status_code == 303

    session = SessionLocal()
    try:
        answered_qid = session.scalars(
            select(Question).where(Question.label == "Question déjà répondue ?")
        ).one().id
        open_qid = session.scalars(
            select(Question).where(Question.label == "Question encore ouverte ?")
        ).one().id
    finally:
        session.close()

    response = client.post(
        f"/missions/{mission_id}/interviews",
        data={"interviewee_name": "Claire Dubois"},
        follow_redirects=False,
    )
    assert response.status_code == 303
    interview_id = response.headers["location"].rsplit("/", 1)[-1]

    # Réponse manuelle existante sur la première question -> ne doit pas être
    # écrasée silencieusement par l'enregistrement.
    response = client.post(
        f"/interviews/{interview_id}/answers/{answered_qid}",
        data={"text": "Réponse manuelle déjà saisie"},
        follow_redirects=False,
    )
    assert response.status_code == 200

    monkeypatch.setattr(
        audio_transcribe, "transcribe_audio", lambda content: "Note vocale transcrite."
    )

    # Action 1 : la transcription (déclenchée automatiquement en JS dès
    # l'arrêt de l'enregistrement) ajoute le texte littéral aux notes libres,
    # sans passer par une revue.
    response = client.post(
        f"/interviews/{interview_id}/notes/transcribe",
        files={"file": ("note.webm", b"fake-audio-bytes", "audio/webm")},
        follow_redirects=False,
    )
    assert response.status_code == 200
    assert response.json() == {"free_notes": "Note vocale transcrite."}

    session = SessionLocal()
    try:
        interview = session.get(Interview, int(interview_id))
        assert interview.free_notes == "Note vocale transcrite."
    finally:
        session.close()

    def fake_extract(questions, text):
        assert text == "Note vocale transcrite."
        return {
            answered_qid: {"text": "Nouvelle proposition (déjà répondue)", "verbatims": []},
            open_qid: {"text": "Réponse pour la question ouverte", "verbatims": ["Citation notes"]},
        }

    monkeypatch.setattr(
        "app.routers.interviews.extract_answers_from_text", fake_extract
    )

    # Action 2 : le bouton "Répartir" analyse le contenu actuel des notes
    # libres et propose une distribution par question, à valider.
    response = client.post(
        f"/interviews/{interview_id}/notes/dispatch",
        data={"free_notes": "Note vocale transcrite."},
        follow_redirects=False,
    )
    assert response.status_code == 200
    assert "Revue de la répartition" in response.text
    assert "Note vocale transcrite." in response.text
    assert "Réponse manuelle déjà saisie" in response.text  # réponse existante affichée
    assert "Réponse pour la question ouverte" in response.text

    proposed_json = json.dumps(
        {
            "answers": [
                {
                    "question_id": answered_qid,
                    "text": "Nouvelle proposition (déjà répondue)",
                    "verbatims": [],
                },
                {
                    "question_id": open_qid,
                    "text": "Réponse pour la question ouverte",
                    "verbatims": ["Citation notes"],
                },
            ],
        }
    )
    # Ne garde que la question encore ouverte -> la réponse déjà répondue doit
    # rester intacte après confirmation.
    response = client.post(
        f"/interviews/{interview_id}/notes/confirm",
        data={"proposed": proposed_json, "keep": str(open_qid)},
        follow_redirects=False,
    )
    assert response.status_code == 303
    assert response.headers["location"].endswith(f"/interviews/{interview_id}?theme=notes")

    session = SessionLocal()
    try:
        interview = session.get(Interview, int(interview_id))
        # La répartition ne duplique pas le texte déjà présent dans les notes.
        assert interview.free_notes == "Note vocale transcrite."

        answered = session.scalars(
            select(Answer).where(Answer.question_id == answered_qid)
        ).one()
        assert answered.text == "Réponse manuelle déjà saisie"
        assert answered.status == "answered"

        opened = session.scalars(
            select(Answer).where(Answer.question_id == open_qid)
        ).one()
        assert opened.text == "Réponse pour la question ouverte"
        assert opened.status == "to_review"
    finally:
        session.close()


def _synthetic_webm_audio() -> bytes:
    """Construit un clip webm/opus minimal (imite un blob MediaRecorder de
    navigateur) pour exercer le vrai décodage + Whisper, sans mock."""
    import av
    import numpy as np

    buf = io.BytesIO()
    container = av.open(buf, mode="w", format="webm")
    stream = container.add_stream("libopus", rate=48000, layout="mono")
    samples = np.zeros((1, 48000), dtype=np.int16)
    frame = av.AudioFrame.from_ndarray(samples, format="s16", layout="mono")
    frame.sample_rate = 48000
    for packet in stream.encode(frame):
        container.mux(packet)
    for packet in stream.encode(None):
        container.mux(packet)
    container.close()
    return buf.getvalue()


@pytest.mark.skipif(
    not audio_transcribe.is_available(), reason="faster-whisper non installé"
)
def test_notes_transcribe_real_pipeline_end_to_end(client: TestClient) -> None:
    """Check sanity : exerce la vraie route HTTP /notes/transcribe avec un
    vrai fichier webm/opus, sans mocker `transcribe_audio` (contrairement à
    tous les autres tests) — c'est précisément ce chemin réel (décodage
    navigateur -> faster-whisper -> mise à jour des notes libres) qui a été
    à l'origine de régressions non détectées par la suite existante."""
    response = client.post(
        "/missions",
        data={
            "name": "Mission Sanity Transcription",
            "description": "Vérifie le pipeline réel de transcription.",
        },
        follow_redirects=False,
    )
    assert response.status_code == 303
    mission_id = response.headers["location"].rsplit("/", 1)[-1]

    response = client.post(
        f"/missions/{mission_id}/interviews",
        data={"interviewee_name": "Sanity Check"},
        follow_redirects=False,
    )
    assert response.status_code == 303
    interview_id = response.headers["location"].rsplit("/", 1)[-1]

    audio_bytes = _synthetic_webm_audio()
    response = client.post(
        f"/interviews/{interview_id}/notes/transcribe",
        files={"file": ("note.webm", audio_bytes, "audio/webm")},
    )
    # Le clip synthétique est un silence pur : avec `vad_filter=True`, Whisper
    # le rejette légitimement ("Aucune parole détectée") au lieu de halluciner
    # du texte dessus (comportement observé sans VAD). Les deux issues sont
    # acceptables ici — ce qu'on vérifie, c'est que le pipeline réel
    # (décodage webm/opus -> VAD -> Whisper -> réponse JSON) tourne de bout
    # en bout via la vraie route HTTP sans crasher et sans réponse mal formée
    # (`{"detail": ...}` ou 500 brute) — précisément la classe de bug déjà
    # rencontrée sur cet endpoint.
    assert response.status_code in (200, 422), response.text
    body = response.json()
    if response.status_code == 200:
        assert "free_notes" in body, response.text
        assert body["free_notes"].strip() != ""
        session = SessionLocal()
        try:
            interview = session.get(Interview, int(interview_id))
            assert interview.free_notes == body["free_notes"]
        finally:
            session.close()
    else:
        assert "error" in body, response.text


@pytest.mark.skipif(
    not audio_transcribe.is_available(), reason="faster-whisper non installé"
)
def test_audio_transcribe_segment_real_pipeline(client: TestClient) -> None:
    """Check sanity équivalent pour /audio/transcribe-segment, utilisé par la
    rotation de segments de record.html pendant un entretien long."""
    audio_bytes = _synthetic_webm_audio()
    response = client.post(
        "/audio/transcribe-segment",
        files={"file": ("segment.webm", audio_bytes, "audio/webm")},
    )
    # Silence pur -> "Aucune parole détectée" est une issue légitime avec
    # vad_filter=True (voir le test équivalent sur /notes/transcribe) ; ce
    # qui compte est que le pipeline réel tourne bout en bout sans crasher.
    assert response.status_code in (200, 422), response.text
    body = response.json()
    if response.status_code == 200:
        assert body["text"].strip() != ""
    else:
        assert "error" in body, response.text


def test_record_interview_backup_saves_file_to_disk(client: TestClient) -> None:
    from app.db import RECORDINGS_DIR

    response = client.post(
        "/missions",
        data={"name": "Mission Backup Audio", "description": "..."},
        follow_redirects=False,
    )
    assert response.status_code == 303
    mission_id = response.headers["location"].rsplit("/", 1)[-1]

    response = client.post(
        f"/missions/{mission_id}/interviews/record/backup",
        files={"file": ("entretien.webm", b"fake-full-interview-bytes", "audio/webm")},
    )
    assert response.status_code == 200, response.text
    body = response.json()
    assert "path" in body, response.text

    saved_file = RECORDINGS_DIR / body["path"]
    try:
        assert saved_file.read_bytes() == b"fake-full-interview-bytes"
    finally:
        saved_file.unlink(missing_ok=True)


def test_global_synthesis_generate_and_autosave(
    client: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    response = client.post(
        "/missions",
        data={"name": "Mission Synthese Globale", "description": "..."},
        follow_redirects=False,
    )
    assert response.status_code == 303
    mission_id = response.headers["location"].rsplit("/", 1)[-1]

    for title in ("Thème A", "Thème B"):
        response = client.post(
            f"/missions/{mission_id}/trame/themes",
            data={"title": title},
            follow_redirects=False,
        )
        assert response.status_code == 303

    session = SessionLocal()
    try:
        theme_ids = {
            t.title: t.id
            for t in session.scalars(
                select(Theme).where(Theme.title.in_(["Thème A", "Thème B"]))
            ).all()
        }
    finally:
        session.close()

    for title, theme_id in theme_ids.items():
        response = client.post(
            f"/missions/{mission_id}/trame/themes/{theme_id}/questions",
            data={
                "label": f"Question {title} ?",
                "qtype": "open",
                "scale_min": "1",
                "scale_max": "5",
                "choices": "",
                "help_text": "",
            },
            follow_redirects=False,
        )
        assert response.status_code == 303

    session = SessionLocal()
    try:
        question_id_by_theme = {q.theme_id: q.id for q in session.scalars(select(Question)).all()}
    finally:
        session.close()

    response = client.post(
        f"/missions/{mission_id}/interviews",
        data={"interviewee_name": "Alice Dupont"},
        follow_redirects=False,
    )
    assert response.status_code == 303
    interview_id = response.headers["location"].rsplit("/", 1)[-1]

    for theme_title, theme_id in theme_ids.items():
        qid = question_id_by_theme[theme_id]
        response = client.post(
            f"/interviews/{interview_id}/answers/{qid}",
            data={"text": f"Réponse sur {theme_title}"},
        )
        assert response.status_code == 200

    monkeypatch.setattr("app.routers.synthese.is_configured", lambda: True)

    def fake_generate_global(mission, material_by_theme):
        assert mission.name == "Mission Synthese Globale"
        # Les deux thèmes de la mission doivent être présents, pas un seul.
        assert len(material_by_theme) == 2
        return {
            "contexte": "- Contexte de test",
            "culture_adn": "- Culture de test",
            "forces_succes": "- Force de test",
            "points_amelioration": "- Point d'amélioration de test",
            "aspirations": "- Aspiration de test",
        }

    monkeypatch.setattr(
        "app.routers.synthese.generate_global_synthesis", fake_generate_global
    )

    response = client.post(f"/missions/{mission_id}/synthese/globale/generate")
    assert response.status_code == 200
    assert "Contexte de test" in response.text

    session = SessionLocal()
    try:
        mission = session.get(Mission, int(mission_id))
        assert mission.global_synthesis is not None
        assert mission.global_synthesis.contexte == "- Contexte de test"
        assert mission.global_synthesis.status == "generated"
    finally:
        session.close()

    # Autosave d'un champ -> passe le statut à "edited".
    response = client.post(
        f"/syntheses/globale/{mission_id}/field",
        data={"field": "aspirations", "value": "- Aspiration modifiée à la main"},
    )
    assert response.status_code == 200

    session = SessionLocal()
    try:
        mission = session.get(Mission, int(mission_id))
        assert mission.global_synthesis.aspirations == "- Aspiration modifiée à la main"
        assert mission.global_synthesis.status == "edited"
    finally:
        session.close()


def test_recommendations_generate_from_global_synthesis(
    client: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    response = client.post(
        "/missions",
        data={"name": "Mission Reco", "description": "Test recommandations"},
        follow_redirects=False,
    )
    assert response.status_code == 303
    mission_id = response.headers["location"].rsplit("/", 1)[-1]

    response = client.post(
        f"/missions/{mission_id}/trame/themes",
        data={"title": "Thème Reco"},
        follow_redirects=False,
    )
    assert response.status_code == 303
    session = SessionLocal()
    try:
        theme_id = session.scalars(select(Theme).where(Theme.title == "Thème Reco")).one().id
    finally:
        session.close()

    response = client.post(
        f"/missions/{mission_id}/trame/themes/{theme_id}/questions",
        data={
            "label": "Question Reco ?",
            "qtype": "open",
            "scale_min": "1",
            "scale_max": "5",
            "choices": "",
            "help_text": "",
        },
        follow_redirects=False,
    )
    assert response.status_code == 303
    session = SessionLocal()
    try:
        question_id = session.scalars(
            select(Question).where(Question.theme_id == theme_id)
        ).one().id
    finally:
        session.close()

    response = client.post(
        f"/missions/{mission_id}/interviews",
        data={"interviewee_name": "Bob Martin"},
        follow_redirects=False,
    )
    assert response.status_code == 303
    interview_id = response.headers["location"].rsplit("/", 1)[-1]
    response = client.post(
        f"/interviews/{interview_id}/answers/{question_id}",
        data={"text": "Réponse pour générer une synthèse globale"},
    )
    assert response.status_code == 200

    monkeypatch.setattr("app.routers.synthese.is_configured", lambda: True)
    monkeypatch.setattr(
        "app.routers.synthese.generate_global_synthesis",
        lambda mission, material_by_theme: {
            "contexte": "- Contexte",
            "culture_adn": "- Culture",
            "forces_succes": "- Forces",
            "points_amelioration": "- Douleurs",
            "aspirations": "- Aspirations",
        },
    )
    response = client.post(f"/missions/{mission_id}/synthese/globale/generate")
    assert response.status_code == 200

    def fake_generate_recommendations(global_synthesis):
        assert global_synthesis.contexte == "- Contexte"
        return [
            {
                "title": "Axe de test",
                "recommendations": [
                    {
                        "title": "Reco de test",
                        "objectif": "Objectif de test",
                        "acteurs": "CODIR",
                        "valeur": 4,
                        "complexite": 2,
                        "proposition_valeur": "Valeur de test",
                        "plan_actions": "- Action 1",
                        "resultats_attendus": "- Résultat 1",
                    }
                ],
            }
        ]

    monkeypatch.setattr(
        "app.routers.synthese.generate_recommendations", fake_generate_recommendations
    )

    response = client.post(f"/missions/{mission_id}/recommandations/generate")
    assert response.status_code == 200
    assert "Axe de test" in response.text
    assert "Reco de test" in response.text

    session = SessionLocal()
    try:
        mission = session.get(Mission, int(mission_id))
        assert len(mission.recommendation_axes) == 1
        axis = mission.recommendation_axes[0]
        assert axis.title == "Axe de test"
        assert len(axis.recommendations) == 1
        reco = axis.recommendations[0]
        assert reco.title == "Reco de test"
        assert reco.valeur == 4
        assert reco.complexite == 2
        recommendation_id = reco.id
    finally:
        session.close()

    # Autosave sur une fiche de recommandation.
    response = client.post(
        f"/recommandations/{recommendation_id}/field",
        data={"field": "acteurs", "value": "CODIR + Factory"},
    )
    assert response.status_code == 200

    session = SessionLocal()
    try:
        reco = session.get(Recommendation, recommendation_id)
        assert reco.acteurs == "CODIR + Factory"
        axis_id = reco.axis_id
    finally:
        session.close()

    # Autosave sur le titre d'un axe.
    response = client.post(
        f"/recommandations/axes/{axis_id}/field",
        data={"field": "title", "value": "Axe renommé"},
    )
    assert response.status_code == 200

    session = SessionLocal()
    try:
        axis = session.get(RecommendationAxis, axis_id)
        assert axis.title == "Axe renommé"
    finally:
        session.close()

    # Champ inconnu -> 400, pas de crash.
    response = client.post(
        f"/recommandations/axes/{axis_id}/field",
        data={"field": "position", "value": "9"},
    )
    assert response.status_code == 400


def test_recommendations_require_global_synthesis_first(client: TestClient) -> None:
    """Garde-fou : pas de synthèse globale -> message d'invitation, pas de crash."""
    response = client.post(
        "/missions",
        data={"name": "Mission Sans Synthese", "description": "..."},
        follow_redirects=False,
    )
    assert response.status_code == 303
    mission_id = response.headers["location"].rsplit("/", 1)[-1]

    response = client.post(
        f"/missions/{mission_id}/trame/themes",
        data={"title": "Thème Vide"},
        follow_redirects=False,
    )
    assert response.status_code == 303

    response = client.post(f"/missions/{mission_id}/recommandations/generate")
    assert response.status_code == 200
    # Apostrophe HTML-échappée par l'autoescape Jinja (d&#39;abord).
    assert "Générez d" in response.text and "abord la synthèse globale" in response.text

    session = SessionLocal()
    try:
        mission = session.get(Mission, int(mission_id))
        assert mission.recommendation_axes == []
    finally:
        session.close()


def _setup_mission_with_one_answer(client: TestClient, name: str) -> tuple[str, str]:
    """Mission + 1 thème + 1 question + 1 entretien répondu. Renvoie
    (mission_id, question_id) — utilitaire pour les tests export/import."""
    response = client.post(
        "/missions", data={"name": name, "description": "..."}, follow_redirects=False
    )
    assert response.status_code == 303
    mission_id = response.headers["location"].rsplit("/", 1)[-1]

    response = client.post(
        f"/missions/{mission_id}/trame/themes", data={"title": "Thème Export"}, follow_redirects=False
    )
    assert response.status_code == 303
    session = SessionLocal()
    try:
        theme_id = session.scalars(
            select(Theme).join(Trame).where(Trame.mission_id == int(mission_id))
        ).one().id
    finally:
        session.close()

    response = client.post(
        f"/missions/{mission_id}/trame/themes/{theme_id}/questions",
        data={
            "label": "Question Export ?", "qtype": "open",
            "scale_min": "1", "scale_max": "5", "choices": "", "help_text": "",
        },
        follow_redirects=False,
    )
    assert response.status_code == 303
    session = SessionLocal()
    try:
        question_id = session.scalars(select(Question).where(Question.theme_id == theme_id)).one().id
    finally:
        session.close()

    response = client.post(
        f"/missions/{mission_id}/interviews", data={"interviewee_name": "Zoé Interviewée"}, follow_redirects=False
    )
    assert response.status_code == 303
    interview_id = response.headers["location"].rsplit("/", 1)[-1]
    response = client.post(
        f"/interviews/{interview_id}/answers/{question_id}",
        data={"text": "Réponse test pour l'export."},
    )
    assert response.status_code == 200

    return mission_id, question_id


_FILLED_ANALYSIS = """
## SYNTHÈSE GLOBALE

### Contexte
- Croissance rapide

### Culture & ADN
Culture morcelée

### Forces & succès
Bonne entraide

### Points d'amélioration
Priorisation floue

### Aspirations (baguette magique)
Vision claire

## RECOMMANDATIONS

#### Axe 1 : Cohérence du cadre

##### Recommandation 1.1 : Revoir les principes
- Objectif : Réduire l'ambiguïté
- Acteurs : CODIR
- Valeur (1-5) : 4
- Complexité (1-5) : 2
- Proposition de valeur : Un cadre clair
- Plan d'actions : Atelier de clarification
- Résultats attendus : Moins de confusion
"""


def test_export_markdown_contains_data_and_analysis_template(client: TestClient) -> None:
    mission_id, _qid = _setup_mission_with_one_answer(client, "Mission Export Markdown")

    response = client.get(f"/missions/{mission_id}/export/interviews")
    assert response.status_code == 200
    assert response.headers["content-type"].startswith("text/markdown")
    assert "attachment" in response.headers["content-disposition"]

    text = response.text
    assert "Réponse test pour l'export." in text
    assert "## SYNTHÈSE GLOBALE" in text
    assert "### Contexte" in text
    assert "## RECOMMANDATIONS" in text
    assert "#### Axe 1" in text
    assert "##### Recommandation 1.1" in text


def test_import_analyse_populates_global_synthesis_and_recommendations(client: TestClient) -> None:
    mission_id, _qid = _setup_mission_with_one_answer(client, "Mission Import Analyse")

    files = {"file": ("analyse.md", _FILLED_ANALYSIS.encode("utf-8"), "text/markdown")}
    response = client.post(
        f"/missions/{mission_id}/import/analyse", files=files, follow_redirects=False
    )
    assert response.status_code == 303
    # Après import, la suite logique du parcours est l'étape 2 (synthèse
    # globale), pas l'écran d'export/import ni l'écran d'export PPT.
    assert response.headers["location"].endswith(f"/missions/{mission_id}/synthese/globale")

    session = SessionLocal()
    try:
        mission = session.get(Mission, int(mission_id))
        assert mission.global_synthesis is not None
        assert mission.global_synthesis.contexte == "- Croissance rapide"
        assert mission.global_synthesis.status == "generated"
        assert len(mission.recommendation_axes) == 1
        axis = mission.recommendation_axes[0]
        assert axis.title == "Cohérence du cadre"
        reco = axis.recommendations[0]
        assert reco.title == "Revoir les principes"
        assert reco.valeur == 4
        assert reco.complexite == 2
    finally:
        session.close()

    # L'aperçu reflète le contenu importé.
    response = client.get(f"/missions/{mission_id}/synthese/apercu")
    assert response.status_code == 200
    assert "Cohérence du cadre" in response.text
    assert "Revoir les principes" in response.text


def test_import_analyse_invalid_file_shows_friendly_error(client: TestClient) -> None:
    mission_id, _qid = _setup_mission_with_one_answer(client, "Mission Import Invalide")

    files = {"file": ("bad.md", b"# Juste du texte\nSans structure reconnue.", "text/markdown")}
    response = client.post(f"/missions/{mission_id}/import/analyse", files=files)
    assert response.status_code == 200
    assert "Structure du fichier non reconnue" in response.text

    session = SessionLocal()
    try:
        mission = session.get(Mission, int(mission_id))
        # Un GlobalSynthesis (vide) est désormais toujours garanti en arrivant
        # sur l'écran Analyse (le sous-onglet IA intégrée en a besoin) —
        # l'import qui échoue ne le remplit simplement pas.
        assert mission.global_synthesis is not None
        assert not mission.global_synthesis.has_content
        assert mission.recommendation_axes == []
    finally:
        session.close()


def test_export_pptx_produces_valid_file(client: TestClient) -> None:
    mission_id, _qid = _setup_mission_with_one_answer(client, "Mission Export PPTX")
    files = {"file": ("analyse.md", _FILLED_ANALYSIS.encode("utf-8"), "text/markdown")}
    client.post(f"/missions/{mission_id}/import/analyse", files=files, follow_redirects=False)

    response = client.get(f"/missions/{mission_id}/export/pptx")
    assert response.status_code == 200
    assert response.headers["content-type"] == (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert len(response.content) > 1000

    prs = Presentation(io.BytesIO(response.content))
    assert len(list(prs.slides)) > 5


def test_pptx_template_upload_and_reuse(client: TestClient) -> None:
    mission_id, _qid = _setup_mission_with_one_answer(client, "Mission Template PPTX")

    # Fichier invalide -> erreur propre, pas de crash, rien d'enregistré.
    response = client.post(
        f"/missions/{mission_id}/pptx-template",
        files={"file": ("fake.pptx", b"not a real pptx", "application/octet-stream")},
    )
    assert response.status_code == 200
    assert "invalide" in response.text.lower() or "corrompu" in response.text.lower()

    # Fichier valide -> accepté, lié à la mission, réutilisé pour l'export.
    template_buf = io.BytesIO()
    Presentation().save(template_buf)
    response = client.post(
        f"/missions/{mission_id}/pptx-template",
        files={"file": ("template.pptx", template_buf.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        follow_redirects=False,
    )
    assert response.status_code == 303

    session = SessionLocal()
    try:
        mission = session.get(Mission, int(mission_id))
        assert mission.pptx_template_path == f"{mission_id}.pptx"
    finally:
        session.close()

    response = client.get(f"/missions/{mission_id}/export/pptx")
    assert response.status_code == 200
    Presentation(io.BytesIO(response.content))  # ne lève pas


def test_export_import_view_is_step_one_of_the_wizard(client: TestClient) -> None:
    mission_id, _qid = _setup_mission_with_one_answer(client, "Mission Wizard Etape 1")

    response = client.get(f"/missions/{mission_id}/synthese/export-import")
    assert response.status_code == 200
    # L'étape 1 s'appelle "Analyse" (évol) — 2 sous-onglets : IA intégrée et
    # export/import manuel.
    assert "Analyse" in response.text
    assert "IA intégrée" in response.text
    assert "Exporter les entretiens" in response.text
    assert "Importer l'analyse" in response.text
    # Les actions d'export/import ne sont plus sur l'écran "Export PPT".
    apercu_response = client.get(f"/missions/{mission_id}/synthese/apercu")
    assert "Importer l'analyse" not in apercu_response.text


def test_wizard_steps_show_done_state_once_populated(client: TestClient) -> None:
    mission_id, _qid = _setup_mission_with_one_answer(client, "Mission Wizard Stepper")

    # Avant tout import/génération : les étapes 2/3 affichent leur numéro,
    # aucune coche nulle part sur la page.
    response = client.get(f"/missions/{mission_id}/synthese/export-import")
    assert response.text.count("✓") == 0

    files = {"file": ("analyse.md", _FILLED_ANALYSIS.encode("utf-8"), "text/markdown")}
    client.post(f"/missions/{mission_id}/import/analyse", files=files, follow_redirects=False)

    # Après import : au moins les étapes 2 (synthèse globale) et 3
    # (recommandations) affichent une coche à la place de leur numéro.
    response = client.get(f"/missions/{mission_id}/synthese/export-import")
    assert response.text.count("✓") >= 2


def test_export_pptx_partial_axis_selection_produces_fewer_slides(client: TestClient) -> None:
    mission_id, _qid = _setup_mission_with_one_answer(client, "Mission Selection Partielle")

    two_axes_analysis = _FILLED_ANALYSIS + """
#### Axe 2 : Deuxième axe

##### Recommandation 2.1 : Une autre recommandation
- Objectif : Autre objectif
- Acteurs : Factory
- Valeur (1-5) : 3
- Complexité (1-5) : 3
- Proposition de valeur : Autre valeur
- Plan d'actions : Autre action
- Résultats attendus : Autre résultat
"""
    files = {"file": ("analyse.md", two_axes_analysis.encode("utf-8"), "text/markdown")}
    client.post(f"/missions/{mission_id}/import/analyse", files=files, follow_redirects=False)

    session = SessionLocal()
    try:
        mission = session.get(Mission, int(mission_id))
        assert len(mission.recommendation_axes) == 2
        first_axis_id = mission.recommendation_axes[0].id
    finally:
        session.close()

    response_full = client.get(f"/missions/{mission_id}/export/pptx")
    assert response_full.status_code == 200
    slides_full = len(list(Presentation(io.BytesIO(response_full.content)).slides))

    # Sélection explicite : uniquement le premier axe, sans sommaire.
    response_partial = client.get(
        f"/missions/{mission_id}/export/pptx",
        params={
            "config_submitted": "1",
            "synthese": "true",
            "axes_overview": "true",
            "matrix": "true",
            "axis": first_axis_id,
        },
    )
    assert response_partial.status_code == 200
    slides_partial = len(list(Presentation(io.BytesIO(response_partial.content)).slides))

    assert slides_partial < slides_full


def _long_text(n_sentences: int) -> str:
    sentence = (
        "Cette phrase est volontairement longue pour simuler une réponse "
        "détaillée d'entretien, avec beaucoup de mots afin de forcer un "
        "retour à la ligne multiple dans la zone de texte de la slide."
    )
    return " ".join(sentence for _ in range(n_sentences))


def _many_axes_analysis(n_axes: int, n_recos_per_axis: int) -> str:
    """Analyse markdown adversariale (US7.1) : synthèse globale avec des
    paragraphes très longs, et N axes de M recommandations chacun, chaque
    champ texte étant lui aussi très long — sert à vérifier que la mise en
    page absorbe des contenus bien plus fournis que le jeu de données de
    test habituel sans faire déborder aucune forme de la slide."""
    lines = ["## SYNTHÈSE GLOBALE", ""]
    for heading in ("Contexte", "Culture & ADN", "Forces & succès", "Points d'amélioration", "Aspirations (baguette magique)"):
        lines += [f"### {heading}", _long_text(4), ""]
    lines += ["## RECOMMANDATIONS", ""]
    for i in range(1, n_axes + 1):
        lines += [f"#### Axe {i} : Axe de transformation numéro {i} avec un intitulé assez long pour tester le débordement", ""]
        for j in range(1, n_recos_per_axis + 1):
            lines += [
                f"##### Recommandation {i}.{j} : Une recommandation avec un titre lui-même particulièrement long et détaillé",
                f"- Objectif : {_long_text(3)}",
                "- Acteurs : Direction Générale, DRH, DSI, Managers de proximité, Représentants du personnel",
                f"- Valeur (1-5) : {(i + j) % 5 + 1}",
                f"- Complexité (1-5) : {(i * j) % 5 + 1}",
                f"- Proposition de valeur : {_long_text(3)}",
                f"- Plan d'actions : {_long_text(4)}",
                f"- Résultats attendus : {_long_text(3)}",
                "",
            ]
    return "\n".join(lines)


def test_export_pptx_geometry_clean_with_long_text_and_many_axes(client: TestClient) -> None:
    """US7.1 : le garde-fou géométrique (build_presentation -> verifier_geometrie)
    ne doit jamais se déclencher, même avec un contenu bien plus dense que le
    jeu de données de test habituel (textes longs, 5 axes x 3 recommandations)."""
    mission_id, _qid = _setup_mission_with_one_answer(client, "Mission Geometrie Contenu Dense")
    analysis = _many_axes_analysis(n_axes=5, n_recos_per_axis=3)
    files = {"file": ("analyse.md", analysis.encode("utf-8"), "text/markdown")}
    response = client.post(f"/missions/{mission_id}/import/analyse", files=files, follow_redirects=False)
    assert response.status_code == 303

    # build_presentation lève RuntimeError si verifier_geometrie détecte une
    # forme hors cadre — une réponse 200 avec un .pptx valide suffit donc à
    # prouver que le garde-fou n'a rien trouvé à redire.
    response = client.get(f"/missions/{mission_id}/export/pptx")
    assert response.status_code == 200
    prs = Presentation(io.BytesIO(response.content))
    assert len(list(prs.slides)) > 20  # titre+sommaire+5 synthèse+axes+matrice+15 recos


def test_export_pptx_geometry_clean_with_4_3_client_template(client: TestClient) -> None:
    """US7.1 : un template client au format 4:3 (au lieu du 16:9 par défaut)
    ne doit pas non plus faire déborder la mise en page — toutes les slides
    lisent prs.slide_width/height dynamiquement plutôt qu'une constante."""
    mission_id, _qid = _setup_mission_with_one_answer(client, "Mission Geometrie Template 4:3")
    files = {"file": ("analyse.md", _FILLED_ANALYSIS.encode("utf-8"), "text/markdown")}
    client.post(f"/missions/{mission_id}/import/analyse", files=files, follow_redirects=False)

    template_buf = io.BytesIO()
    template_prs = Presentation()
    template_prs.slide_width = Inches(10)
    template_prs.slide_height = Inches(7.5)
    template_prs.save(template_buf)
    client.post(
        f"/missions/{mission_id}/pptx-template",
        files={"file": ("template_4_3.pptx", template_buf.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        follow_redirects=False,
    )

    response = client.get(f"/missions/{mission_id}/export/pptx")
    assert response.status_code == 200
    prs = Presentation(io.BytesIO(response.content))
    assert Emu(prs.slide_width).inches == pytest.approx(10)


def _slide_texts(slide) -> list[str]:
    """Tous les paragraphes texte d'une slide, toutes zones de texte
    confondues — utilitaire pour chercher un titre ou une puce sans dépendre
    de la forme exacte qui la porte."""
    texts = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p in shape.text_frame.paragraphs:
            texts.append("".join(run.text for run in p.runs))
    return texts


def test_export_pptx_axes_overview_paginates_when_many_axes(client: TestClient) -> None:
    """US6.4 (pagination auto) : au-delà du nombre d'axes qui tient
    lisiblement sur une slide, la vue d'ensemble des axes se poursuit sur
    des slides « (k/n) » plutôt que d'écraser les cartes en rangées
    illisibles — aucun axe ne doit être perdu ni dupliqué au passage."""
    mission_id, _qid = _setup_mission_with_one_answer(client, "Mission Pagination Axes")
    analysis = _many_axes_analysis(n_axes=10, n_recos_per_axis=1)
    files = {"file": ("analyse.md", analysis.encode("utf-8"), "text/markdown")}
    response = client.post(f"/missions/{mission_id}/import/analyse", files=files, follow_redirects=False)
    assert response.status_code == 303

    response = client.get(f"/missions/{mission_id}/export/pptx")
    assert response.status_code == 200
    prs = Presentation(io.BytesIO(response.content))

    title_re = re.compile(r"^Les recommandations sont construites autour de ces axes \((\d+)/(\d+)\)$")
    axes_slides = []
    for slide in prs.slides:
        for text in _slide_texts(slide):
            m = title_re.match(text)
            if m:
                axes_slides.append((slide, m.group(1), m.group(2)))
                break

    assert len(axes_slides) >= 2, "attendu plusieurs slides de vue d'ensemble des axes paginées"
    assert {n for _, _, n in axes_slides} == {axes_slides[0][2]}  # même total "n" partout

    numbers = []
    for slide, _, _ in axes_slides:
        for text in _slide_texts(slide):
            m = re.fullmatch(r"#(\d+)", text)
            if m:
                numbers.append(int(m.group(1)))
    assert sorted(numbers) == list(range(1, 11))  # les 10 axes, une seule fois chacun


def test_export_pptx_recommendation_bullet_overflow_spills_to_continuation_slide(client: TestClient) -> None:
    """US6.4 (pagination auto) : un plan d'actions trop long pour tenir dans
    sa colonne, même à la police minimale, doit se poursuivre sur une slide
    « (suite — Plan d'actions) » au lieu de déborder silencieusement de sa
    zone de texte (ce que `verifier_geometrie` ne peut pas détecter)."""
    mission_id, _qid = _setup_mission_with_one_answer(client, "Mission Pagination Plan Actions")

    bullets = [f"Action numéro {i} : {_long_text(1)}" for i in range(1, 31)]
    analysis = f"""
## SYNTHÈSE GLOBALE

### Contexte
Contexte bref.

## RECOMMANDATIONS

#### Axe 1 : Cohérence du cadre

##### Recommandation 1.1 : Revoir les principes
- Objectif : Réduire l'ambiguïté
- Acteurs : CODIR
- Valeur (1-5) : 4
- Complexité (1-5) : 2
- Proposition de valeur : Un cadre clair
- Plan d'actions : {bullets[0]}
{chr(10).join(bullets[1:])}
"""
    files = {"file": ("analyse.md", analysis.encode("utf-8"), "text/markdown")}
    response = client.post(f"/missions/{mission_id}/import/analyse", files=files, follow_redirects=False)
    assert response.status_code == 303

    response = client.get(f"/missions/{mission_id}/export/pptx")
    assert response.status_code == 200
    prs = Presentation(io.BytesIO(response.content))

    def is_relevant(slide) -> bool:
        texts = _slide_texts(slide)
        return any("Revoir les principes" in t or "(suite — Plan d'actions)" in t for t in texts)

    relevant_slides = [slide for slide in prs.slides if is_relevant(slide)]
    assert len(relevant_slides) >= 2, "attendu la slide principale + au moins une slide de continuation"

    all_bullets = [
        text[len("•  "):]
        for slide in relevant_slides
        for text in _slide_texts(slide)
        if text.startswith("•  ") and text != "•  —"  # placeholder du champ "Résultats attendus" laissé vide
    ]
    assert len(all_bullets) == len(bullets)  # aucune puce perdue


def _soffice_path() -> str | None:
    import shutil

    found = shutil.which("soffice") or shutil.which("soffice.exe")
    if found:
        return found
    for candidate in (
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ):
        if Path(candidate).exists():
            return candidate
    return None


@pytest.mark.skipif(_soffice_path() is None, reason="LibreOffice non installé")
def test_export_pptx_renders_cleanly_in_a_real_engine(client: TestClient, tmp_path: Path) -> None:
    """US7.2 (vérification visuelle) : `verifier_geometrie()` ne détecte que
    des formes hors-cadre — pas un fichier que PowerPoint/LibreOffice refuse
    d'ouvrir ou tronque silencieusement, la classe de bug ayant justement
    motivé l'incrément 5 (cf. CLAUDE.md, « template client illisible par
    PowerPoint » découvert via l'automation COM, pas par la suite de tests).
    Un `.pptx` peut parser sans erreur avec python-pptx (parseur tolérant)
    tout en étant rejeté ou mal rendu par un vrai moteur — ce test automatise
    le réflexe de vérification réelle : convertir l'export en PDF via
    LibreOffice et vérifier qu'il produit bien une page par slide attendue,
    plutôt que de se fier au seul comptage `len(prs.slides)`."""
    mission_id, _qid = _setup_mission_with_one_answer(client, "Mission Rendu Reel LibreOffice")
    analysis = _many_axes_analysis(n_axes=6, n_recos_per_axis=2)
    files = {"file": ("analyse.md", analysis.encode("utf-8"), "text/markdown")}
    client.post(f"/missions/{mission_id}/import/analyse", files=files, follow_redirects=False)

    response = client.get(f"/missions/{mission_id}/export/pptx")
    assert response.status_code == 200
    prs = Presentation(io.BytesIO(response.content))
    expected_slides = len(list(prs.slides))

    pptx_path = tmp_path / "export.pptx"
    pptx_path.write_bytes(response.content)
    result = subprocess.run(
        [_soffice_path(), "--headless", "--convert-to", "pdf", "--outdir", str(tmp_path), str(pptx_path)],
        capture_output=True, timeout=120,
    )
    assert result.returncode == 0, result.stderr.decode(errors="replace")
    pdf_path = tmp_path / "export.pdf"
    assert pdf_path.exists(), "LibreOffice n'a produit aucun PDF — l'export ne s'ouvrirait probablement pas non plus dans PowerPoint"

    pdf_bytes = pdf_path.read_bytes()
    assert len(pdf_bytes) > 2000, "PDF quasi vide — rendu suspect"
    # Comptage de pages par regex plutôt qu'une dépendance PDF supplémentaire
    # (aucune lib de lecture PDF dans ce venv) : suffisant pour un PDF simple
    # généré par LibreOffice, pas cense être un parseur PDF général.
    page_count = len(re.findall(rb"/Type\s*/Page[^s]", pdf_bytes))
    assert page_count == expected_slides, (
        f"{page_count} page(s) rendue(s) par LibreOffice pour {expected_slides} slide(s) exportée(s) "
        "— le fichier n'est probablement pas ouvrable proprement par un vrai lecteur PowerPoint"
    )
