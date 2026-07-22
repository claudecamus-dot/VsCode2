"""Génération de l'export PowerPoint (évol) — restitue à l'identique de
l'aperçu web (`synthese/apercu.html`) : slide de titre, sommaire, une slide
par catégorie de synthèse globale, une vue d'ensemble des axes, une matrice
effort/valeur (graphique natif PowerPoint) puis une slide par recommandation
(gabarit fixe calqué sur un rapport de restitution réel).

Si `mission.pptx_template_path` est renseigné, la présentation démarre à
partir de ce template client (hérite thème/masters/logo) ; sinon une
présentation vierge en 16:9, stylée via `pptx_deck` (skill pptx-deck, copié
tel quel dans ce module pour ne pas dépendre du chemin d'installation du
skill — cf. son propre en-tête).
"""
from __future__ import annotations

import os
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from ..models import Mission
from . import pptx_deck as D

# --- Cadres photo (têtes de chapitre, P3) : skill pptx-framed-image (greffé,
# présent dans .claude/skills/). Import gardé — si le skill/Pillow manque, les
# intercalaires retombent proprement sur leur version texte-seul (cf. _slide_chapitre).
_FRAMED_OK = False
try:  # pragma: no cover - dépend de la présence du skill + Pillow
    import sys as _sys
    _FRAMED_SCRIPTS = Path(__file__).resolve().parents[2] / ".claude" / "skills" / "pptx-framed-image" / "scripts"
    if str(_FRAMED_SCRIPTS) not in _sys.path:
        _sys.path.insert(0, str(_FRAMED_SCRIPTS))
    from framed_image import place_image_in_frame as _place_image_in_frame  # type: ignore
    from framed_image import cover_crop_to_aspect as _cover_crop_to_aspect  # type: ignore
    import nature_images as _nature_images  # type: ignore
    import stock_images as _stock_images  # type: ignore  # fetch Openverse (vraies photos CC0)
    _IMG_CACHE = Path(__file__).resolve().parents[2] / "data" / "pptx_chapitre_images"
    _FRAMED_OK = True
except Exception:  # skill absent, Pillow non installé, etc. -> repli texte-seul
    _FRAMED_OK = False

MARGIN = 0.6

# --------------------------------------------------------------------------- #
# Repères "forme" pour l'éditeur web par onglets (aperçu.html) — mêmes
# contraintes géométriques (largeur, hauteur max, échelle typo) que les
# fonctions de slide ci-dessous, dupliquées ici en constantes plutôt que
# recalculées dynamiquement : l'éditeur enregistre un champ à la fois, hors
# contexte d'une Presentation réelle (pas de mission/axes/inclusions connus à
# cet instant). Restent donc des repères indicatifs, pas une garantie — le
# garde-fou qui compte vraiment reste `D.verifier_geometrie()` à l'export.
# plan_actions/resultats_attendus utilisent en vrai l'espace *restant* après
# les blocs précédents (variable) ; ici on prend une estimation généreuse
# mais fixe, cohérente avec une slide "normale".
_W_IN, _H_IN = 10.0, 5.625  # dims du template OCTO de marque (16:9) — FIELD_SHAPE (hints web) aligné dessus
# Fiche reco en encarts arrondis (2026-07-22) : le contenu vit DANS des cartes
# (carte gauche + encart proposition + carte plan), les largeurs utiles sont donc
# les largeurs de carte moins les marges internes (pad 0.2 ×2 + liseré 0.05).
# 3.15 (resserrée) : la carte droite porte plan+résultats, souvent de longues
# puces — elle a besoin de la largeur (cf. _slide_recommendation).
_CARD_L_W = 3.15
_LEFT_W = _CARD_L_W - 0.45
_RIGHT_W = (_W_IN - MARGIN - (MARGIN + _CARD_L_W + 0.3)) - 0.45
# Slide de synthèse enrichie (claim + visuel + encart) : largeur de la carte de puces
# une fois la bande photo réservée à droite (2.7in), sinon pleine largeur (repli).
_SYNTH_VIS_W = 2.7
_SYNTH_AREA_W = (
    ((_W_IN - MARGIN - _SYNTH_VIS_W) - 0.3 - (MARGIN + 0.3))
    if _FRAMED_OK else (_W_IN - 2 * (MARGIN + 0.3))
)

# Template OCTO de marque, versionné (masters/layouts/thème + police Outfit) : défaut de
# build_presentation. Un template client (mission.pptx_template_path) reste prioritaire ;
# le chrome (logo/pied de page/n° de slide) survit via _pick_layout (« titre seul »).
OCTO_TEMPLATE_PATH = Path(__file__).resolve().parent.parent / "assets" / "template-octo.pptx"

FIELD_SHAPE = {
    "objectif": dict(width_in=_LEFT_W, max_h_in=1.1),
    "acteurs": dict(width_in=_LEFT_W, max_h_in=0.5),
    # resultats_attendus vit dans le bandeau bas PLEINE LARGEUR de la fiche
    # (encart gris, 2026-07-22) : taille FIXE small, tronqué à 2 lignes — hint en
    # mode size_pt/max_lignes (annonce la troncature, ne promet pas un shrink que
    # la slide ne fait pas ; même dérive corrigée que es_headline/difficulty_label).
    "resultats_attendus": dict(width_in=_W_IN - 2 * MARGIN - 0.45, size_pt=D.TYPE["small"], max_lignes=2),
    # 0.80 = prop_h(1.00) - 0.20 de _slide_recommendation — recalé revue adversariale.
    "proposition_valeur": dict(width_in=_RIGHT_W, max_h_in=0.80),
    "plan_actions": dict(width_in=_RIGHT_W, max_h_in=2.0),
    # 1 ligne : le titre de fiche est tronqué à l'ellipse (le complet vit sur la
    # vue d'ensemble des axes) — recalé revue adversariale.
    "reco_title": dict(width_in=_W_IN - 2 * MARGIN, size_pt=D.TYPE["title"], max_lignes=1),
    "axis_title": dict(width_in=_W_IN - 2 * MARGIN - 2.0, max_h_in=1.1, size_max=D.TYPE["h3"]),
    # Slide enrichie (carte de puces à gauche, visuel à droite, 1re puce en encart) :
    # largeur = carte réduite du visuel ; hauteur = zone au-dessus de l'encart « à retenir ».
    "synthese_categorie": dict(width_in=_SYNTH_AREA_W - 0.48, max_h_in=1.9, size_max=D.TYPE["body"]),
    # Un quadrant SWOT = ~demi-largeur de la zone de contenu ; la hauteur de la
    # zone de PUCES (pas de la carte) = row_h - titre - paddings ≈ 1.9 in sur un
    # deck vierge (cf. _slide_swot) — pas la demi-hauteur brute (~2.2), qui
    # surestimait le budget du repère de ~20 % et rendait le fit-hint trompeur.
    # Matrice SWOT (skill swot-matrix) : cellule teintée = gouttière 0.30 à gauche +
    # bandeau d'axe en haut ; largeur de puces = demi-grille moins paddings.
    "swot_quadrant": dict(width_in=(_W_IN - MARGIN - (MARGIN + 0.30) - 0.22) / 2 - 0.36, max_h_in=1.2, size_max=D.TYPE["small"]),
    # Executive summary (piste F) : panneau pleine largeur (constat + points) et
    # bande cyan « key message » — mêmes contraintes que la slide (cf.
    # _slide_executive_summary), pour un fit-hint fidèle dans l'aperçu.
    # headline / key_message : rendus à taille FIXE (h3) et tronqués à 2 lignes par
    # _slide_executive_summary -> hint en mode size_pt/max_lignes (annonce la
    # troncature, ne PROMET pas de réduction de police que la slide ne fait pas).
    # Constat revue adversariale 2026-07-21 : l'ancien mode max_h_in promettait un
    # shrink inexistant (même classe de dérive que le fit-hint SWOT déjà corrigé).
    "es_headline": dict(width_in=_W_IN - 2 * (MARGIN + 0.3) - 0.48, size_pt=D.TYPE["h3"], max_lignes=2),
    "es_points": dict(width_in=_W_IN - 2 * (MARGIN + 0.3) - 0.48, max_h_in=2.5, size_max=D.TYPE["body"]),
    "es_key_message": dict(width_in=_W_IN - 2 * (MARGIN + 0.3) - 0.48, size_pt=D.TYPE["h3"], max_lignes=2),
    # Difficulté (planche §D.1) : libellé d'une carte, taille fixe body, tronqué à
    # 3 lignes par _slide_difficultes -> hint size_pt/max_lignes (honnête). Largeur
    # réduite du chip de rang à gauche (2*pad + rang_w 0.46 + gap 0.16 = 0.98).
    "difficulty_label": dict(width_in=_W_IN - 2 * (MARGIN + 0.3) - 0.98, size_pt=D.TYPE["body"], max_lignes=3),
}


def field_fit_hint(field_key: str, text: str) -> str:
    """Message court indiquant comment `text` sera rendu à l'export pour le
    champ `field_key` (police retenue, nombre de lignes, troncature
    éventuelle) — s'appuie sur les mêmes fonctions d'ajustement
    (`D.ajuster_police` / `D.tronquer_a_lignes` / `D.estimer_lignes`) que le
    générateur, appliquées aux contraintes de forme de `FIELD_SHAPE` (voir
    note du module). Chaîne vide si le champ est inconnu ou vide — pas de
    repère à afficher plutôt qu'un repère trompeur."""
    spec = FIELD_SHAPE.get(field_key)
    text = (text or "").strip()
    if spec is None or not text:
        return ""

    width_in = spec["width_in"]

    if "max_lignes" in spec:
        size = spec["size_pt"]
        lignes = D.estimer_lignes(text, width_in, size)
        if lignes > spec["max_lignes"]:
            return f"⚠ trop long — sera tronqué à {spec['max_lignes']} lignes à l'export"
        return f"{lignes} ligne(s) à {size:.0f}pt à l'export"

    max_h_in = spec["max_h_in"]
    size_max = spec.get("size_max", D.TYPE["body"])
    size_min = D.TYPE["tiny"]

    def budget_ok(taille, lignes_max):
        # Même réserve d'une demi-ligne que le budget réel de _add_bulleted_text
        # (sinon le hint annonce à la frontière une taille que l'export réduit).
        return lignes_max * _per_line_height_in(taille) <= max_h_in - 0.5 * _per_line_height_in(taille)

    size, lignes = D.ajuster_police([text], width_in, size_max, size_min, budget_ok=budget_ok)
    if lignes * _per_line_height_in(size) > max_h_in:
        return f"⚠ très long — sera réduit à {size:.0f}pt et tronqué à l'export"
    if size < size_max - 0.5:
        return f"{lignes} ligne(s) — police réduite à {size:.0f}pt pour tenir à l'export"
    return f"{lignes} ligne(s) à {size:.0f}pt à l'export"


def _dims(prs: Presentation) -> tuple[float, float]:
    return Emu(prs.slide_width).inches, Emu(prs.slide_height).inches


def _clear_slides(prs: Presentation) -> None:
    """Retire toutes les slides d'une présentation chargée depuis un template
    client — on ne veut hériter que masters/layouts/thème. Sans ça, un
    template qui est un vrai exemple de deck (pas un .potx vierge) ferait
    apparaître tout son contenu d'origine avant le nôtre. `python-pptx`
    n'expose pas de suppression de slide côté API publique ; on vide
    directement la liste XML des slides — mais il faut aussi lâcher la
    relation (r:id) de chaque slide sur la part présentation, sans quoi le
    fichier réserialisé contient des relations pointant vers des parts
    devenues orphelines : invisible pour python-pptx (parseur tolérant),
    mais PowerPoint refuse ensuite d'ouvrir le fichier (constaté via
    l'automation COM — l'export semblait « valide » côté tests avant ça)."""
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.get(qn("r:id")))
        sld_id_lst.remove(sld_id)


def _has_title_placeholder(layout) -> bool:
    try:
        return any(
            ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
            for ph in layout.placeholders
        )
    except Exception:
        return False


def _pick_layout(prs: Presentation, preferred: int = 6):
    """Choisit un layout pour du contenu personnalisé, en respectant au mieux
    le template injecté (point 6) : un layout nommé "Title Only"/"Section" (le
    logo/footer du master survit, pas de placeholder de corps qui entrerait en
    collision avec notre mise en page), sinon le layout avec titre le moins
    chargé en autres placeholders, sinon le comportement historique (repli
    toujours valide même sur un template client mal structuré)."""
    layouts = list(prs.slide_layouts)
    # « titre seul » = le layout de contenu OCTO (idx0 titre, garde logo/pied/n° de slide).
    for kw in ("titre seul", "title only", "section"):
        for layout in layouts:
            if kw in (layout.name or "").lower() and _has_title_placeholder(layout):
                return layout
    with_title = [l for l in layouts if _has_title_placeholder(l)]
    if with_title:
        return min(with_title, key=lambda l: len(l.placeholders))
    return layouts[preferred] if preferred < len(layouts) else layouts[-1]


def _new_slide(prs: Presentation, title: str, max_title_lines: int = 2):
    """Crée une slide de contenu et pose son titre. Renvoie
    `(slide, w_in, h_in, content_top)` — `content_top` est calculé à partir
    de la position/hauteur réelle du placeholder de titre (natif du
    template) et du nombre de lignes qu'occupera effectivement `title` une
    fois replié, plutôt qu'une constante suppposant un titre sur une seule
    ligne : un titre de longueur normale (~50 caractères) suffit à passer
    sur 2 lignes et, avec une position de contenu figée, à chevaucher la
    zone en dessous — ce qu'une constante ne peut pas anticiper."""
    slide = prs.slides.add_slide(_pick_layout(prs))
    w_in, h_in = _dims(prs)
    title_shape = slide.shapes.title
    # Le placeholder de titre natif hérite position/police du template —
    # préféré à une zone de texte dessinée à la main dès qu'il existe sur le
    # layout choisi. On fige sa taille de police sur D.TYPE["title"] (au lieu
    # de laisser le style hérité, potentiellement bien plus grand) : ça reste
    # cohérent avec l'unique échelle typographique du reste du deck, et ça
    # rend le nombre de lignes prévisible (donc calculable) plutôt que soumis
    # à un style de thème inconnu.
    if title_shape is not None:
        if getattr(prs, "_i2d_synthetic", False):
            # Présentation vierge (pas de template client) : le placeholder
            # de titre hérité du modèle par défaut de python-pptx est
            # dimensionné pour un slide 10x7.5in (4:3) — trop étroit une fois
            # la slide passée en 16:9. On le repositionne explicitement sur
            # CETTE slide (jamais sur le layout/master : leurs placeholders
            # sont résolus par héritage et se sont révélés instables à muter
            # directement avec python-pptx — cf. essai précédent).
            title_shape.left = Inches(MARGIN)
            title_shape.top = Inches(0.3)
            title_shape.width = Inches(w_in - 2 * MARGIN)
            title_shape.height = Inches(1.1)
        title_w_in = Emu(title_shape.width).inches if title_shape.width is not None else (w_in - 2 * MARGIN)
        title_top_in = Emu(title_shape.top).inches if title_shape.top is not None else 0.3
        size = D.TYPE["title"]
        max_lignes = max_title_lines
        if D.estimer_lignes(title, title_w_in, size) > max_lignes:
            title = D.tronquer_a_lignes(title, title_w_in, size, max_lignes)
        title_shape.text = title
        tf = title_shape.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        for p in tf.paragraphs:
            for run in p.runs:
                run.font.size = Pt(size)
                run.font.bold = True
                if D.POLICE:  # police du deck (thème), pas l'Outfit hérité du layout
                    run.font.name = D.POLICE
        lignes = D.estimer_lignes(title, title_w_in, size)
        needed_h = lignes * _per_line_height_in(size) + 0.15
        # Pas de barre d'accent avant le titre : les decks OCTO réels (VSCode4) n'en
        # ont pas — titre navy + logo suffisent (retrait demandé 2026-07-22, charte VSCode4).
        # Hauteur RÉELLE du texte de titre (plancher 0.55), pas la boîte du
        # placeholder — indépendant de la hauteur de boîte du template. Réserve
        # connue (revue adversariale) : un template CLIENT à boîte de titre haute
        # ancrée middle/bottom pourrait voir le contenu remonter dans sa zone de
        # titre — à re-vérifier au premier template client réel.
        content_top = title_top_in + max(needed_h, 0.55) + 0.25
    else:
        D.add_text(
            slide, MARGIN, 0.35, w_in - 2 * MARGIN, 0.7,
            [(title, {"size": D.TYPE["title"], "bold": True, "color": D.INK})],
        )
        content_top = 1.4
    return slide, w_in, h_in, content_top


def _bullet_lines(text: str) -> list[str]:
    lines = []
    for raw in (text or "").splitlines():
        line = raw.strip().lstrip("-•").strip()
        if line:
            lines.append(line)
    return lines


def _per_line_height_in(size_pt: float) -> float:
    # Calibrage empirique (skill pptx-deck) : ~0.17in/ligne à 10.5pt, un peu
    # plus large ici pour couvrir l'espacement inter-puces (space_after).
    return size_pt * 0.017 + 4 / 72


def _add_bulleted_text(
    slide, l, t, w, h, text: str, size: float | None = None,
    anchor=MSO_ANCHOR.TOP, size_max: float = D.TYPE["body"], size_min: float = D.TYPE["tiny"],
    paginate: bool = False,
) -> list[str]:
    """Pose une liste à puces dans la zone donnée. Par défaut (`paginate=False`,
    comportement historique) la police est réduite jusqu'à `size_min` pour
    tenter de tout faire tenir, sans garantie — un texte très long peut
    déborder silencieusement de sa zone (indétectable par
    `D.verifier_geometrie`, qui ne voit que les bords des formes).

    `paginate=True` change la garantie : si même `size_min` ne suffit pas, les
    puces qui ne tiennent pas sont retenues (pas rendues) et renvoyées à
    l'appelant plutôt que de déborder — à charge pour lui de les poser sur
    une slide de continuation (voir `_emit_bullet_overflow`). Renvoie la
    liste des puces non rendues (vide si tout tient)."""
    lines = _bullet_lines(text) or ["—"]

    def budget_ok(taille, _lignes_max):
        # Le budget inclut la MÊME réserve d'une demi-ligne que la pagination —
        # sinon l'ajusteur valide une taille que la pagination recoupe ensuite
        # (explosion de slides de suite constatée le 2026-07-22) : on préfère
        # rétrécir la police d'un cran que couper une puce.
        total = sum(D.estimer_lignes(line, w, taille) for line in lines)
        return total * _per_line_height_in(taille) <= h - 0.5 * _per_line_height_in(taille)

    if size is None:
        size, _ = D.ajuster_police(lines, w, size_max, size_min, budget_ok=budget_ok)

    overflow: list[str] = []
    if paginate:
        # Capacité MINORÉE d'une demi-ligne : l'estimation de repli est optimiste
        # pour du français — sans réserve, le dernier bloc d'une carte sortait du
        # cadre au vrai rendu (défaut récurrent 2026-07-22). Et un item SEUL plus
        # haut que la zone est insécable pour paginer_items — il débordait en
        # silence (attrapé par verifier_debordements_texte) : désormais tronqué à
        # l'ellipse ici, son texte COMPLET partant sur la slide de suite.
        line_h = _per_line_height_in(size)
        capacite = max(line_h, h - 0.5 * line_h)
        if sum(D.estimer_lignes(li, w, size) for li in lines) * line_h > capacite:
            pages = D.paginer_items(
                lines, lambda li: D.estimer_lignes(li, w, size) * line_h,
                capacite_in=capacite,
            )
            lines = pages[0]
            overflow = [li for page in pages[1:] for li in page]
            if len(lines) == 1 and D.estimer_lignes(lines[0], w, size) * line_h > capacite:
                complet = lines[0]
                lines = [D.tronquer_a_lignes(complet, w, size, max(1, int(capacite / line_h)))]
                overflow = [complet] + overflow

    paragraphs = [(f"•  {line}", {"size": size, "color": D.INK, "space_after": 4}) for line in lines]

    if anchor == MSO_ANCHOR.MIDDLE:
        total_lines = sum(D.estimer_lignes(line, w, size) for line in lines)
        # +0.5 ligne de marge dans la boîte centrée — même logique que partout :
        # l'estimation nominale est optimiste, la boîte exacte faisait peindre la
        # dernière ligne hors boîte (constat verifier_debordements_texte).
        content_h = min(h, (total_lines + 0.5) * _per_line_height_in(size))
        box_t = t + max(0.0, (h - content_h) / 2)
        D.add_text(slide, l, box_t, w, content_h, paragraphs)
    else:
        D.add_text(slide, l, t, w, h, paragraphs)
    return overflow


def _emit_bullet_overflow(prs: Presentation, base_title: str, field_label: str, overflow_lines: list[str]) -> None:
    """Pose les puces qui n'ont pas tenu sur la slide d'origine (voir
    `_add_bulleted_text(paginate=True)`) sur une ou plusieurs slides de
    continuation pleine largeur — chacune dispose de bien plus d'espace que
    la colonne étroite d'origine, donc peut se voir attribuer sa propre
    police (recalculée, pas figée à `size_min`)."""
    remaining = "\n".join(overflow_lines)
    page_no = 1
    while remaining:
        suffix = f" {page_no}" if page_no > 1 else ""
        slide, w_in, h_in, top = _new_slide(prs, f"{base_title} (suite — {field_label}){suffix}")
        w = w_in - 2 * (MARGIN + 0.3)
        h = h_in - top - 0.5
        overflow = _add_bulleted_text(slide, MARGIN + 0.3, top, w, h, remaining, paginate=True)
        nouveau = "\n".join(overflow)
        if nouveau == remaining:
            # Garde de PROGRESSION (revue adversariale) : une puce insécable plus
            # haute qu'une slide de suite entière revient intégralement en
            # overflow (le chemin de troncature repousse le texte COMPLET) —
            # sans cette garde, boucle infinie + slides sans fin. La slide qui
            # vient d'être posée montre déjà tout ce qui tient, à l'ellipse.
            break
        remaining = nouveau
        page_no += 1


# --------------------------------------------------------------------------- #
# Slides
# --------------------------------------------------------------------------- #
def _add_measured_field(
    slide, l, t, w, label: str, text: str, max_h: float,
    size_max: float = D.TYPE["body"], size_min: float = D.TYPE["tiny"],
    bold: bool = False, italic: bool = False,
) -> float:
    """Pose un libellé (petit, gras, discret) puis son contenu juste en
    dessous, en adaptant la taille de police du contenu à `max_h`
    (D.ajuster_police) et en tronquant en tout dernier recours — jamais de
    débordement dans le bloc suivant même avec une réponse d'entretien très
    longue. Renvoie la hauteur réellement occupée (libellé + contenu), à
    utiliser pour empiler le bloc suivant à la bonne position."""
    label_h = 0.3
    D.add_text(slide, l, t, w, label_h, [(label, {"size": D.TYPE["small"], "bold": True, "color": D.MUTED})])
    body = ((text or "").strip()) or "—"
    body_max_h = max(0.2, max_h - label_h)

    def budget_ok(taille, lignes_max):
        return lignes_max * _per_line_height_in(taille) <= body_max_h

    size, lignes_max = D.ajuster_police([body], w, size_max, size_min, budget_ok=budget_ok)
    if lignes_max * _per_line_height_in(size) > body_max_h:
        max_lignes = max(1, int(body_max_h / _per_line_height_in(size)))
        body = D.tronquer_a_lignes(body, w, size, max_lignes)
        lignes_max = max_lignes
    # +0.6 ligne de marge : la boîte était dimensionnée à l'estimation EXACTE —
    # au vrai repli PowerPoint (un peu plus gourmand), la dernière ligne sortait
    # du cadre (constat verifier_debordements_texte sur les fiches, 2026-07-22).
    body_h = (lignes_max + 0.6) * _per_line_height_in(size)
    D.add_text(slide, l, t + label_h, w, body_h, [(body, {"size": size, "bold": bold, "italic": italic, "color": D.INK})])
    return label_h + body_h


def _layout_by_name(prs: Presentation, *keywords: str):
    """Premier layout dont le nom contient l'un des mots-clés (insensible à la
    casse) — repérage robuste des layouts de marque OCTO (« 40 - Couverture »,
    « 50 - Chapitre ») par nom plutôt que par indice. None si aucun ne matche."""
    for layout in prs.slide_layouts:
        name = (layout.name or "").lower()
        if any(kw in name for kw in keywords):
            return layout
    return None


# Refonte P2 — structure narrative : 4 chapitres. Chaque section du deck appartient
# à un chapitre ; un intercalaire (layout « 50 - Chapitre ») ouvre chaque chapitre
# qui a du contenu, et le sommaire quali regroupe les sections sous ces intitulés
# narratifs (couleur = repère de navigation, reprise sur l'intercalaire).
_CH_RETENIR, _CH_DIAGNOSTIC, _CH_PAROLE, _CH_TRAJECTOIRE = 0, 1, 2, 3
# (intitulé, couleur, scène image de l'intercalaire — cf. _slide_chapitre P3).
# (intitulé, couleur, scène, sous-titre italique — format d'intercalaire VSCode3 :
# le placeholder de titre porte titre + sous-titre, cf. _slide_chapitre).
_CHAPITRES = [
    ("Ce qu'il faut retenir", "#00D2DD", "sunset", "L'essentiel de la mission, en une page."),
    ("Le diagnostic", "#0E2356", "mountains", "Ce que les entretiens révèlent, sans détour."),
    ("La parole des équipes", "#138086", "forest", "Les mots des acteurs, tels quels."),
    ("La trajectoire proposée", "#6a3d9a", "ocean", "Où aller, et par quoi commencer."),
]


def _sans_puce(paragraph) -> None:
    """Retire l'indentation de puce héritée (marL/indent) et la puce elle-même —
    reproduit tel quel le helper du générateur de référence VSCode3. Cause réelle
    du « 01 » qui wrappe dans le petit encart numéro du layout Chapitre : le style
    de liste hérité pose marL=0.5in dans un encart de ~0.55in. python-pptx n'expose
    pas ces attributs -> manipulation XML directe."""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", "0")
    pPr.set("indent", "0")
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    # Forcer explicitement l'absence de puce : notre template hérite un caractère de
    # puce à un niveau que le retrait ci-dessus ne couvre pas (un ◉ résiduel
    # apparaissait avant le numéro) — buNone le neutralise à coup sûr.
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))
# Scène -> requête photo Openverse (vraies photos CC0, comme les decks OCTO réels /
# VSCode3-4). Repli sur la génération procédurale `nature_images` (nom = la scène).
_SCENE_REQUETE = {
    "sunset": "sunset sky",
    "mountains": "mountains landscape",
    "forest": "green forest sunlight",
    "ocean": "turquoise water",
}


def _slide_cover(prs: Presentation, mission: Mission) -> None:
    """Couverture : sur un template OCTO, remplit les placeholders du layout de
    marque « 40 - Couverture » (titre/sous-titre/date) — la mise en forme (police
    Outfit, tailles, éventuelle photo de couverture) vient du template. Repli sur un
    titre dessiné centré si le layout n'existe pas (deck synthétique sans marque)."""
    subtitle = "Synthèse transverse & recommandations"
    date_str = datetime.now(timezone.utc).strftime("%d/%m/%Y")
    cover = _layout_by_name(prs, "couverture", "cover")
    if cover is not None:
        slide = prs.slides.add_slide(cover)
        phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
        if 0 in phs:
            phs[0].text_frame.text = mission.name
            D.appliquer_police(phs[0].text_frame)
        if 1 in phs:
            phs[1].text_frame.text = subtitle
            D.appliquer_police(phs[1].text_frame)
        # idx2 = « OCTO Technology » : le texte d'invite du layout ne REND pas —
        # laissé vide, la pastille affichait « | date » avec un trou à gauche
        # (constat utilisateur 2026-07-22) ; à remplir explicitement. idx3 = date.
        if 2 in phs:
            phs[2].text_frame.text = "OCTO Technology"
            D.appliquer_police(phs[2].text_frame)
        if 3 in phs:
            phs[3].text_frame.text = date_str
            D.appliquer_police(phs[3].text_frame)
        return
    _slide_title(prs, mission, subtitle, date_str)


def _slide_title(prs: Presentation, mission: Mission, subtitle: str, date_str: str) -> None:
    slide = prs.slides.add_slide(_pick_layout(prs))
    w_in, h_in = _dims(prs)
    title_w = w_in - 2.0
    size = 32
    # Le sous-titre se place sous la hauteur RÉELLE du titre (un nom long passe sur
    # 2 lignes sur une slide étroite type OCTO 10in) plutôt qu'à un offset fixe qui
    # supposait un titre sur 1 ligne — sinon la 2e ligne chevauche le sous-titre.
    title_lines = max(1, min(3, D.estimer_lignes(mission.name, title_w - 0.2, size)))
    title_h = title_lines * _per_line_height_in(size)
    cy = h_in * 0.38
    D.add_text(
        slide, 1.0, cy, title_w, title_h,
        [(mission.name, {"size": size, "bold": True, "color": D.INK, "align": PP_ALIGN.CENTER})],
    )
    y = cy + title_h + 0.15
    D.add_text(
        slide, 1.0, y, title_w, 0.5,
        [(subtitle, {"size": D.TYPE["h2"], "color": D.MUTED, "align": PP_ALIGN.CENTER})],
    )
    D.add_text(
        slide, 1.0, y + 0.55, title_w, 0.4,
        [(date_str, {"size": D.TYPE["small"], "color": D.MUTED, "align": PP_ALIGN.CENTER})],
    )


def _slide_sommaire(prs: Presentation, ch_sections: list[list[str]]) -> None:
    """Sommaire quali (P2) : les chapitres AYANT du contenu, chacun avec sa pastille
    couleur + intitulé narratif + les sections qu'il regroupe (repère de navigation,
    couleur reprise sur l'intercalaire) — au lieu d'une liste plate de sections."""
    slide, w_in, h_in, top = _new_slide(prs, "Sommaire")
    active = [ci for ci, subs in enumerate(ch_sections) if subs]
    # Grille 2×2 de badges GOUTTE (teardrop) à contour — signature du sommaire des
    # decks OCTO réels (VSCode4) : numéro dans la goutte, intitulé narratif + sections
    # en regard. Remplissage colonne par colonne (01,02 à gauche ; 03,04 à droite).
    area_t = top + 0.2
    row_h = (h_in - 0.5 - area_t) / 2
    col_w = (w_in - 2 * MARGIN) / 2
    d = 0.92  # diamètre du badge goutte
    for idx, ci in enumerate(active):
        subs = ch_sections[ci]
        label, color = _CHAPITRES[ci][0], _CHAPITRES[ci][1]
        col, r = idx // 2, idx % 2
        cell_x = MARGIN + col * col_w
        cell_y = area_t + r * row_h
        bx = cell_x + 0.15
        by = cell_y + (row_h - d) / 2
        D.add_teardrop(slide, bx, by, d, f"{idx + 1:02d}", color, size=D.TYPE["h2"])
        tx = bx + d + 0.3
        tw = col_w - (bx - cell_x) - d - 0.3 - 0.2
        D.add_text(
            slide, tx, cell_y, max(1.0, tw), row_h,
            [
                (label, {"size": D.TYPE["h3"], "bold": True, "color": D.INK, "space_after": 4}),
                (" · ".join(subs), {"size": D.TYPE["small"], "color": D.MUTED}),
            ],
            anchor=MSO_ANCHOR.MIDDLE,
        )


def _find_teardrop_frame(shapes):
    """`(left, top, width, height, geom)` du cadre photo teardrop d'un layout
    (le layout « 50 - Chapitre » place son cadre en top-level, pas dans un groupe),
    ou None. Même principe que pptx-framed-image.frame_geometry, cas non groupé."""
    for sh in shapes:
        spPr = getattr(sh._element, "spPr", None)
        if spPr is None:
            continue
        g = spPr.find(qn("a:prstGeom"))
        if g is not None and g.get("prst") == "teardrop":
            return sh.left, sh.top, sh.width, sh.height, g
    return None


def _resoudre_image_cachee(base: str, scene: str, seed: int, aspect: float,
                           px_w: int, px_h: int, requete: str):
    """Résout l'image d'une zone : vraie photo Openverse CC0 si le fetch est permis
    et réussit, repli procédural sinon — avec un cache SÉPARÉ PAR SOURCE
    (`…_photo.png` / `…_proc.png`). Leçon 2026-07-22 (demande « ne pas générer des
    images ») : avec un seul slot de cache, un repli procédural écrit pendant un run
    de tests (PPTX_NO_PHOTO_FETCH=1, posé par conftest) restait servi À VIE par le
    serveur — le slot photo n'était jamais retenté. Ici un échec de fetch remplit le
    slot `_proc` sans jamais occuper `_photo` : le prochain run en ligne re-tente la
    vraie photo. Renvoie le Path de l'image utilisable."""
    import hashlib
    _IMG_CACHE.mkdir(parents=True, exist_ok=True)
    no_fetch = os.environ.get("PPTX_NO_PHOTO_FETCH") == "1"
    # Le slot photo dépend AUSSI de la requête : affiner une requête doit
    # re-déclencher un fetch, pas resservir l'ancien résultat mis en cache.
    qh = hashlib.md5(requete.encode("utf-8")).hexdigest()[:6]
    photo = _IMG_CACHE / f"{base}_{qh}_photo.png"
    proc = _IMG_CACHE / f"{base}_proc.png"
    if not no_fetch:
        if photo.exists():
            return photo
        # Échelle de retry (constat 2026-07-22, slides 7/10 restées procédurales) :
        # les échecs Openverse sont INTERMITTENTS (SSL sporadique) → 2 tentatives ;
        # et une requête précise peut n'avoir AUCUN résultat CC0 pour cet aspect →
        # repli sur la requête simplifiée « {scene} photography » avant d'abandonner.
        ar = "wide" if aspect > 1.15 else "tall" if aspect < 0.85 else "square"
        variantes = [requete]
        simple = f"{scene} photography"
        if simple != requete:
            variantes.append(simple)
        for req in variantes:
            for _tentative in range(2):
                try:
                    brut = _IMG_CACHE / f"_brut_{scene}_{seed}.jpg"
                    _stock_images.fetch_to(str(brut), req, seed=seed, aspect_ratio=ar)
                    _cover_crop_to_aspect(str(brut), str(photo), aspect)
                    return photo
                except Exception:
                    continue  # réseau/API KO : tentative/variante suivante
        # tout a échoué : repli procédural ci-dessous, slot photo intact
    if not proc.exists():
        _nature_images.generate_to(str(proc), scene, px_w, px_h, seed=seed)
    return proc


def _remplir_cadre_chapitre(slide, cadre, scene: str, seed: int = 0) -> None:
    """Remplit le cadre teardrop d'un intercalaire avec une image à l'aspect exact
    du cadre — vraie photo Openverse via _resoudre_image_cachee (repli procédural).
    Silencieux sur échec : l'intercalaire reste lisible sans image."""
    if not _FRAMED_OK or cadre is None:
        return
    try:
        left, top, width, height, geom = cadre
        aspect = Emu(width).inches / Emu(height).inches
        px_w = 900
        px_h = max(1, int(round(px_w / aspect)))
        path = _resoudre_image_cachee(
            f"{scene}_{seed}_{px_w}x{px_h}", scene, seed, aspect, px_w, px_h,
            requete=_SCENE_REQUETE.get(scene, scene),
        )
        _place_image_in_frame(slide, str(path), left, top, width, height, geom=geom)
    except Exception:
        pass  # repli : intercalaire sans image, jamais un export cassé


def _slide_chapitre(prs: Presentation, numero: int, titre: str, color: str,
                    scene: str | None = None, sous_titre: str | None = None) -> None:
    """Intercalaire de chapitre — reproduit le format du générateur de référence
    VSCode3 (« 50 - Chapitre ») : idx0 = titre coloré + sous-titre italique gris ;
    idx1 = numéro DANS l'encart logo (17pt, marges à zéro, centré, sans puce — ce qui
    empêche « 01 » de wrapper) ; cadre photo teardrop rempli. Repli dessiné sinon."""
    layout = _layout_by_name(prs, "chapitre") if _FRAMED_OK else None
    if layout is not None:
        slide = prs.slides.add_slide(layout)
        phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
        if 0 in phs:
            # idx0 = titre (couleur de chapitre) + sous-titre italique gris.
            tf0 = phs[0].text_frame
            tf0.text = titre
            for r in tf0.paragraphs[0].runs:
                r.font.color.rgb = D.rgb(color)
            if sous_titre:
                p2 = tf0.add_paragraph()
                p2.text = sous_titre
                for r in p2.runs:
                    r.font.size = Pt(D.TYPE["small"])
                    r.font.italic = True
                    r.font.color.rgb = D.rgb(D.MUTED)
            D.appliquer_police(tf0)  # police du deck, pas l'Outfit du layout
        if 1 in phs:
            # idx1 = numéro DANS l'encart logo (format VSCode3) : marges à zéro + 17pt
            # + centré + sans puce — ce qui empêche « 01 » de wrapper dans le petit
            # encart. Reproduit à l'identique du générateur de référence VSCode3.
            tf1 = phs[1].text_frame
            tf1.text = f"{numero:02d}"
            tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = 0
            tf1.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in tf1.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                _sans_puce(p)
                for r in p.runs:
                    r.font.size = Pt(17)
                    r.font.color.rgb = D.rgb(color)
            D.appliquer_police(tf1)  # police du deck, pas l'Outfit du layout
        _remplir_cadre_chapitre(slide, _find_teardrop_frame(slide.slide_layout.shapes),
                                scene or "mountains")
        return
    slide = prs.slides.add_slide(_pick_layout(prs))
    w_in, h_in = _dims(prs)
    cy = h_in * 0.30
    D.add_text(
        slide, MARGIN + 0.3, cy, 3.0, 1.3,
        [(f"{numero:02d}", {"size": D.TYPE["kpi"], "bold": True, "color": color})],
    )
    ty = cy + 1.35
    D.add_rect(slide, MARGIN + 0.36, ty, 1.2, 0.06, fill=color)
    D.add_text(
        slide, MARGIN + 0.3, ty + 0.18, w_in - 2 * MARGIN - 0.6, 1.0,
        [(titre, {"size": D.TYPE["title"], "bold": True, "color": D.INK})],
    )


# Enrichissement synthèse (ask design 2026-07-22) : pattern claim + visuel + encart
# des decks OCTO réels (VSCode4). Scène/requête photo par catégorie (repli procédural
# offline, comme les têtes de chapitre) — clé = libellé exact passé par build_presentation.
# Scènes NATURE (comme les têtes de chapitre) : rendu procédural fiable hors ligne
# ET vraie photo Openverse en prod — cohérent avec l'imagerie de marque du deck.
# (scène, requête photo, seed distinct pour varier des intercalaires).
_SYNTHESE_VISUEL = {
    # « photography » dans la requête : Openverse mélange photos et illustrations —
    # sans ce biais, une requête générique peut renvoyer un clipart (constat
    # pptx-verify 2026-07-22 : « mountains landscape » → illustration Fuji).
    "Contexte": ("mountains", "mountain landscape photography", 11),
    "Culture & ADN": ("forest", "forest sunlight nature photography", 12),
    "Forces & succès": ("sunset", "sunset sky photography", 13),
    "Points d'amélioration": ("ocean", "ocean waves photography", 14),
    "Aspirations (baguette magique)": ("sunset", "sunrise horizon photography", 15),
}


def _slide_synthese_categorie(prs: Presentation, label: str, content: str) -> None:
    """Slide de catégorie de synthèse, ENRICHIE (claim + visuel + encart) : puces à
    gauche dans une carte, photo métier à droite, 1re puce promue en encart « à
    retenir » cyan en bas — au lieu d'un titre + puces sur fond vide. Repli propre
    (carte pleine largeur, pas d'encart) si l'infra image manque ou si la catégorie
    n'a qu'une puce. Même pattern que _slide_executive_summary."""
    slide, w_in, h_in, top = _new_slide(prs, f"Synthèse globale — {label}")
    accent = (D.theme_colors(prs).get("accent3") or "#00D2DD")  # cyan OCTO
    area_l = MARGIN + 0.3
    has_vis = _FRAMED_OK
    vis_w = _SYNTH_VIS_W
    vis_l = w_in - MARGIN - vis_w
    area_w = (vis_l - 0.3 - area_l) if has_vis else (w_in - 2 * (MARGIN + 0.3))
    pad = 0.24
    band_h, band_gap = 0.9, 0.3
    band_t = h_in - 0.5 - band_h

    lines = _bullet_lines(content) or ["—"]
    # 1re puce -> encart « à retenir » si au moins 2 puces (sinon tout dans la carte).
    retenir = lines[0] if len(lines) >= 2 else None
    rest = lines[1:] if retenir else lines

    zone_bottom = (band_t - band_gap) if retenir else (h_in - 0.5)
    avail = max(0.0, zone_bottom - top)

    body = D.TYPE["body"]
    rest_text = "\n".join(rest) or "—"
    # La carte occupe TOUTE la zone (même hauteur que le visuel à droite → colonnes
    # équilibrées, pas de vide sous une carte trop courte) ; puces centrées verticalement.
    card_h = avail
    D.add_card(slide, area_l, top, area_w, card_h, accent)
    _add_bulleted_text(
        slide, area_l + pad, top + pad, area_w - 2 * pad, max(0.0, card_h - 2 * pad),
        rest_text, anchor=MSO_ANCHOR.MIDDLE, size_max=body, size_min=D.TYPE["small"],
        paginate=True,
    )

    if has_vis:
        scene, requete, seed = _SYNTHESE_VISUEL.get(label, ("mountains", "mountains landscape", 11))
        if not _image_dans_zone(slide, vis_l, top, vis_w, avail, scene, requete, seed=seed):
            D.add_rect(slide, vis_l, top, vis_w, avail, fill=accent, rounded=True, radius=0.06)

    if retenir:
        # Encart « à retenir » gris (même composant add_encart que l'executive summary
        # — cohérence de composant §5, sobriété §3/§7, motif VSCode4).
        msg = D.tronquer_a_lignes(retenir, area_w - 0.6, D.TYPE["h3"], 2)
        D.add_encart(slide, area_l, band_t, area_w, band_h, msg, accent=accent)


# SWOT : Forces/Faiblesses = interne (vert/rouge), Opportunités/Menaces =
# externe (bleu/ambre). Couleurs sémantiques prises dans D.PALETTE (design
# system : différenciation par liseré de carte, pas de dégradé/ombre).
_SWOT_QUADRANTS = [
    ("forces", "Forces", "#1e6b34"),
    ("faiblesses", "Faiblesses", "#b3261e"),
    ("opportunites", "Opportunités", "#2c5cc5"),
    ("menaces", "Menaces", "#b8860b"),
]

# Badge-icône par quadrant : flèches directionnelles (bloc Arrows, monochrome,
# rendu fiable — cf. l'usage de « → » sur les decks OCTO réels VSCode4). Sémantique
# de la grille : interne haut/bas (↑ force / ↓ faiblesse), externe haut/bas
# (↗ opportunité / ↘ menace). bold=False au badge (certains glyphes « tofu » en gras).
_SWOT_ICONS = {"forces": "↑", "faiblesses": "↓", "opportunites": "↗", "menaces": "↘"}


def _image_dans_zone(slide, left, top, width, height, scene: str, requete: str,
                     seed: int = 0) -> bool:
    """Pose une photo (Openverse CC0, repli procédural offline) cover-croppée à
    l'aspect de la zone, dans un rectangle — pattern « claim + visuel » (P3b, repéré
    sur les decks OCTO réels VSCode4). PPTX_NO_PHOTO_FETCH=1 force le procédural
    (tests). Renvoie True si posée, False sinon (silencieux — jamais un export cassé)."""
    if not _FRAMED_OK:
        return False
    try:
        aspect = width / height
        px_w = 900
        px_h = max(1, int(round(px_w / aspect)))
        path = _resoudre_image_cachee(
            f"zone_{scene}_{seed}_{px_w}x{px_h}", scene, seed, aspect, px_w, px_h,
            requete=requete,
        )
        pic = slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                       Inches(width), Inches(height))
        _clip_octo_frame(pic)  # cadre OCTO round2DiagRect (format image VSCode3)
        return True
    except Exception:
        return False


def _clip_octo_frame(pic) -> None:
    """Clippe une image au cadre OCTO « round2DiagRect » (rectangle à 2 coins diagonaux
    arrondis, 2 vifs) — format « image encadrée » des slides de contenu VSCode3, au lieu
    d'un rectangle plat. python-pptx n'expose pas prstGeom sur une image -> XML direct."""
    try:
        spPr = pic._element.spPr
        for tag in ("a:prstGeom", "a:custGeom"):
            for el in spPr.findall(qn(tag)):
                spPr.remove(el)
        geom = spPr.makeelement(qn("a:prstGeom"), {"prst": "round2DiagRect"})
        av = geom.makeelement(qn("a:avLst"), {})
        for name, val in (("adj1", "33000"), ("adj2", "0")):
            av.append(av.makeelement(qn("a:gd"), {"name": name, "fmla": f"val {val}"}))
        geom.append(av)
        spPr.append(geom)
    except Exception:
        pass  # image en rectangle plat si le clip échoue — jamais un export cassé


# Couleurs des cartes de points clés de l'exec summary (format VSCode3 :
# Doctrine bleu / Méthode vert / Maturité ambre / Posture rouge).
_EXEC_CARD_COLORS = ["#2c5cc5", "#1e6b34", "#b8860b", "#b3261e"]


def _slide_executive_summary(prs: Presentation, es) -> None:
    """Slide d'ouverture « Executive Summary » (piste F restitution, 2026-07-21) :
    un panneau constat + points clés, et une bande cyan « key message » (le
    so-what) en bas — pattern relevé sur les vraies restitutions OCTO (Executive
    Summary + bande de message à retenir), cf.
    docs/reflexions/restitution-mission.md §F. Placée juste après le sommaire."""
    slide, w_in, h_in, top = _new_slide(prs, "Executive Summary")
    area_l = MARGIN + 0.3
    area_w = w_in - 2 * (MARGIN + 0.3)
    headline = (getattr(es, "headline", "") or "").strip()
    key_message = (getattr(es, "key_message", "") or "").strip()
    points = _bullet_lines(getattr(es, "points", "") or "")

    # Le GROUPE (claim + sous-claim + cartes) est CENTRÉ verticalement dans la
    # bande — claim en haut + cartes plaquées en bas laissaient un grand vide au
    # milieu (constat utilisateur 2026-07-22 « contenu mieux centré »). On mesure
    # donc chaque bloc AVANT de dessiner.
    hl = km = ""
    hl_h = km_h = cards_h = 0.0
    gap_claim = 0.14
    gap_cards = 0.4
    if headline:
        hl = D.tronquer_a_lignes(headline, area_w, D.TYPE["h2"], 2)
        hl_h = D.estimer_lignes(hl, area_w, D.TYPE["h2"]) * _per_line_height_in(D.TYPE["h2"])
    if key_message:
        km = D.tronquer_a_lignes(key_message, area_w, D.TYPE["body"], 2)
        km_h = D.estimer_lignes(km, area_w, D.TYPE["body"]) * _per_line_height_in(D.TYPE["body"])
    n = min(len(points), 4)
    gap = 0.2
    cpad = 0.18
    col_w = (area_w - gap * (n - 1)) / n if n else area_w
    if points:
        lh = _per_line_height_in(D.TYPE["small"])
        max_lines = max(2, max(D.estimer_lignes(pt, col_w - 2 * cpad, D.TYPE["small"])
                               for pt in points[:n]))
        cards_h = min(1.6, 2 * cpad + max_lines * lh + 0.1)

    band = (h_in - 0.5) - top
    total = hl_h + (gap_claim if hl and km else 0.0) + km_h + (gap_cards if points else 0.0) + cards_h
    y = top + max(0.0, (band - total) / 2)

    # Claim (headline) — navy bold, pleine largeur (format VSCode3).
    if hl:
        D.add_text(slide, area_l, y, area_w, hl_h,
                   [(hl, {"size": D.TYPE["h2"], "bold": True, "color": D.INK})])
        y += hl_h + gap_claim
    # Sous-claim (key_message) — italique gris : le « so-what ».
    if km:
        D.add_text(slide, area_l, y, area_w, km_h,
                   [(km, {"size": D.TYPE["body"], "italic": True, "color": D.MUTED})])
        y += km_h
    y += gap_cards

    # Points clés en CARTES COULEUR — signature VSCode3. Carte blanche + liseré
    # couleur, texte centré, tronqué à ce qui tient (jamais de débordement).
    if points:
        for i, pt in enumerate(points[:n]):
            cx = area_l + i * (col_w + gap)
            color = _EXEC_CARD_COLORS[i % len(_EXEC_CARD_COLORS)]
            D.add_card(slide, cx, y, col_w, cards_h, color)
            D.add_text(
                slide, cx + cpad, y + cpad, col_w - 2 * cpad, cards_h - 2 * cpad,
                [(D.tronquer_a_lignes(pt, col_w - 2 * cpad, D.TYPE["small"], max_lines),
                  {"size": D.TYPE["small"], "color": D.INK})],
                anchor=MSO_ANCHOR.MIDDLE,
            )


def _label_axe_vertical(slide, cx: float, cy: float, longueur: float,
                        epaisseur: float, texte: str,
                        size: float | None = None) -> None:
    """Label d'axe roté 270° (lecture bas→haut), centré sur (cx, cy). `longueur`
    = dimension le long du texte (≈ hauteur de la ligne couverte), `epaisseur` =
    largeur de la gouttière. Sert aux libellés de ligne INTERNE/EXTERNE de la
    matrice SWOT et à l'axe Valeur de la matrice de priorisation."""
    box = D.add_text(
        slide, cx - longueur / 2, cy - epaisseur / 2, longueur, epaisseur,
        [(texte, {"size": size or D.TYPE["small"], "bold": True, "color": D.MUTED})],
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )
    box.rotation = 270


def _slide_swot(prs: Presentation, swot) -> None:
    """Matrice SWOT 2×2 — cf. skill `swot-matrix`. Ce n'est PAS quatre cartes
    posées côte à côte : c'est une matrice dont les deux axes sont explicites —
    lignes INTERNE (Forces/Faiblesses) / EXTERNE (Opportunités/Menaces) dans la
    gouttière gauche (labels rotés), colonnes FAVORABLE (Forces/Opportunités) /
    DÉFAVORABLE (Faiblesses/Menaces) au-dessus. Chaque quadrant est une CELLULE
    TEINTÉE (fond = melanger_blanc de sa couleur) : le fond rempli rend le vide
    sous les puces intentionnel, au lieu de la carte blanche sur-étirée que
    pptx-verify signalait. Grille figée par les axes : Forces (h-g), Faiblesses
    (h-d), Opportunités (b-g), Menaces (b-d)."""
    slide, w_in, h_in, top = _new_slide(prs, "Matrice SWOT")
    gutter = 0.30   # gouttière gauche : labels de ligne INTERNE/EXTERNE (rotés)
    axis_h = 0.30   # bandeau haut : labels de colonne FAVORABLE/DÉFAVORABLE
    gap = 0.22
    pad = 0.16
    area_l = MARGIN + gutter
    area_w = w_in - MARGIN - area_l
    area_t = top + axis_h
    area_h = h_in - area_t - 0.45
    col_w = (area_w - gap) / 2
    row_h = (area_h - gap) / 2
    cells = [(0, 0), (1, 0), (0, 1), (1, 1)]
    title_h = 0.40

    # Axe horizontal (effet sur l'objectif) : FAVORABLE (vert) / DÉFAVORABLE (rouge).
    for ci, (lbl, col) in enumerate((("FAVORABLE", D.OK), ("DÉFAVORABLE", D.WARN))):
        D.add_text(
            slide, area_l + ci * (col_w + gap), top, col_w, axis_h,
            [(lbl, {"size": D.TYPE["tiny"], "bold": True, "color": col})],
            anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
        )
    # Axe vertical (origine) : INTERNE / EXTERNE (neutre — l'origine n'est pas +/-).
    # `longueur` bornée (< 2×cx) pour que le cadre NON roté du label — celui que
    # verifier_geometrie contrôle — reste dans la slide ; le label roté visuel, lui,
    # tient dans la gouttière quoi qu'il arrive.
    for ri, lbl in enumerate(("INTERNE", "EXTERNE")):
        cy = area_t + ri * (row_h + gap) + row_h / 2
        _label_axe_vertical(slide, MARGIN + gutter / 2, cy, min(row_h, 1.3), gutter, lbl)

    for (key, label, color), (col, row) in zip(_SWOT_QUADRANTS, cells):
        cl = area_l + col * (col_w + gap)
        ct = area_t + row * (row_h + gap)
        # Cellule teintée + liseré coloré (style de carte du deck). Le fond rempli
        # supprime l'effet « carte blanche vide » sous des puces courtes.
        D.add_rect(slide, cl, ct, col_w, row_h,
                   fill=D.melanger_blanc(color, 0.90),
                   line=D.melanger_blanc(color, 0.55), line_w=1.0,
                   rounded=True, radius=0.05)
        D.add_rect(slide, cl, ct, 0.06, row_h, fill=color, rounded=True, radius=0.5)
        # En-tête : badge icône + titre coloré du quadrant.
        badge_d = 0.30
        hy = ct + pad
        D.add_badge(slide, cl + pad + 0.04, hy, badge_d, _SWOT_ICONS[key],
                    color, size=D.TYPE["small"], bold=False, radius=0.28)
        D.add_text(
            slide, cl + pad + 0.04 + badge_d + 0.12, hy,
            col_w - 2 * pad - badge_d - 0.16, title_h,
            [(label, {"size": D.TYPE["h3"], "bold": True, "color": color})],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # paginate=True : un quadrant trop long est TRONQUÉ à la cellule plutôt que
        # de déborder sur le voisin. max(0.0, …) : jamais négatif.
        _add_bulleted_text(
            slide, cl + pad + 0.04, ct + pad + title_h + 0.04, col_w - 2 * pad - 0.04,
            max(0.0, row_h - (pad + title_h + 0.04) - pad),
            getattr(swot, key) or "—",
            anchor=MSO_ANCHOR.TOP, size_max=D.TYPE["small"], size_min=D.TYPE["tiny"],
            paginate=True,
        )


def _slide_difficultes(prs: Presentation, difficulties) -> None:
    """Planche « Difficultés identifiées » (piste F restitution) — une carte par
    difficulté (rang + constat), chacune pouvant porter un verbatim en encadré
    citation (l'« insert citation » prévu de longue date, cf.
    docs/reflexions/restitution-mission.md §D.1). Cartes empilées et DIMENSIONNÉES
    à leur contenu (comme _slide_verbatims), on s'arrête avant de déborder du
    cadre (garantit verifier_geometrie)."""
    slide, w_in, h_in, top = _new_slide(prs, "Difficultés identifiées")
    pad, gap = 0.18, 0.16
    area_l = MARGIN + 0.3
    area_w = w_in - 2 * (MARGIN + 0.3)
    area_bottom = h_in - 0.5
    accent = "#b8860b"  # ambre : signal « point d'attention »
    teal = "#138086"    # citation, cohérent avec la planche verbatims
    size = D.TYPE["body"]
    q_size = D.TYPE["small"]
    line_h = _per_line_height_in(size)
    q_line_h = _per_line_height_in(q_size)
    # Rang en chip numéroté (ambre) à gauche de la carte, au lieu du préfixe « N. »
    # dans le libellé — le texte du constat démarre après le chip (largeur réduite,
    # reflétée dans FIELD_SHAPE["difficulty_label"]).
    rang_w, rang_h = 0.46, 0.30
    lab_x = area_l + pad + rang_w + 0.16
    lab_w = area_w - 2 * pad - rang_w - 0.16
    for i, d in enumerate(difficulties, 1):
        label = (getattr(d, "label", "") or "").strip()
        if not label:
            continue
        lab_lines = min(3, max(1, D.estimer_lignes(label, lab_w, size)))
        v = getattr(d, "verbatim", None)
        quote = ""
        if v is not None and (getattr(v, "quote", "") or "").strip():
            who = (getattr(getattr(v, "interview", None), "interviewee_name", "") or "Anonyme").strip() or "Anonyme"
            quote = f"«  {v.quote.strip()}  » — {who}"
        q_lines = min(2, max(1, D.estimer_lignes(quote, lab_w, q_size))) if quote else 0
        head_block = max(rang_h, lab_lines * line_h)
        card_h = pad + head_block + (0.06 + q_lines * q_line_h if quote else 0.0) + pad
        if top + card_h > area_bottom and i > 1:  # au moins la 1re carte, sinon stop
            break
        if top + card_h > area_bottom:
            card_h = max(0.0, area_bottom - top)  # 1re carte trop haute : bornée au cadre
        D.add_card(slide, area_l, top, area_w, card_h, accent)
        D.add_chip(slide, area_l + pad, top + pad, rang_w, rang_h, str(i), accent,
                   size=D.TYPE["small"])
        D.add_text(
            slide, lab_x, top + pad, lab_w, head_block,
            [(D.tronquer_a_lignes(label, lab_w, size, lab_lines),
              {"size": size, "bold": True, "color": D.INK})],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        if quote:
            D.add_text(
                slide, lab_x, top + pad + head_block + 0.06,
                lab_w, q_lines * q_line_h,
                [(D.tronquer_a_lignes(quote, lab_w, q_size, q_lines),
                  {"size": q_size, "italic": True, "color": teal})],
            )
        top += card_h + gap


def _slide_verbatims(prs: Presentation, verbatims) -> None:
    """Planche « Paroles d'acteurs » (Palier 2) — une carte-citation par
    verbatim retenu (attribution en libellé discret, citation en corps italique),
    empilées depuis le haut, chaque carte DIMENSIONNÉE À SON CONTENU (2 lignes de
    citation au plus) plutôt qu'étirée à `area_h / n` — sinon une citation d'une
    ligne laisse un grand vide dans sa carte (constat pptx-verify). Le surplus se
    reporte en blanc en bas de slide. On s'arrête avant de déborder du cadre
    (garantit le garde-fou géométrie) — l'onglet aperçu invite à 2-4 citations."""
    slide, w_in, h_in, top = _new_slide(prs, "Paroles d'acteurs")
    pad, gap = 0.18, 0.18
    label_h = 0.3
    area_l = MARGIN + 0.3
    area_w = w_in - 2 * (MARGIN + 0.3)
    area_bottom = h_in - 0.5
    size = D.TYPE["body"]
    line_h = _per_line_height_in(size)
    y = top
    for v in verbatims:
        quote = f"«  {(v.quote or '').strip()}  »"
        q_lines = min(3, max(1, D.estimer_lignes(quote, area_w - 2 * pad, size)))
        card_h = pad + label_h + q_lines * line_h + pad
        if y + card_h > area_bottom:  # ne jamais déborder du cadre
            break
        D.add_card(slide, area_l, y, area_w, card_h, "#138086")
        who = (getattr(v.interview, "interviewee_name", "") or "Anonyme").strip() or "Anonyme"
        _add_measured_field(
            slide, area_l + pad, y + pad, area_w - 2 * pad,
            label=who, text=quote, max_h=label_h + q_lines * line_h,
            size_max=size, size_min=D.TYPE["tiny"], italic=True,
        )
        y += card_h + gap


_AXES_ROW_H_MAX = 1.1
_AXES_ROW_GAP = 0.15
# En dessous de cette hauteur de ligne, le chiffre "#N" (D.TYPE["kpi"]=44pt)
# chevauche visuellement le titre de l'axe à côté — verifier_geometrie() ne
# peut pas le détecter (il ne vérifie que les bords des formes, pas le rendu
# du texte à l'intérieur) : mieux vaut paginer sur une slide suivante que
# de laisser les cartes devenir illisibles avec beaucoup d'axes.
_AXES_ROW_H_MIN = 0.75


def _axes_row_h(n: int, band_h: float) -> float:
    return min(_AXES_ROW_H_MAX, (band_h - _AXES_ROW_GAP * (n - 1)) / max(1, n))


def _slide_axes_overview(prs: Presentation, axes: list, palette: list[str]) -> None:
    title = "Les recommandations sont construites autour de ces axes"
    # Sert UNIQUEMENT à décider combien d'axes tiennent par page (1.4in ~
    # hauteur de contenu typique après un titre sur une ligne, cf. _new_slide) ;
    # chaque page recalcule ensuite sa hauteur réellement disponible à partir
    # de SON PROPRE titre (avec suffixe) une fois la slide créée, donc ce
    # découpage préalable ne peut jamais faire déborder une carte — au pire
    # (titre passé à 2 lignes à cause du suffixe) la page rend des rangées
    # un peu plus basses que prévu, jamais hors-cadre.
    w_in, h_in = _dims(prs)
    band_h_estimate = h_in - 1.4 - 0.5
    row_h_estimate = max(_AXES_ROW_H_MIN, _axes_row_h(len(axes), band_h_estimate))
    pages = D.paginer_items(
        list(enumerate(axes)), lambda _item: row_h_estimate + _AXES_ROW_GAP,
        capacite_in=band_h_estimate + _AXES_ROW_GAP,
    )
    for k, page in enumerate(pages):
        suffix = f" ({k + 1}/{len(pages)})" if len(pages) > 1 else ""
        slide, w_in, h_in, top = _new_slide(prs, title + suffix)
        band_h = h_in - top - 0.5
        row_h = _axes_row_h(len(page), band_h)
        total_h = len(page) * row_h + _AXES_ROW_GAP * (len(page) - 1)
        # Centré verticalement dans la bande plutôt que plaqué en haut : avec
        # peu d'axes (1-3) sur la page, row_h plafonne à 1.1in et laisse
        # sinon un grand vide sous les cartes.
        y = top + max(0.0, (band_h - total_h) / 2)
        for i, axis in page:
            accent = palette[i % len(palette)]
            D.add_card(slide, MARGIN, y, w_in - 2 * MARGIN, row_h, accent)
            # Pastille teardrop numérotée (signature OCTO du sommaire) plutôt que
            # « #N » nu, et les INTITULÉS des recos en 2e ligne plutôt qu'un simple
            # compte « 2 recommandations » — la rangée était à moitié vide et le
            # texte creux (constat utilisateur 2026-07-22, slide 16).
            td = min(0.62, row_h - 0.16)
            D.add_teardrop(slide, MARGIN + 0.28, y + (row_h - td) / 2, td,
                           str(i + 1), accent, size=D.TYPE["h3"])
            text_x = MARGIN + 0.28 + td + 0.3
            text_w = w_in - MARGIN - text_x - 0.3
            recos_txt = "   ·   ".join(
                f"{i + 1}.{j + 1}  {r.title}" for j, r in enumerate(axis.recommendations)
            )
            D.add_text(
                slide, text_x, y, text_w, row_h,
                [
                    (axis.title, {"size": D.TYPE["h3"], "bold": True, "color": D.INK,
                                  "space_after": 4}),
                    (D.tronquer_a_lignes(recos_txt, text_w, D.TYPE["small"], 2),
                     {"size": D.TYPE["small"], "color": D.MUTED}),
                ],
                anchor=MSO_ANCHOR.MIDDLE,
            )
            y += row_h + _AXES_ROW_GAP


# Quadrants de la matrice de priorisation (skill priority-matrix) : le SENS de
# chaque quadrant est écrit dessus — c'est ce qui transforme un nuage de points en
# outil de décision. (label, couleur) par position (colonne, ligne) de la grille.
_PRIO_QUADRANTS = {
    # Libellés courts exprès : à `small` bold ils doivent tenir sur UNE ligne dans
    # une demi-grille (« CHANTIERS STRUCTURANTS » wrappait derrière les bulles).
    (0, 0): ("QUICK WINS", "#1e6b34"),          # valeur haute, effort faible
    (1, 0): ("CHANTIERS DE FOND", "#2c5cc5"),   # valeur haute, effort fort
    (0, 1): ("OPPORTUNISTES", "#6b7280"),       # valeur basse, effort faible
    (1, 1): ("À DIFFÉRER", "#b8860b"),          # valeur basse, effort fort
}


def _slide_matrice_effort_valeur(prs: Presentation, axes: list,
                                 palette: list[str]) -> None:
    """Matrice de priorisation valeur/effort DESSINÉE (skill priority-matrix) — le
    graphique scatter natif PowerPoint rendait « très mauvais » (constat utilisateur
    2026-07-22 : marqueurs Excel minuscules gris, légende cryptique ◆■▲, aucun
    quadrant). Ici : 4 quadrants teintés dont le SENS est écrit dessus, une bulle
    par reco (couleur = axe, même palette que la vue d'ensemble ; numéro dedans),
    bulles co-localisées déployées en éventail, légende par axe à droite."""
    slide, w_in, h_in, top = _new_slide(prs, "Matrice de priorisation — valeur / effort")
    # Zone de tracé (gouttière gauche pour le label d'axe Y roté, bande basse pour X).
    pl = MARGIN + 0.45
    pt = top + 0.10
    pb = h_in - 0.72
    ph = pb - pt
    pw = 3.95  # plot un peu plus étroit : la légende porte les intitulés COMPLETS
    lx = pl + pw + 0.3   # légende à droite
    lw = w_in - MARGIN - lx
    qw, qh = pw / 2, ph / 2

    # Quadrants teintés + libellé de sens dans chaque coin — en `small`, pas
    # `tiny` : lisibilité relevée par l'utilisateur (2026-07-22, slide 17).
    for (col, row), (lbl, color) in _PRIO_QUADRANTS.items():
        qx, qy = pl + col * qw, pt + row * qh
        D.add_rect(slide, qx, qy, qw, qh, fill=D.melanger_blanc(color, 0.93),
                   line=D.melanger_blanc(color, 0.70), line_w=0.75)
        D.add_text(
            slide, qx + 0.10, qy + 0.06, qw - 0.20, 0.26,
            [(lbl, {"size": D.TYPE["small"], "bold": True,
                    "color": D.melanger_blanc(color, 0.15)})],
            align=PP_ALIGN.LEFT if col == 0 else PP_ALIGN.RIGHT,
        )
    # Labels d'axes : X sous la zone, Y roté dans la gouttière gauche.
    D.add_text(slide, pl, pb + 0.08, pw, 0.3,
               [("Complexité (effort) →", {"size": D.TYPE["small"], "bold": True,
                                           "color": D.MUTED})],
               align=PP_ALIGN.CENTER)
    _label_axe_vertical(slide, MARGIN + 0.18, pt + ph / 2, min(ph, 1.5), 0.3,
                        "Valeur (impact) →")

    # Bulles : les scores sont des entiers 1-5, les collisions sont la norme — y
    # compris ENTRE scores voisins (constat pptx-verify : deux bulles de scores
    # adjacents se chevauchaient, la 2e masquait le numéro de la 1re). Résolution
    # par LIGNE (même valeur → même y, les lignes sont espacées de ph/5 > d) :
    # balayage gauche→droite qui impose un écart minimal à partir des positions
    # cibles, puis recalage si la ligne déborde à droite.
    d = 0.46  # bulle agrandie + numéro en `small` (lisibilité, 2026-07-22)
    gap = 0.06
    lignes_bulles: dict[int, list] = {}
    for i, axis in enumerate(axes):
        for j, r in enumerate(axis.recommendations):
            c = max(1, min(5, r.complexite or 3))
            v = max(1, min(5, r.valeur or 3))
            lignes_bulles.setdefault(v, []).append((c, f"{i + 1}.{j + 1}", i))
    # Les bulles vivent SOUS la bande des libellés de quadrant (0.38 réservé en
    # haut) : à valeur=5 elles venaient sinon recouvrir le libellé (constat rendu).
    ph_bulles = ph - 0.38
    for v, membres in lignes_bulles.items():
        membres.sort(key=lambda m: (m[0], m[1]))
        by = pb - (v - 0.5) / 5 * ph_bulles - d / 2
        xs: list[float] = []
        for c, _num, _ai in membres:
            cible = pl + (c - 0.5) / 5 * pw - d / 2
            xs.append(cible if not xs else max(cible, xs[-1] + d + gap))
        depassement = xs[-1] - (pl + pw - d - 0.02)
        if depassement > 0:  # recale toute la ligne dans la zone (sans re-chevaucher)
            xs = [max(pl + 0.02, x - depassement) for x in xs]
        for x, (c, num, ai) in zip(xs, membres):
            D.add_badge(slide, x, by, d, num, palette[ai % len(palette)],
                        size=D.TYPE["small"], bold=True, radius=0.5)

    # Légende ENCADRÉE (carte) portant les intitulés COMPLETS — demande 2026-07-22 :
    # « réduire la taille du texte afin qu'il apparaisse complètement et à
    # encadrer ». Taille `tiny` partout, repli sur 2 lignes max par item (mesuré,
    # jamais tronqué à 1 ligne comme avant), hauteur de chaque item MESURÉE.
    # Une ligne par RECO uniquement — pas d'intitulés d'axes (ils vivent en
    # toutes lettres sur la vue d'ensemble ; ici la pastille couleur suffit à
    # porter l'axe, comme les bulles) : c'est ce qui permet aux 8 titres de reco
    # COMPLETS de tenir (4 titres d'axes en plus faisaient sauter l'axe 4).
    # -0.60 (pas -0.50) : le chrome n° de page du master OCTO démarre à y≈5.09 /
    # x≈9.25 — à -0.50 le coin bas-droit de la carte (blanc + bordure) peignait
    # par-dessus (revue adversariale, mesuré sur le master ; même garde que la
    # fiche reco).
    leg_bottom = h_in - 0.60
    D.add_card(slide, lx, pt, lw, leg_bottom - pt)
    lpad = 0.12
    tx = lx + lpad
    tw = lw - 2 * lpad
    # Shrink-to-fit — JAMAIS droper une reco (à taille fixe, l'estimation
    # pessimiste s'accumulait sur 8 items et « 4.2 » sautait alors qu'il
    # restait de la place réelle) : on cherche la plus grande taille <= tiny
    # qui fait tenir TOUTES les recos à l'estimation PESSIMISTE (celle du
    # vérificateur — à l'estimation nominale, un item limite wrappait hors
    # boîte). Cascade complète (revue adversariale : l'ancien `while t_leg > 7.5`
    # sortait SANS avoir évalué 7.5, et le garde-fou du rendu dropait alors des
    # recos sur titres extrêmes) : tailles 9→7.5 à 2 lignes/item, puis dernier
    # cran 7.5 pt à 1 ligne/item (titre tronqué à l'ellipse — un titre coupé
    # vaut toujours mieux qu'une reco absente).
    dispo = (leg_bottom - lpad) - (pt + lpad)
    t_leg = D.TYPE["tiny"]
    lignes_leg = 2
    while True:
        lh_leg = _per_line_height_in(t_leg)
        besoin = 0.0
        for i, axis in enumerate(axes):
            for j, r in enumerate(axis.recommendations):
                item = D.tronquer_a_lignes(f"{i + 1}.{j + 1}  {r.title}", tw - 0.24, t_leg, lignes_leg)
                besoin += D.estimer_lignes(item, tw - 0.24, t_leg, cpi_ref=10.7) * lh_leg + 0.03
            besoin += 0.02
        if besoin <= dispo:
            break
        if t_leg > 7.5:
            t_leg -= 0.5
        elif lignes_leg == 2:
            lignes_leg = 1
        elif t_leg > 6.5:
            t_leg -= 0.5  # dernier étage : 1 ligne, 7.5→6.5 (loge ~16 recos)
        else:
            break  # plafond structurel ~18 recos — au-delà le garde-fou du rendu coupe
    y = pt + lpad
    plein = False  # garde-fou ultime : coupe TOUT le reste (pas d'items suivants
    # rendus après un trou — des numéros manquants au milieu seraient trompeurs)
    for i, axis in enumerate(axes):
        color = palette[i % len(palette)]
        for j, r in enumerate(axis.recommendations):
            item = D.tronquer_a_lignes(f"{i + 1}.{j + 1}  {r.title}", tw - 0.24, t_leg, lignes_leg)
            h_item = D.estimer_lignes(item, tw - 0.24, t_leg, cpi_ref=10.7) * lh_leg
            if y + h_item > leg_bottom - lpad:
                plein = True
                break
            D.add_rect(slide, tx, y + 0.03, 0.12, 0.12, fill=color, rounded=True, radius=0.5)
            D.add_text(slide, tx + 0.24, y, tw - 0.24, h_item,
                       [(item, {"size": t_leg, "color": D.INK})])
            y += h_item + 0.03
        if plein:
            break
        y += 0.02


def _slide_recommendation(prs: Presentation, axis: object, index: str, reco: object,
                          accent: str | None = None) -> None:
    """Fiche recommandation en ENCARTS ARRONDIS format OCTO (demande 2026-07-22 —
    les sections flottaient sur fond blanc) : colonne gauche (objectif / acteurs /
    jauges / résultats) dans une carte arrondie au liseré couleur d'AXE (identité,
    même palette que la vue d'ensemble et la matrice de priorisation) ; colonne
    droite en deux encarts empilés — PROPOSITION DE VALEUR en encart gris arrondi
    (le « so-what » de la fiche, même composant que l'exec/synthèse) puis PLAN
    D'ACTIONS en carte arrondie. OBJECTIF/ACTEURS gardent la hauteur MESURÉE
    (_add_measured_field) pour s'empiler sans déborder ; l'encart proposition est
    à hauteur FIXE (rythme identique de fiche en fiche, l'espace gris restant est
    intentionnel — même principe que les cellules teintées de la SWOT)."""
    # Titre sur UNE ligne (tronqué à l'ellipse) : le titre complet vit sur la vue
    # d'ensemble des axes — la ligne gagnée ici est ce qui permet à la carte
    # droite de loger plan + résultats sans slide de suite systématique.
    slide, w_in, h_in, top = _new_slide(prs, f"{index} — {reco.title}", max_title_lines=1)
    accent = accent or (D.theme_colors(prs).get("accent1") or D.PALETTE[0])
    pad = 0.2
    lis = 0.05  # dégagement du liseré de carte
    # Carte gauche resserrée (3.15) au profit de la droite : les jauges 0.65×2 y
    # tiennent, et le plan (souvent UNE longue puce) a besoin de largeur.
    card_l_w = 3.15
    right_x = MARGIN + card_l_w + 0.3
    right_w = w_in - right_x - MARGIN
    # Bandeau RÉSULTATS pleine largeur en bas (encart gris, motif « à retenir »
    # des synthèses) : une longue phrase tient en 2 lignes sur ~8.7in là où elle
    # en demandait 5 dans une colonne — c'est ce qui rend la fiche tenable sans
    # slide de suite systématique (mesuré : colonne droite 1.9in dispo pour 2.4in
    # de besoin, aucune allocation ne pouvait gagner). -0.60 : badge n° de page.
    # Prédicat aligné sur le CONTENU rendu (_bullet_lines filtre les marqueurs
    # de puce vides) : « - » seul réservait 0.72in pour un bandeau « — » vide.
    a_resultats = bool(_bullet_lines(reco.resultats_attendus or ""))
    strip_h = 0.72 if a_resultats else 0.0
    band_h = h_in - top - 0.60 - strip_h - (0.12 if a_resultats else 0.0)
    bottom = top + band_h

    # ---- Colonne gauche : carte arrondie, liseré couleur d'axe ----
    D.add_card(slide, MARGIN, top, card_l_w, band_h, accent)
    lx = MARGIN + pad + lis
    lw = card_l_w - 2 * pad - lis
    y = top + pad
    y += _add_measured_field(slide, lx, y, lw, "OBJECTIF", reco.objectif, max_h=1.1)
    y += 0.10
    y += _add_measured_field(slide, lx, y, lw, "ACTEURS", reco.acteurs, max_h=0.5)
    y += 0.10
    D.add_text(slide, lx, y, lw, 0.26, [("CRITÈRES DE PRIORISATION", {"size": D.TYPE["small"], "bold": True, "color": D.MUTED})])
    # Chips « Valeur N/5 » / « Complexité N/5 » (couleurs sémantiques OK/WARN) au
    # lieu des jauges donut : la carte gauche a perdu ~0.8in au profit du bandeau
    # résultats — les donuts s'y écrasaient (labels hors carte, constat rendu
    # 2026-07-22). Une ligne de chips porte la même information en 0.3in.
    chip_h = 0.32
    chip_w = min(1.30, (lw - 0.15) / 2)
    # Clamp DANS la carte (revue adversariale) : objectif+acteurs au max de leurs
    # hauteurs poussaient les chips sous le bord bas — le clamp les garde dans la
    # carte (au pire chevauche le label CRITÈRES, jamais le bandeau résultats).
    chips_y = min(y + 0.32, top + band_h - pad - chip_h)
    D.add_chip(slide, lx, chips_y, chip_w, chip_h,
               f"Valeur {reco.valeur}/5", D.OK, size=D.TYPE["tiny"])
    D.add_chip(slide, lx + chip_w + 0.15, chips_y, chip_w, chip_h,
               f"Complexité {reco.complexite}/5", D.WARN, size=D.TYPE["tiny"])

    # ---- Colonne droite : encart « proposition » + carte « plan + résultats » ----
    # prop_h 1.10 (était 1.35) : la carte droite porte TROIS blocs — au-delà, la
    # zone résultats devenait fictive (~0.1in) et son texte peignait PAR-DESSUS le
    # cadre (le « texte sort du cadre » relevé par l'utilisateur, objectivé par
    # verifier_debordements_texte).
    prop_h = 1.00
    D.add_rect(slide, right_x, top, right_w, prop_h, fill=D.ENCART_BG,
               rounded=True, radius=0.12)
    D.add_rect(slide, right_x, top, 0.06, prop_h, fill=accent, rounded=True, radius=0.5)
    _add_measured_field(
        slide, right_x + pad + lis, top + 0.10, right_w - 2 * pad - lis,
        "PROPOSITION DE VALEUR", reco.proposition_valeur, max_h=prop_h - 0.20,
        bold=True, italic=True,
    )
    plan_top = top + prop_h + 0.15
    plan_h = bottom - plan_top
    D.add_card(slide, right_x, plan_top, right_w, plan_h, accent)
    rcx = right_x + pad + lis
    rcw = right_w - 2 * pad - lis
    r_top = plan_top + pad
    r_bottom = plan_top + plan_h - pad
    D.add_text(slide, rcx, r_top, rcw, 0.26,
               [("PLAN D'ACTIONS", {"size": D.TYPE["small"], "bold": True, "color": D.MUTED})])
    # Le plan a TOUTE la carte (les résultats vivent dans le bandeau bas) —
    # shrink-first, suite en dernier recours seulement.
    plan_overflow = _add_bulleted_text(
        slide, rcx, r_top + 0.26, rcw, r_bottom - (r_top + 0.26),
        reco.plan_actions, paginate=True,
    )

    # ---- Bandeau bas pleine largeur : RÉSULTATS ATTENDUS (encart gris) ----
    # Une seule zone MIDDLE (libellé + texte dans la même boîte) : sur ~8.7in de
    # large, 2 lignes en `small` logent ~200 caractères — pas de pagination,
    # troncature à l'ellipse en tout dernier recours (FIELD_SHAPE l'annonce).
    if a_resultats:
        sy = top + band_h + 0.12
        sw = w_in - 2 * MARGIN
        D.add_rect(slide, MARGIN, sy, sw, strip_h, fill=D.ENCART_BG,
                   rounded=True, radius=0.12)
        D.add_rect(slide, MARGIN, sy, 0.06, strip_h, fill=accent, rounded=True, radius=0.5)
        scx = MARGIN + pad + lis
        scw = sw - 2 * pad - lis
        res_txt = " — ".join(_bullet_lines(reco.resultats_attendus)) or "—"
        D.add_text(
            slide, scx, sy, scw, strip_h,
            [("RÉSULTATS ATTENDUS", {"size": D.TYPE["tiny"], "bold": True,
                                     "color": D.MUTED, "space_after": 2}),
             # cpi PESSIMISTE (10.7) : à l'estimation nominale un texte limite
             # (~180 car.) repassait à 3 lignes au vrai rendu et sortait du
             # bandeau — hors du champ des vérificateurs (ancre MIDDLE).
             (D.tronquer_a_lignes(res_txt, scw, D.TYPE["small"], 2, cpi_ref=10.7),
              {"size": D.TYPE["small"], "color": D.INK})],
            anchor=MSO_ANCHOR.MIDDLE,
        )

    base_title = f"{index} — {reco.title}"
    if plan_overflow:
        _emit_bullet_overflow(prs, base_title, "Plan d'actions", plan_overflow)


# --------------------------------------------------------------------------- #
# Façade
# --------------------------------------------------------------------------- #
def build_presentation(
    mission: Mission,
    template_path: Path | None = None,
    include_sommaire: bool = True,
    include_executive_summary: bool = True,
    include_synthese: bool = True,
    include_difficultes: bool = True,
    include_swot: bool = True,
    include_verbatims: bool = True,
    include_axes_overview: bool = True,
    include_matrix: bool = True,
    include_axis_ids: set[int] | None = None,
) -> Presentation:
    """`include_axis_ids=None` inclut les fiches de recommandation de tous les
    axes (comportement par défaut/rétrocompatible) ; un set (même vide)
    restreint aux axes dont l'id y figure — la vue d'ensemble des axes et la
    matrice effort/valeur restent, elles, toujours complètes (ce sont des
    slides de synthèse, pas de détail par axe)."""
    if template_path and Path(template_path).exists():
        prs = Presentation(str(template_path))
        _clear_slides(prs)
    elif OCTO_TEMPLATE_PATH.exists():
        # Défaut : le template de marque OCTO (chrome + layouts + thème + Outfit).
        prs = Presentation(str(OCTO_TEMPLATE_PATH))
        _clear_slides(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(_W_IN)
        prs.slide_height = Inches(_H_IN)
        prs._i2d_synthetic = True

    # Police effective du deck. On PRÉFÈRE la police du THÈME (fontScheme, Arial sur
    # OCTO) à celle des placeholders (Outfit) : Outfit n'étant pas installée, elle est
    # rendue en substitution — c'est la cause du « la police ne matche pas la référence »
    # (bmad-iap-cadrage-synthese utilise, lui, la police du thème). Repli sur la police
    # des placeholders puis héritage. None sur le deck synthétique (inchangé).
    if getattr(prs, "_i2d_synthetic", False):
        D.set_police(None)
    else:
        D.set_police(D.police_theme(prs) or D.police_marque(prs))

    # Ancre la palette catégorielle des axes sur la couleur de marque du
    # template injecté, sans jamais remplacer toute la palette par elle
    # (une palette catégorielle reste plus lisible pour distinguer N axes).
    brand_accent = D.theme_colors(prs).get("accent1")
    palette = ([brand_accent] + D.PALETTE) if brand_accent else D.PALETTE

    _slide_cover(prs, mission)

    gs = mission.global_synthesis
    swot = mission.swot
    executive_summary = mission.executive_summary
    difficulties = [d for d in mission.difficulties if (d.label or "").strip()]
    verbatims = mission.selected_verbatims
    axes = list(mission.recommendation_axes)
    selected_axes = [a for a in axes if include_axis_ids is None or a.id in include_axis_ids]

    # Sections présentes, groupées par chapitre (P2 — structure narrative). Un
    # intercalaire ouvre chaque chapitre qui a du contenu ; le sommaire quali les liste.
    ch_sections: list[list[str]] = [[] for _ in _CHAPITRES]
    if include_executive_summary and executive_summary and executive_summary.has_content:
        ch_sections[_CH_RETENIR].append("Executive Summary")
    if include_synthese and gs and gs.has_content:
        ch_sections[_CH_DIAGNOSTIC].append("Synthèse globale")
    if include_difficultes and difficulties:
        ch_sections[_CH_DIAGNOSTIC].append("Difficultés")
    if include_swot and swot and swot.has_content:
        ch_sections[_CH_DIAGNOSTIC].append("Matrice SWOT")
    if include_verbatims and verbatims:
        ch_sections[_CH_PAROLE].append("Paroles d'acteurs")
    # « Recommandations » dès qu'il y a des axes à détailler (les fiches reco
    # s'émettent indépendamment des toggles overview/matrice) OU la vue d'ensemble.
    if (axes and include_axes_overview) or selected_axes:
        ch_sections[_CH_TRAJECTOIRE].append("Recommandations")
    if axes and include_matrix:
        ch_sections[_CH_TRAJECTOIRE].append("Matrice de priorisation")

    if include_sommaire and any(ch_sections):
        _slide_sommaire(prs, ch_sections)

    numero = 0

    def _chapitre(ci: int) -> None:
        nonlocal numero
        numero += 1
        label, color, scene, sous_titre = _CHAPITRES[ci]
        _slide_chapitre(prs, numero, label, color, scene, sous_titre=sous_titre)

    # Chapitre 1 — Ce qu'il faut retenir
    if ch_sections[_CH_RETENIR]:
        _chapitre(_CH_RETENIR)
        _slide_executive_summary(prs, executive_summary)

    # Chapitre 2 — Le diagnostic
    if ch_sections[_CH_DIAGNOSTIC]:
        _chapitre(_CH_DIAGNOSTIC)
        if include_synthese and gs and gs.has_content:
            categories = [
                ("Contexte", gs.contexte),
                ("Culture & ADN", gs.culture_adn),
                ("Forces & succès", gs.forces_succes),
                ("Points d'amélioration", gs.points_amelioration),
                ("Aspirations (baguette magique)", gs.aspirations),
            ]
            for label, content in categories:
                if (content or "").strip():
                    _slide_synthese_categorie(prs, label, content)
        if include_difficultes and difficulties:
            _slide_difficultes(prs, difficulties)
        if include_swot and swot and swot.has_content:
            _slide_swot(prs, swot)

    # Chapitre 3 — La parole des équipes
    if ch_sections[_CH_PAROLE]:
        _chapitre(_CH_PAROLE)
        _slide_verbatims(prs, verbatims)

    # Chapitre 4 — La trajectoire proposée
    if ch_sections[_CH_TRAJECTOIRE]:
        _chapitre(_CH_TRAJECTOIRE)
        if axes and include_axes_overview:
            _slide_axes_overview(prs, axes, palette)
        if axes and include_matrix:
            _slide_matrice_effort_valeur(prs, axes, palette)
        for i, axis in enumerate(axes):
            if axis not in selected_axes:
                continue
            for j, reco in enumerate(axis.recommendations):
                # accent = couleur d'axe (identité) — même palette que la vue
                # d'ensemble et les bulles de la matrice de priorisation.
                _slide_recommendation(prs, axis, f"{i + 1}.{j + 1}", reco,
                                      accent=palette[i % len(palette)])

    # Garde-fou géométrique (US7.1) : un texte trop long ou un template client
    # aux dimensions inattendues peut faire déborder une forme de la slide —
    # mieux vaut échouer bruyamment ici qu'exporter un .pptx visuellement cassé.
    problemes = D.verifier_geometrie(prs)
    if problemes:
        raise RuntimeError(
            "Export PPT : formes hors cadre détectées —\n" + "\n".join(problemes)
        )

    return prs
