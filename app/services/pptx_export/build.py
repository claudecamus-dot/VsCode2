"""Façade de l'export PPT : `build_presentation` assemble le deck complet
(couverture, sommaire, intercalaires, slides de contenu) puis fait respecter
le garde-fou géométrique. Extrait de pptx_export.py (découpage du gros
module, finding audit 2026-07-24) — code déplacé tel quel."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from ...models import Mission
from .. import pptx_deck as D
from .base import _H_IN, _W_IN, OCTO_TEMPLATE_PATH, _clear_slides
from .slides_cadre import (
    _CH_DIAGNOSTIC,
    _CH_PAROLE,
    _CH_RETENIR,
    _CH_TRAJECTOIRE,
    _CHAPITRES,
    _slide_chapitre,
    _slide_cover,
    _slide_sommaire,
)
from .slides_diagnostic import (
    _slide_difficultes,
    _slide_executive_summary,
    _slide_swot,
    _slide_synthese_categorie,
    _slide_verbatims,
)
from .slides_trajectoire import (
    _slide_axes_overview,
    _slide_matrice_effort_valeur,
    _slide_recommendation,
)


def build_presentation(
    mission: Mission,
    template_path: Path | None = None,
    include_sommaire: bool = True,
    include_executive_summary: bool = True,
    include_synthese: bool = True,
    include_difficultes: bool = True,
    include_swot: bool = True,
    include_verbatims: bool = True,
    include_axes_overview: bool = True,
    include_matrix: bool = True,
    include_axis_ids: set[int] | None = None,
) -> Presentation:
    """`include_axis_ids=None` inclut les fiches de recommandation de tous les
    axes (comportement par défaut/rétrocompatible) ; un set (même vide)
    restreint aux axes dont l'id y figure — la vue d'ensemble des axes et la
    matrice effort/valeur restent, elles, toujours complètes (ce sont des
    slides de synthèse, pas de détail par axe)."""
    if template_path and Path(template_path).exists():
        prs = Presentation(str(template_path))
        _clear_slides(prs)
    elif OCTO_TEMPLATE_PATH.exists():
        # Défaut : le template de marque OCTO (chrome + layouts + thème + Outfit).
        prs = Presentation(str(OCTO_TEMPLATE_PATH))
        _clear_slides(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(_W_IN)
        prs.slide_height = Inches(_H_IN)
        prs._i2d_synthetic = True

    # Police effective du deck. On PRÉFÈRE la police du THÈME (fontScheme, Arial sur
    # OCTO) à celle des placeholders (Outfit) : Outfit n'étant pas installée, elle est
    # rendue en substitution — c'est la cause du « la police ne matche pas la référence »
    # (bmad-iap-cadrage-synthese utilise, lui, la police du thème). Repli sur la police
    # des placeholders puis héritage. None sur le deck synthétique (inchangé).
    if getattr(prs, "_i2d_synthetic", False):
        D.set_police(None)
    else:
        D.set_police(D.police_theme(prs) or D.police_marque(prs))

    # Ancre la palette catégorielle des axes sur la couleur de marque du
    # template injecté, sans jamais remplacer toute la palette par elle
    # (une palette catégorielle reste plus lisible pour distinguer N axes).
    brand_accent = D.theme_colors(prs).get("accent1")
    palette = ([brand_accent] + D.PALETTE) if brand_accent else D.PALETTE

    _slide_cover(prs, mission)

    gs = mission.global_synthesis
    swot = mission.swot
    executive_summary = mission.executive_summary
    difficulties = [d for d in mission.difficulties if (d.label or "").strip()]
    verbatims = mission.selected_verbatims
    axes = list(mission.recommendation_axes)
    selected_axes = [a for a in axes if include_axis_ids is None or a.id in include_axis_ids]

    # Sections présentes, groupées par chapitre (P2 — structure narrative). Un
    # intercalaire ouvre chaque chapitre qui a du contenu ; le sommaire quali les liste.
    ch_sections: list[list[str]] = [[] for _ in _CHAPITRES]
    if include_executive_summary and executive_summary and executive_summary.has_content:
        ch_sections[_CH_RETENIR].append("Executive Summary")
    if include_synthese and gs and gs.has_content:
        ch_sections[_CH_DIAGNOSTIC].append("Synthèse globale")
    if include_difficultes and difficulties:
        ch_sections[_CH_DIAGNOSTIC].append("Difficultés")
    if include_swot and swot and swot.has_content:
        ch_sections[_CH_DIAGNOSTIC].append("Matrice SWOT")
    if include_verbatims and verbatims:
        ch_sections[_CH_PAROLE].append("Paroles d'acteurs")
    # « Recommandations » dès qu'il y a des axes à détailler (les fiches reco
    # s'émettent indépendamment des toggles overview/matrice) OU la vue d'ensemble.
    if (axes and include_axes_overview) or selected_axes:
        ch_sections[_CH_TRAJECTOIRE].append("Recommandations")
    if axes and include_matrix:
        ch_sections[_CH_TRAJECTOIRE].append("Matrice de priorisation")

    if include_sommaire and any(ch_sections):
        _slide_sommaire(prs, ch_sections)

    numero = 0

    def _chapitre(ci: int) -> None:
        nonlocal numero
        numero += 1
        label, color, scene, sous_titre = _CHAPITRES[ci]
        _slide_chapitre(prs, numero, label, color, scene, sous_titre=sous_titre)

    # Chapitre 1 — Ce qu'il faut retenir
    if ch_sections[_CH_RETENIR]:
        _chapitre(_CH_RETENIR)
        _slide_executive_summary(prs, executive_summary)

    # Chapitre 2 — Le diagnostic
    if ch_sections[_CH_DIAGNOSTIC]:
        _chapitre(_CH_DIAGNOSTIC)
        if include_synthese and gs and gs.has_content:
            categories = [
                ("Contexte", gs.contexte),
                ("Culture & ADN", gs.culture_adn),
                ("Forces & succès", gs.forces_succes),
                ("Points d'amélioration", gs.points_amelioration),
                ("Aspirations (baguette magique)", gs.aspirations),
            ]
            for label, content in categories:
                if (content or "").strip():
                    _slide_synthese_categorie(prs, label, content)
        if include_difficultes and difficulties:
            _slide_difficultes(prs, difficulties)
        if include_swot and swot and swot.has_content:
            _slide_swot(prs, swot)

    # Chapitre 3 — La parole des équipes
    if ch_sections[_CH_PAROLE]:
        _chapitre(_CH_PAROLE)
        _slide_verbatims(prs, verbatims)

    # Chapitre 4 — La trajectoire proposée
    if ch_sections[_CH_TRAJECTOIRE]:
        _chapitre(_CH_TRAJECTOIRE)
        if axes and include_axes_overview:
            _slide_axes_overview(prs, axes, palette)
        if axes and include_matrix:
            _slide_matrice_effort_valeur(prs, axes, palette)
        for i, axis in enumerate(axes):
            if axis not in selected_axes:
                continue
            for j, reco in enumerate(axis.recommendations):
                # accent = couleur d'axe (identité) — même palette que la vue
                # d'ensemble et les bulles de la matrice de priorisation.
                _slide_recommendation(prs, axis, f"{i + 1}.{j + 1}", reco,
                                      accent=palette[i % len(palette)])

    # Garde-fou géométrique (US7.1) : un texte trop long ou un template client
    # aux dimensions inattendues peut faire déborder une forme de la slide —
    # mieux vaut échouer bruyamment ici qu'exporter un .pptx visuellement cassé.
    problemes = D.verifier_geometrie(prs)
    if problemes:
        raise RuntimeError(
            "Export PPT : formes hors cadre détectées —\n" + "\n".join(problemes)
        )

    return prs
