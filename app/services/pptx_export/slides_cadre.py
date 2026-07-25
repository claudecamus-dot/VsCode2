"""Slides de cadre du deck : couverture, titre de repli, sommaire quali,
intercalaires de chapitre — et la structure narrative en 4 chapitres.
Extrait de pptx_export.py (découpage du gros module, finding audit
2026-07-24) — code déplacé tel quel."""
from __future__ import annotations

from datetime import datetime, timezone

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from ...models import Mission
from .. import pptx_deck as D
from .base import (
    MARGIN,
    _dims,
    _layout_by_name,
    _new_slide,
    _per_line_height_in,
    _pick_layout,
    _sans_puce,
)
from .images import _FRAMED_OK, _find_teardrop_frame, _remplir_cadre_chapitre

# Refonte P2 — structure narrative : 4 chapitres. Chaque section du deck appartient
# à un chapitre ; un intercalaire (layout « 50 - Chapitre ») ouvre chaque chapitre
# qui a du contenu, et le sommaire quali regroupe les sections sous ces intitulés
# narratifs (couleur = repère de navigation, reprise sur l'intercalaire).
_CH_RETENIR, _CH_DIAGNOSTIC, _CH_PAROLE, _CH_TRAJECTOIRE = 0, 1, 2, 3
# (intitulé, couleur, scène image de l'intercalaire — cf. _slide_chapitre P3).
# (intitulé, couleur, scène, sous-titre italique — format d'intercalaire VSCode3 :
# le placeholder de titre porte titre + sous-titre, cf. _slide_chapitre).
_CHAPITRES = [
    ("Ce qu'il faut retenir", "#00D2DD", "sunset", "L'essentiel de la mission, en une page."),
    ("Le diagnostic", "#0E2356", "mountains", "Ce que les entretiens révèlent, sans détour."),
    ("La parole des équipes", "#138086", "forest", "Les mots des acteurs, tels quels."),
    ("La trajectoire proposée", "#6a3d9a", "ocean", "Où aller, et par quoi commencer."),
]


def _slide_cover(prs: Presentation, mission: Mission) -> None:
    """Couverture : sur un template OCTO, remplit les placeholders du layout de
    marque « 40 - Couverture » (titre/sous-titre/date) — la mise en forme (police
    Outfit, tailles, éventuelle photo de couverture) vient du template. Repli sur un
    titre dessiné centré si le layout n'existe pas (deck synthétique sans marque)."""
    subtitle = "Synthèse transverse & recommandations"
    date_str = datetime.now(timezone.utc).strftime("%d/%m/%Y")
    cover = _layout_by_name(prs, "couverture", "cover")
    if cover is not None:
        slide = prs.slides.add_slide(cover)
        phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
        if 0 in phs:
            phs[0].text_frame.text = mission.name
            D.appliquer_police(phs[0].text_frame)
        if 1 in phs:
            phs[1].text_frame.text = subtitle
            D.appliquer_police(phs[1].text_frame)
        # idx2 = « OCTO Technology » : le texte d'invite du layout ne REND pas —
        # laissé vide, la pastille affichait « | date » avec un trou à gauche
        # (constat utilisateur 2026-07-22) ; à remplir explicitement. idx3 = date.
        if 2 in phs:
            phs[2].text_frame.text = "OCTO Technology"
            D.appliquer_police(phs[2].text_frame)
        if 3 in phs:
            phs[3].text_frame.text = date_str
            D.appliquer_police(phs[3].text_frame)
        return
    _slide_title(prs, mission, subtitle, date_str)


def _slide_title(prs: Presentation, mission: Mission, subtitle: str, date_str: str) -> None:
    slide = prs.slides.add_slide(_pick_layout(prs))
    w_in, h_in = _dims(prs)
    title_w = w_in - 2.0
    size = 32
    # Le sous-titre se place sous la hauteur RÉELLE du titre (un nom long passe sur
    # 2 lignes sur une slide étroite type OCTO 10in) plutôt qu'à un offset fixe qui
    # supposait un titre sur 1 ligne — sinon la 2e ligne chevauche le sous-titre.
    title_lines = max(1, min(3, D.estimer_lignes(mission.name, title_w - 0.2, size)))
    title_h = title_lines * _per_line_height_in(size)
    cy = h_in * 0.38
    D.add_text(
        slide, 1.0, cy, title_w, title_h,
        [(mission.name, {"size": size, "bold": True, "color": D.INK, "align": PP_ALIGN.CENTER})],
    )
    y = cy + title_h + 0.15
    D.add_text(
        slide, 1.0, y, title_w, 0.5,
        [(subtitle, {"size": D.TYPE["h2"], "color": D.MUTED, "align": PP_ALIGN.CENTER})],
    )
    D.add_text(
        slide, 1.0, y + 0.55, title_w, 0.4,
        [(date_str, {"size": D.TYPE["small"], "color": D.MUTED, "align": PP_ALIGN.CENTER})],
    )


def _slide_sommaire(prs: Presentation, ch_sections: list[list[str]]) -> None:
    """Sommaire quali (P2) : les chapitres AYANT du contenu, chacun avec sa pastille
    couleur + intitulé narratif + les sections qu'il regroupe (repère de navigation,
    couleur reprise sur l'intercalaire) — au lieu d'une liste plate de sections."""
    slide, w_in, h_in, top = _new_slide(prs, "Sommaire")
    active = [ci for ci, subs in enumerate(ch_sections) if subs]
    # Grille 2×2 de badges GOUTTE (teardrop) à contour — signature du sommaire des
    # decks OCTO réels (VSCode4) : numéro dans la goutte, intitulé narratif + sections
    # en regard. Remplissage colonne par colonne (01,02 à gauche ; 03,04 à droite).
    area_t = top + 0.2
    row_h = (h_in - 0.5 - area_t) / 2
    col_w = (w_in - 2 * MARGIN) / 2
    d = 0.92  # diamètre du badge goutte
    for idx, ci in enumerate(active):
        subs = ch_sections[ci]
        label, color = _CHAPITRES[ci][0], _CHAPITRES[ci][1]
        col, r = idx // 2, idx % 2
        cell_x = MARGIN + col * col_w
        cell_y = area_t + r * row_h
        bx = cell_x + 0.15
        by = cell_y + (row_h - d) / 2
        D.add_teardrop(slide, bx, by, d, f"{idx + 1:02d}", color, size=D.TYPE["h2"])
        tx = bx + d + 0.3
        tw = col_w - (bx - cell_x) - d - 0.3 - 0.2
        D.add_text(
            slide, tx, cell_y, max(1.0, tw), row_h,
            [
                (label, {"size": D.TYPE["h3"], "bold": True, "color": D.INK, "space_after": 4}),
                (" · ".join(subs), {"size": D.TYPE["small"], "color": D.MUTED}),
            ],
            anchor=MSO_ANCHOR.MIDDLE,
        )


def _slide_chapitre(prs: Presentation, numero: int, titre: str, color: str,
                    scene: str | None = None, sous_titre: str | None = None) -> None:
    """Intercalaire de chapitre — reproduit le format du générateur de référence
    VSCode3 (« 50 - Chapitre ») : idx0 = titre coloré + sous-titre italique gris ;
    idx1 = numéro DANS l'encart logo (17pt, marges à zéro, centré, sans puce — ce qui
    empêche « 01 » de wrapper) ; cadre photo teardrop rempli. Repli dessiné sinon."""
    layout = _layout_by_name(prs, "chapitre") if _FRAMED_OK else None
    if layout is not None:
        slide = prs.slides.add_slide(layout)
        phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
        if 0 in phs:
            # idx0 = titre (couleur de chapitre) + sous-titre italique gris.
            tf0 = phs[0].text_frame
            tf0.text = titre
            for r in tf0.paragraphs[0].runs:
                r.font.color.rgb = D.rgb(color)
            if sous_titre:
                p2 = tf0.add_paragraph()
                p2.text = sous_titre
                for r in p2.runs:
                    r.font.size = Pt(D.TYPE["small"])
                    r.font.italic = True
                    r.font.color.rgb = D.rgb(D.MUTED)
            D.appliquer_police(tf0)  # police du deck, pas l'Outfit du layout
        if 1 in phs:
            # idx1 = numéro DANS l'encart logo (format VSCode3) : marges à zéro + 17pt
            # + centré + sans puce — ce qui empêche « 01 » de wrapper dans le petit
            # encart. Reproduit à l'identique du générateur de référence VSCode3.
            tf1 = phs[1].text_frame
            tf1.text = f"{numero:02d}"
            tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = 0
            tf1.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in tf1.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                _sans_puce(p)
                for r in p.runs:
                    r.font.size = Pt(17)
                    r.font.color.rgb = D.rgb(color)
            D.appliquer_police(tf1)  # police du deck, pas l'Outfit du layout
        _remplir_cadre_chapitre(slide, _find_teardrop_frame(slide.slide_layout.shapes),
                                scene or "mountains")
        return
    slide = prs.slides.add_slide(_pick_layout(prs))
    w_in, h_in = _dims(prs)
    cy = h_in * 0.30
    D.add_text(
        slide, MARGIN + 0.3, cy, 3.0, 1.3,
        [(f"{numero:02d}", {"size": D.TYPE["kpi"], "bold": True, "color": color})],
    )
    ty = cy + 1.35
    D.add_rect(slide, MARGIN + 0.36, ty, 1.2, 0.06, fill=color)
    D.add_text(
        slide, MARGIN + 0.3, ty + 0.18, w_in - 2 * MARGIN - 0.6, 1.0,
        [(titre, {"size": D.TYPE["title"], "bold": True, "color": D.INK})],
    )
