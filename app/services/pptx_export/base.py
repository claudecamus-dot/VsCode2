"""Constantes de mise en page et helpers bas niveau de l'export PPT :
repères de forme FIELD_SHAPE + field_fit_hint (éditeur web), création de
slide/titre, listes à puces mesurées et paginées, champs mesurés, labels
d'axe rotés. Extrait de pptx_export.py (découpage du gros module, finding
audit 2026-07-24) — code déplacé tel quel."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from .. import pptx_deck as D
from .images import _FRAMED_OK

MARGIN = 0.6

# --------------------------------------------------------------------------- #
# Repères "forme" pour l'éditeur web par onglets (aperçu.html) — mêmes
# contraintes géométriques (largeur, hauteur max, échelle typo) que les
# fonctions de slide ci-dessous, dupliquées ici en constantes plutôt que
# recalculées dynamiquement : l'éditeur enregistre un champ à la fois, hors
# contexte d'une Presentation réelle (pas de mission/axes/inclusions connus à
# cet instant). Restent donc des repères indicatifs, pas une garantie — le
# garde-fou qui compte vraiment reste `D.verifier_geometrie()` à l'export.
# plan_actions/resultats_attendus utilisent en vrai l'espace *restant* après
# les blocs précédents (variable) ; ici on prend une estimation généreuse
# mais fixe, cohérente avec une slide "normale".
_W_IN, _H_IN = 10.0, 5.625  # dims du template OCTO de marque (16:9) — FIELD_SHAPE (hints web) aligné dessus
# Fiche reco en encarts arrondis (2026-07-22) : le contenu vit DANS des cartes
# (carte gauche + encart proposition + carte plan), les largeurs utiles sont donc
# les largeurs de carte moins les marges internes (pad 0.2 ×2 + liseré 0.05).
# 3.15 (resserrée) : la carte droite porte plan+résultats, souvent de longues
# puces — elle a besoin de la largeur (cf. _slide_recommendation).
_CARD_L_W = 3.15
_LEFT_W = _CARD_L_W - 0.45
_RIGHT_W = (_W_IN - MARGIN - (MARGIN + _CARD_L_W + 0.3)) - 0.45
# Slide de synthèse enrichie (claim + visuel + encart) : largeur de la carte de puces
# une fois la bande photo réservée à droite (2.7in), sinon pleine largeur (repli).
_SYNTH_VIS_W = 2.7
_SYNTH_AREA_W = (
    ((_W_IN - MARGIN - _SYNTH_VIS_W) - 0.3 - (MARGIN + 0.3))
    if _FRAMED_OK else (_W_IN - 2 * (MARGIN + 0.3))
)

# Template OCTO de marque, versionné (masters/layouts/thème + police Outfit) : défaut de
# build_presentation. Un template client (mission.pptx_template_path) reste prioritaire ;
# le chrome (logo/pied de page/n° de slide) survit via _pick_layout (« titre seul »).
OCTO_TEMPLATE_PATH = Path(__file__).resolve().parents[2] / "assets" / "template-octo.pptx"

# Plancher de réduction d'un titre de slide : un titre trop long pour son budget
# de lignes voit sa police descendre jusqu'ici (h3) au lieu d'être tronqué à
# l'ellipse (demande 2026-07-23) ; si même le plancher ne suffit pas, le titre
# replie sur des lignes supplémentaires et le contenu descend (content_top est
# calculé sur les lignes réelles dans _new_slide).
_TITLE_SIZE_MIN = D.TYPE["h3"]

FIELD_SHAPE = {
    "objectif": dict(width_in=_LEFT_W, max_h_in=1.1),
    "acteurs": dict(width_in=_LEFT_W, max_h_in=0.5),
    # resultats_attendus vit dans le bandeau bas PLEINE LARGEUR de la fiche
    # (encart gris, 2026-07-22) : taille FIXE small, tronqué à 2 lignes — hint en
    # mode size_pt/max_lignes (annonce la troncature, ne promet pas un shrink que
    # la slide ne fait pas ; même dérive corrigée que es_headline/difficulty_label).
    "resultats_attendus": dict(width_in=_W_IN - 2 * MARGIN - 0.45, size_pt=D.TYPE["small"], max_lignes=2),
    # 0.80 = prop_h(1.00) - 0.20 de _slide_recommendation — recalé revue adversariale.
    "proposition_valeur": dict(width_in=_RIGHT_W, max_h_in=0.80),
    "plan_actions": dict(width_in=_RIGHT_W, max_h_in=2.0),
    # 1 ligne préférée : un titre long est rendu en police réduite (jusqu'à
    # _TITLE_SIZE_MIN) par _new_slide plutôt que tronqué — size_min aligne le
    # hint sur ce shrink (jamais de troncature annoncée).
    "reco_title": dict(width_in=_W_IN - 2 * MARGIN, size_pt=D.TYPE["title"], max_lignes=1,
                       size_min=_TITLE_SIZE_MIN),
    "axis_title": dict(width_in=_W_IN - 2 * MARGIN - 2.0, max_h_in=1.1, size_max=D.TYPE["h3"]),
    # Slide enrichie (carte de puces à gauche, visuel à droite, 1re puce en encart) :
    # largeur = carte réduite du visuel ; hauteur = zone au-dessus de l'encart « à retenir ».
    "synthese_categorie": dict(width_in=_SYNTH_AREA_W - 0.48, max_h_in=1.9, size_max=D.TYPE["body"]),
    # Un quadrant SWOT = ~demi-largeur de la zone de contenu ; la hauteur de la
    # zone de PUCES (pas de la carte) = row_h - titre - paddings ≈ 1.9 in sur un
    # deck vierge (cf. _slide_swot) — pas la demi-hauteur brute (~2.2), qui
    # surestimait le budget du repère de ~20 % et rendait le fit-hint trompeur.
    # Matrice SWOT (skill swot-matrix) : cellule teintée = gouttière 0.30 à gauche +
    # bandeau d'axe en haut ; largeur de puces = demi-grille moins paddings.
    "swot_quadrant": dict(width_in=(_W_IN - MARGIN - (MARGIN + 0.30) - 0.22) / 2 - 0.36, max_h_in=1.2, size_max=D.TYPE["small"]),
    # Executive summary (piste F) : panneau pleine largeur (constat + points) et
    # bande cyan « key message » — mêmes contraintes que la slide (cf.
    # _slide_executive_summary), pour un fit-hint fidèle dans l'aperçu.
    # headline / key_message : rendus à taille FIXE (h3) et tronqués à 2 lignes par
    # _slide_executive_summary -> hint en mode size_pt/max_lignes (annonce la
    # troncature, ne PROMET pas de réduction de police que la slide ne fait pas).
    # Constat revue adversariale 2026-07-21 : l'ancien mode max_h_in promettait un
    # shrink inexistant (même classe de dérive que le fit-hint SWOT déjà corrigé).
    "es_headline": dict(width_in=_W_IN - 2 * (MARGIN + 0.3) - 0.48, size_pt=D.TYPE["h3"], max_lignes=2),
    "es_points": dict(width_in=_W_IN - 2 * (MARGIN + 0.3) - 0.48, max_h_in=2.5, size_max=D.TYPE["body"]),
    "es_key_message": dict(width_in=_W_IN - 2 * (MARGIN + 0.3) - 0.48, size_pt=D.TYPE["h3"], max_lignes=2),
    # Difficulté (planche §D.1) : libellé d'une carte, taille fixe body, tronqué à
    # 3 lignes par _slide_difficultes -> hint size_pt/max_lignes (honnête). Largeur
    # réduite du chip de rang à gauche (2*pad + rang_w 0.46 + gap 0.16 = 0.98).
    "difficulty_label": dict(width_in=_W_IN - 2 * (MARGIN + 0.3) - 0.98, size_pt=D.TYPE["body"], max_lignes=3),
}

def field_fit_hint(field_key: str, text: str) -> str:
    """Message court indiquant comment `text` sera rendu à l'export pour le
    champ `field_key` (police retenue, nombre de lignes, troncature
    éventuelle) — s'appuie sur les mêmes fonctions d'ajustement
    (`D.ajuster_police` / `D.tronquer_a_lignes` / `D.estimer_lignes`) que le
    générateur, appliquées aux contraintes de forme de `FIELD_SHAPE` (voir
    note du module). Chaîne vide si le champ est inconnu ou vide — pas de
    repère à afficher plutôt qu'un repère trompeur."""
    spec = FIELD_SHAPE.get(field_key)
    text = (text or "").strip()
    if spec is None or not text:
        return ""

    width_in = spec["width_in"]

    if "max_lignes" in spec:
        size = spec["size_pt"]
        size_min = spec.get("size_min")
        if size_min is not None:
            # Champ jamais tronqué (titre de slide) : même réduction de police
            # que _new_slide, puis repli sur des lignes supplémentaires.
            while size > size_min and D.estimer_lignes(text, width_in, size) > spec["max_lignes"]:
                size = max(size_min, size - 1.0)
            lignes = D.estimer_lignes(text, width_in, size)
            if size != spec["size_pt"] or lignes > spec["max_lignes"]:
                return (f"⚠ long — rendu en {size:.0f}pt sur {lignes} ligne(s) "
                        "à l'export (police réduite, jamais tronqué)")
            return f"{lignes} ligne(s) à {size:.0f}pt à l'export"
        lignes = D.estimer_lignes(text, width_in, size)
        if lignes > spec["max_lignes"]:
            return f"⚠ trop long — sera tronqué à {spec['max_lignes']} lignes à l'export"
        return f"{lignes} ligne(s) à {size:.0f}pt à l'export"

    max_h_in = spec["max_h_in"]
    size_max = spec.get("size_max", D.TYPE["body"])
    size_min = D.TYPE["tiny"]

    budget_ok = _budget_lignes(max_h_in)  # réserve ½ ligne, cohérente avec l'export

    size, lignes = D.ajuster_police([text], width_in, size_max, size_min, budget_ok=budget_ok)
    if lignes * _per_line_height_in(size) > max_h_in:
        return f"⚠ très long — sera réduit à {size:.0f}pt et tronqué à l'export"
    if size < size_max - 0.5:
        return f"{lignes} ligne(s) — police réduite à {size:.0f}pt pour tenir à l'export"
    return f"{lignes} ligne(s) à {size:.0f}pt à l'export"


def _dims(prs: Presentation) -> tuple[float, float]:
    return Emu(prs.slide_width).inches, Emu(prs.slide_height).inches


def _clear_slides(prs: Presentation) -> None:
    """Retire toutes les slides d'une présentation chargée depuis un template
    client — on ne veut hériter que masters/layouts/thème. Sans ça, un
    template qui est un vrai exemple de deck (pas un .potx vierge) ferait
    apparaître tout son contenu d'origine avant le nôtre. `python-pptx`
    n'expose pas de suppression de slide côté API publique ; on vide
    directement la liste XML des slides — mais il faut aussi lâcher la
    relation (r:id) de chaque slide sur la part présentation, sans quoi le
    fichier réserialisé contient des relations pointant vers des parts
    devenues orphelines : invisible pour python-pptx (parseur tolérant),
    mais PowerPoint refuse ensuite d'ouvrir le fichier (constaté via
    l'automation COM — l'export semblait « valide » côté tests avant ça)."""
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.get(qn("r:id")))
        sld_id_lst.remove(sld_id)


def _has_title_placeholder(layout) -> bool:
    try:
        return any(
            ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
            for ph in layout.placeholders
        )
    except Exception:
        return False


def _pick_layout(prs: Presentation, preferred: int = 6):
    """Choisit un layout pour du contenu personnalisé, en respectant au mieux
    le template injecté (point 6) : un layout nommé "Title Only"/"Section" (le
    logo/footer du master survit, pas de placeholder de corps qui entrerait en
    collision avec notre mise en page), sinon le layout avec titre le moins
    chargé en autres placeholders, sinon le comportement historique (repli
    toujours valide même sur un template client mal structuré)."""
    layouts = list(prs.slide_layouts)
    # « titre seul » = le layout de contenu OCTO (idx0 titre, garde logo/pied/n° de slide).
    for kw in ("titre seul", "title only", "section"):
        for layout in layouts:
            if kw in (layout.name or "").lower() and _has_title_placeholder(layout):
                return layout
    with_title = [l for l in layouts if _has_title_placeholder(l)]
    if with_title:
        return min(with_title, key=lambda l: len(l.placeholders))
    return layouts[preferred] if preferred < len(layouts) else layouts[-1]


def _new_slide(prs: Presentation, title: str, max_title_lines: int = 2):
    """Crée une slide de contenu et pose son titre. Renvoie
    `(slide, w_in, h_in, content_top)` — `content_top` est calculé à partir
    de la position/hauteur réelle du placeholder de titre (natif du
    template) et du nombre de lignes qu'occupera effectivement `title` une
    fois replié, plutôt qu'une constante suppposant un titre sur une seule
    ligne : un titre de longueur normale (~50 caractères) suffit à passer
    sur 2 lignes et, avec une position de contenu figée, à chevaucher la
    zone en dessous — ce qu'une constante ne peut pas anticiper."""
    slide = prs.slides.add_slide(_pick_layout(prs))
    w_in, h_in = _dims(prs)
    title_shape = slide.shapes.title
    # Le placeholder de titre natif hérite position/police du template —
    # préféré à une zone de texte dessinée à la main dès qu'il existe sur le
    # layout choisi. On fige sa taille de police sur D.TYPE["title"] (au lieu
    # de laisser le style hérité, potentiellement bien plus grand) : ça reste
    # cohérent avec l'unique échelle typographique du reste du deck, et ça
    # rend le nombre de lignes prévisible (donc calculable) plutôt que soumis
    # à un style de thème inconnu.
    if title_shape is not None:
        if getattr(prs, "_i2d_synthetic", False):
            # Présentation vierge (pas de template client) : le placeholder
            # de titre hérité du modèle par défaut de python-pptx est
            # dimensionné pour un slide 10x7.5in (4:3) — trop étroit une fois
            # la slide passée en 16:9. On le repositionne explicitement sur
            # CETTE slide (jamais sur le layout/master : leurs placeholders
            # sont résolus par héritage et se sont révélés instables à muter
            # directement avec python-pptx — cf. essai précédent).
            title_shape.left = Inches(MARGIN)
            title_shape.top = Inches(0.3)
            title_shape.width = Inches(w_in - 2 * MARGIN)
            title_shape.height = Inches(1.1)
        title_w_in = Emu(title_shape.width).inches if title_shape.width is not None else (w_in - 2 * MARGIN)
        title_top_in = Emu(title_shape.top).inches if title_shape.top is not None else 0.3
        # Jamais de troncature d'un titre de slide (demande 2026-07-23) : un
        # titre trop long pour `max_title_lines` voit d'abord sa police réduite
        # (par pas de 1 pt jusqu'au plancher _TITLE_SIZE_MIN) ; si même le
        # plancher ne suffit pas, il replie sur des lignes supplémentaires —
        # content_top suit les lignes réelles, le contenu descend d'autant.
        size = D.TYPE["title"]
        while (size > _TITLE_SIZE_MIN
               and D.estimer_lignes(title, title_w_in, size) > max_title_lines):
            size = max(_TITLE_SIZE_MIN, size - 1.0)
        title_shape.text = title
        tf = title_shape.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        for p in tf.paragraphs:
            for run in p.runs:
                run.font.size = Pt(size)
                run.font.bold = True
                if D.POLICE:  # police du deck (thème), pas l'Outfit hérité du layout
                    run.font.name = D.POLICE
        lignes = D.estimer_lignes(title, title_w_in, size)
        needed_h = lignes * _per_line_height_in(size) + 0.15
        # Garde template client (defer revue adversariale) : si l'ancre verticale
        # du titre est MIDDLE/BOTTOM (posée sur la slide — l'héritage renvoie
        # None, traité comme TOP), le texte peut flotter bas dans sa boîte —
        # réserver alors toute la hauteur de boîte. Sans effet sur OCTO (TOP).
        if tf.vertical_anchor in (MSO_ANCHOR.MIDDLE, MSO_ANCHOR.BOTTOM):
            box_h = Emu(title_shape.height).inches if title_shape.height is not None else 0.7
            needed_h = max(needed_h, box_h)
        # Pas de barre d'accent avant le titre : les decks OCTO réels (VSCode4) n'en
        # ont pas — titre navy + logo suffisent (retrait demandé 2026-07-22, charte VSCode4).
        # Hauteur RÉELLE du texte de titre (plancher 0.55), pas la boîte du
        # placeholder — indépendant de la hauteur de boîte du template. Réserve
        # connue (revue adversariale) : un template CLIENT à boîte de titre haute
        # ancrée middle/bottom pourrait voir le contenu remonter dans sa zone de
        # titre — à re-vérifier au premier template client réel.
        content_top = title_top_in + max(needed_h, 0.55) + 0.25
    else:
        D.add_text(
            slide, MARGIN, 0.35, w_in - 2 * MARGIN, 0.7,
            [(title, {"size": D.TYPE["title"], "bold": True, "color": D.INK})],
        )
        content_top = 1.4
    return slide, w_in, h_in, content_top


def _bullet_lines(text: str) -> list[str]:
    lines = []
    for raw in (text or "").splitlines():
        line = raw.strip().lstrip("-•").strip()
        if line:
            lines.append(line)
    return lines


def _per_line_height_in(size_pt: float) -> float:
    # Calibrage empirique (skill pptx-deck) : ~0.17in/ligne à 10.5pt, un peu
    # plus large ici pour couvrir l'espacement inter-puces (space_after).
    return size_pt * 0.017 + 4 / 72


def _budget_lignes(max_h_in: float, reserve_lignes: float = 0.5):
    """Fabrique le `budget_ok` standard de D.ajuster_police : `lignes_max` lignes à
    `taille` pt tiennent dans `max_h_in`, moins une réserve exprimée en fraction de
    ligne (½ ligne par défaut — la même que la pagination de _add_bulleted_text,
    sinon le hint annonce à la frontière une taille que l'export réduit).

    Factorisation du finding audit 2026-07-24 (« budget_ok redéfini 3 fois ») : les
    deux définitions de MÊME forme passent par ici ; celle de _add_bulleted_text
    (somme d'estimations par ligne) garde sa logique propre, documentée sur place."""
    def budget_ok(taille, lignes_max):
        return lignes_max * _per_line_height_in(taille) <= max_h_in - reserve_lignes * _per_line_height_in(taille)
    return budget_ok


def _add_bulleted_text(
    slide, l, t, w, h, text: str, size: float | None = None,
    anchor=MSO_ANCHOR.TOP, size_max: float = D.TYPE["body"], size_min: float = D.TYPE["tiny"],
    paginate: bool = False,
) -> list[str]:
    """Pose une liste à puces dans la zone donnée. Par défaut (`paginate=False`,
    comportement historique) la police est réduite jusqu'à `size_min` pour
    tenter de tout faire tenir, sans garantie — un texte très long peut
    déborder silencieusement de sa zone (indétectable par
    `D.verifier_geometrie`, qui ne voit que les bords des formes).

    `paginate=True` change la garantie : si même `size_min` ne suffit pas, les
    puces qui ne tiennent pas sont retenues (pas rendues) et renvoyées à
    l'appelant plutôt que de déborder — à charge pour lui de les poser sur
    une slide de continuation (voir `_emit_bullet_overflow`). Renvoie la
    liste des puces non rendues (vide si tout tient)."""
    lines = _bullet_lines(text) or ["—"]

    def budget_ok(taille, _lignes_max):
        # Le budget inclut la MÊME réserve d'une demi-ligne que la pagination —
        # sinon l'ajusteur valide une taille que la pagination recoupe ensuite
        # (explosion de slides de suite constatée le 2026-07-22) : on préfère
        # rétrécir la police d'un cran que couper une puce.
        total = sum(D.estimer_lignes(line, w, taille) for line in lines)
        return total * _per_line_height_in(taille) <= h - 0.5 * _per_line_height_in(taille)

    if size is None:
        size, _ = D.ajuster_police(lines, w, size_max, size_min, budget_ok=budget_ok)

    overflow: list[str] = []
    if paginate:
        # Capacité MINORÉE d'une demi-ligne : l'estimation de repli est optimiste
        # pour du français — sans réserve, le dernier bloc d'une carte sortait du
        # cadre au vrai rendu (défaut récurrent 2026-07-22). Et un item SEUL plus
        # haut que la zone est insécable pour paginer_items — il débordait en
        # silence (attrapé par verifier_debordements_texte) : désormais tronqué à
        # l'ellipse ici, son texte COMPLET partant sur la slide de suite.
        line_h = _per_line_height_in(size)
        capacite = max(line_h, h - 0.5 * line_h)
        if sum(D.estimer_lignes(li, w, size) for li in lines) * line_h > capacite:
            pages = D.paginer_items(
                lines, lambda li: D.estimer_lignes(li, w, size) * line_h,
                capacite_in=capacite,
            )
            lines = pages[0]
            overflow = [li for page in pages[1:] for li in page]
            if len(lines) == 1 and D.estimer_lignes(lines[0], w, size) * line_h > capacite:
                complet = lines[0]
                lines = [D.tronquer_a_lignes(complet, w, size, max(1, int(capacite / line_h)))]
                overflow = [complet] + overflow

    paragraphs = [(f"•  {line}", {"size": size, "color": D.INK, "space_after": 4}) for line in lines]

    if anchor == MSO_ANCHOR.MIDDLE:
        total_lines = sum(D.estimer_lignes(line, w, size) for line in lines)
        # +0.5 ligne de marge dans la boîte centrée — même logique que partout :
        # l'estimation nominale est optimiste, la boîte exacte faisait peindre la
        # dernière ligne hors boîte (constat verifier_debordements_texte).
        content_h = min(h, (total_lines + 0.5) * _per_line_height_in(size))
        box_t = t + max(0.0, (h - content_h) / 2)
        D.add_text(slide, l, box_t, w, content_h, paragraphs)
    else:
        D.add_text(slide, l, t, w, h, paragraphs)
    return overflow


def _emit_bullet_overflow(prs: Presentation, base_title: str, field_label: str, overflow_lines: list[str]) -> None:
    """Pose les puces qui n'ont pas tenu sur la slide d'origine (voir
    `_add_bulleted_text(paginate=True)`) sur une ou plusieurs slides de
    continuation pleine largeur — chacune dispose de bien plus d'espace que
    la colonne étroite d'origine, donc peut se voir attribuer sa propre
    police (recalculée, pas figée à `size_min`)."""
    remaining = "\n".join(overflow_lines)
    page_no = 1
    while remaining:
        suffix = f" {page_no}" if page_no > 1 else ""
        slide, w_in, h_in, top = _new_slide(prs, f"{base_title} (suite — {field_label}){suffix}")
        w = w_in - 2 * (MARGIN + 0.3)
        h = h_in - top - 0.5
        overflow = _add_bulleted_text(slide, MARGIN + 0.3, top, w, h, remaining, paginate=True)
        nouveau = "\n".join(overflow)
        if nouveau == remaining:
            # Garde de PROGRESSION (revue adversariale) : une puce insécable plus
            # haute qu'une slide de suite entière revient intégralement en
            # overflow (le chemin de troncature repousse le texte COMPLET) —
            # sans cette garde, boucle infinie + slides sans fin. La slide qui
            # vient d'être posée montre déjà tout ce qui tient, à l'ellipse.
            break
        remaining = nouveau
        page_no += 1


def _add_measured_field(
    slide, l, t, w, label: str, text: str, max_h: float,
    size_max: float = D.TYPE["body"], size_min: float = D.TYPE["tiny"],
    bold: bool = False, italic: bool = False,
) -> float:
    """Pose un libellé (petit, gras, discret) puis son contenu juste en
    dessous, en adaptant la taille de police du contenu à `max_h`
    (D.ajuster_police) et en tronquant en tout dernier recours — jamais de
    débordement dans le bloc suivant même avec une réponse d'entretien très
    longue. Renvoie la hauteur réellement occupée (libellé + contenu), à
    utiliser pour empiler le bloc suivant à la bonne position."""
    label_h = 0.3
    D.add_text(slide, l, t, w, label_h, [(label, {"size": D.TYPE["small"], "bold": True, "color": D.MUTED})])
    body = ((text or "").strip()) or "—"
    body_max_h = max(0.2, max_h - label_h)

    budget_ok = _budget_lignes(body_max_h, reserve_lignes=0.0)  # comportement historique : sans réserve

    size, lignes_max = D.ajuster_police([body], w, size_max, size_min, budget_ok=budget_ok)
    if lignes_max * _per_line_height_in(size) > body_max_h:
        max_lignes = max(1, int(body_max_h / _per_line_height_in(size)))
        body = D.tronquer_a_lignes(body, w, size, max_lignes)
        lignes_max = max_lignes
    # +0.6 ligne de marge : la boîte était dimensionnée à l'estimation EXACTE —
    # au vrai repli PowerPoint (un peu plus gourmand), la dernière ligne sortait
    # du cadre (constat verifier_debordements_texte sur les fiches, 2026-07-22).
    body_h = (lignes_max + 0.6) * _per_line_height_in(size)
    D.add_text(slide, l, t + label_h, w, body_h, [(body, {"size": size, "bold": bold, "italic": italic, "color": D.INK})])
    return label_h + body_h


def _layout_by_name(prs: Presentation, *keywords: str):
    """Premier layout dont le nom contient l'un des mots-clés (insensible à la
    casse) — repérage robuste des layouts de marque OCTO (« 40 - Couverture »,
    « 50 - Chapitre ») par nom plutôt que par indice. None si aucun ne matche."""
    for layout in prs.slide_layouts:
        name = (layout.name or "").lower()
        if any(kw in name for kw in keywords):
            return layout
    return None


def _sans_puce(paragraph) -> None:
    """Retire l'indentation de puce héritée (marL/indent) et la puce elle-même —
    reproduit tel quel le helper du générateur de référence VSCode3. Cause réelle
    du « 01 » qui wrappe dans le petit encart numéro du layout Chapitre : le style
    de liste hérité pose marL=0.5in dans un encart de ~0.55in. python-pptx n'expose
    pas ces attributs -> manipulation XML directe."""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", "0")
    pPr.set("indent", "0")
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    # Forcer explicitement l'absence de puce : notre template hérite un caractère de
    # puce à un niveau que le retrait ci-dessus ne couvre pas (un ◉ résiduel
    # apparaissait avant le numéro) — buNone le neutralise à coup sûr.
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def _label_axe_vertical(slide, cx: float, cy: float, longueur: float,
                        epaisseur: float, texte: str,
                        size: float | None = None) -> None:
    """Label d'axe roté 270° (lecture bas→haut), centré sur (cx, cy). `longueur`
    = dimension le long du texte (≈ hauteur de la ligne couverte), `epaisseur` =
    largeur de la gouttière. Sert aux libellés de ligne INTERNE/EXTERNE de la
    matrice SWOT et à l'axe Valeur de la matrice de priorisation."""
    box = D.add_text(
        slide, cx - longueur / 2, cy - epaisseur / 2, longueur, epaisseur,
        [(texte, {"size": size or D.TYPE["small"], "bold": True, "color": D.MUTED})],
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )
    box.rotation = 270
